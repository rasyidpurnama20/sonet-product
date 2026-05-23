from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── COLOR PALETTE ──────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG       = rgb("0A0A0F")   # near-black bg
HARARI_R = rgb("E63946")   # Harari red
NG_B     = rgb("00B4D8")   # Ng cyan-blue
GOLD     = rgb("FFD166")   # accent gold
WHITE    = rgb("FFFFFF")
GRAY     = rgb("AAAAAA")
CARD_BG  = rgb("14141E")
LINE_COL = rgb("2A2A3E")
GREEN    = rgb("06D6A0")

# ── HELPERS ────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, color, line=None):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def htxt(slide, text, l, t, w, h, size=18, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Multi-line helper — splits on \\n into separate paragraphs."""
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — EPIC COVER
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)

# left red half + right blue half (subtle)
rect(sl, 0,   0, 6.66, 7.5, rgb("0D0508"))
rect(sl, 6.66,0, 6.67, 7.5, rgb("030D12"))

# center divider glow line
rect(sl, 6.5, 0, 0.08, 7.5, rgb("333355"))

# top bar
rect(sl, 0, 0, 13.33, 0.06, GOLD)

# badge
rect(sl, 4.5, 0.35, 4.33, 0.55, rgb("1A1A2E"), line=GOLD)
txt(sl, "⚡  WEBINAR NASIONAL  ⚡", 4.5, 0.35, 4.33, 0.55,
    size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# main title
htxt(sl, "BATTLE OF MINDS", 1.0, 1.1, 11.33, 1.4,
     size=62, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# subtitle
htxt(sl, "AI: ANCAMAN atau PELUANG?", 1.0, 2.55, 11.33, 0.8,
     size=26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

# vs badge center
rect(sl, 5.9, 3.45, 1.53, 0.7, rgb("1A1A2E"), line=GOLD)
txt(sl, "VS", 5.9, 3.45, 1.53, 0.7,
    size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# left: Harari
rect(sl, 0.4, 3.35, 5.2, 0.85, rgb("2A0810"), line=HARARI_R)
txt(sl, "🔴  YUVAL NOAH HARARI", 0.5, 3.38, 5.0, 0.8,
    size=17, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)

# right: Ng
rect(sl, 7.73, 3.35, 5.2, 0.85, rgb("031420"), line=NG_B)
txt(sl, "🔵  ANDREW NG", 7.83, 3.38, 5.0, 0.8,
    size=17, bold=True, color=NG_B, align=PP_ALIGN.CENTER)

# left subtitle
htxt(sl, "Sejarawan · Penulis Sapiens\nPerspektif Kritis & Historis",
     0.4, 4.32, 5.2, 0.9, size=12, italic=True, color=GRAY,
     align=PP_ALIGN.CENTER)

# right subtitle
htxt(sl, "Pendiri DeepLearning.AI · Coursera\nPerspektif Optimis & Pragmatis",
     7.73, 4.32, 5.2, 0.9, size=12, italic=True, color=GRAY,
     align=PP_ALIGN.CENTER)

# topics
htxt(sl,
     "Topik 1: AI dalam Dunia Pendidikan   |   Topik 2: AI dalam Dunia Kerja",
     1.0, 5.5, 11.33, 0.55, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# bottom bar
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / RUNDOWN
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, GOLD)
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)

txt(sl, "📋  AGENDA WEBINAR", 0.5, 0.2, 12.0, 0.65,
    size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
rect(sl, 4.0, 0.9, 5.33, 0.05, LINE_COL)

items = [
    ("🔥", "01", "OPENING — Pertanyaan Pemantik", "Apakah AI ancaman atau peluang bagimu?"),
    ("🔴", "02", "PERSPEKTIF HARARI", "AI dalam Pendidikan & Karir — Sisi Kritis"),
    ("🔵", "03", "PERSPEKTIF ANDREW NG", "AI dalam Pendidikan & Karir — Sisi Optimis"),
    ("⚡", "04", "DUEL PERSPEKTIF", "Face-to-face: 4 tema utama dibenturkan"),
    ("💬", "05", "SESI INTERAKTIF", "Poll, Q&A, & Refleksi Peserta"),
    ("🚀", "06", "ACTION PLAN", "Langkah konkret untuk kamu mulai hari ini"),
]

for i, (icon, num, title, sub) in enumerate(items):
    y = 1.1 + i * 0.97
    rect(sl, 0.5, y, 0.7, 0.75, CARD_BG, line=LINE_COL)
    txt(sl, num, 0.5, y, 0.7, 0.75,
        size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, 1.35, y, 11.0, 0.75, CARD_BG, line=LINE_COL)
    txt(sl, f"{icon}  {title}", 1.5, y+0.04, 6.5, 0.38,
        size=15, bold=True, color=WHITE)
    txt(sl, sub, 1.5, y+0.38, 10.5, 0.35,
        size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — POLL PEMBUKA (Interaktif)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, GOLD)
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)

rect(sl, 4.0, 0.18, 5.33, 0.55, rgb("1A1A2E"), line=GOLD)
txt(sl, "⚡  PERTANYAAN PEMBUKA", 4.0, 0.18, 5.33, 0.55,
    size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

htxt(sl, "Menurutmu, AI itu\nANCAMAN atau PELUANG?",
     0.6, 0.9, 12.13, 1.6, size=40, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "🗳️  Ketik jawabanmu di kolom chat sekarang!", 0.6, 2.6, 12.13, 0.5,
    size=16, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# option A
rect(sl, 0.6, 3.25, 5.5, 1.6, rgb("2A0810"), line=HARARI_R)
txt(sl, "🔴  A", 0.8, 3.35, 5.1, 0.5,
    size=28, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)
htxt(sl, "ANCAMAN\nAI akan mengambil pekerjaan\n& merusak tatanan sosial",
     0.8, 3.85, 5.1, 0.9, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# option B
rect(sl, 7.23, 3.25, 5.5, 1.6, rgb("031420"), line=NG_B)
txt(sl, "🔵  B", 7.43, 3.35, 5.1, 0.5,
    size=28, bold=True, color=NG_B, align=PP_ALIGN.CENTER)
htxt(sl, "PELUANG\nAI akan membuka lapangan\nkerja & akselerasi kemajuan",
     7.43, 3.85, 5.1, 0.9, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# option C center
rect(sl, 5.3, 3.25, 2.73, 1.6, rgb("1A1A10"), line=GOLD)
txt(sl, "⚡  C", 5.3, 3.35, 2.73, 0.5,
    size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
htxt(sl, "KEDUANYA\nTergantung bagaimana\nkita menyikapinya",
     5.3, 3.85, 2.73, 0.9, size=12, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "💡  Tidak ada jawaban salah — semua perspektif valid dan akan kita bahas hari ini!",
    0.6, 5.1, 12.13, 0.5, size=13, italic=True,
    color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HARARI QUOTE CARD (Pembuka sesi Harari)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, rgb("0D0508"))
rect(sl, 0, 0, 13.33, 0.06, HARARI_R)
rect(sl, 0, 7.44, 13.33, 0.06, HARARI_R)

# left red strip
rect(sl, 0, 0.06, 0.25, 7.38, HARARI_R)

# big quote mark
txt(sl, "\u201c", 0.5, 0.3, 3.0, 2.5, size=160, bold=True,
    color=rgb("3A0A10"), align=PP_ALIGN.LEFT)

htxt(sl,
     "AI bukan sekadar alat.\nIA adalah kekuatan yang bisa\nbelajar, berubah, dan memanipulasi.",
     0.6, 1.0, 10.5, 2.5, size=34, bold=True,
     color=WHITE, align=PP_ALIGN.LEFT)

rect(sl, 0.6, 3.65, 6.0, 0.05, HARARI_R)

htxt(sl, "— Yuval Noah Harari\nSejarawan & Penulis Nexus (2024)",
     0.6, 3.8, 8.0, 0.9, size=16, italic=True,
     color=HARARI_R, align=PP_ALIGN.LEFT)

# right info box
rect(sl, 8.5, 1.5, 4.3, 4.0, rgb("1A0508"), line=HARARI_R)
htxt(sl,
     "🔴  YUVAL NOAH HARARI\n\n"
     "📚  Sapiens (2011)\n"
     "📚  Homo Deus (2015)\n"
     "📚  21 Lessons (2018)\n"
     "📚  Nexus (2024)\n\n"
     "🎓  Hebrew University Jerusalem\n"
     "🎓  Cambridge — Existential Risk",
     8.7, 1.6, 3.9, 3.8, size=12, color=WHITE)

txt(sl, "\"Jika manusia tidak hati-hati, AI akan mengontrol\nnarasi — dan siapa yang mengontrol narasi, mengontrol dunia.\"",
    0.6, 5.0, 12.0, 1.1, size=14, italic=True,
    color=GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HARARI: AI PENDIDIKAN (simple cards)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, HARARI_R)
rect(sl, 0, 7.44, 13.33, 0.06, HARARI_R)
rect(sl, 0, 0.06, 0.25, 7.38, HARARI_R)

txt(sl, "🔴  HARARI  —  AI dalam Dunia Pendidikan",
    0.4, 0.18, 12.5, 0.6, size=22, bold=True, color=HARARI_R)
rect(sl, 0.4, 0.82, 12.5, 0.04, rgb("2A0810"))

cards = [
    ("⚠️", "DAMPAK", HARARI_R,
     "Anak-anak tumbuh lebih banyak berinteraksi\n"
     "dengan AI daripada manusia — generasi pertama\n"
     "yang pikirannya dibentuk oleh algoritma."),
    ("🧠", "TANTANGAN", rgb("FF6B6B"),
     "AI bisa menulis esai lebih baik dari siswa.\n"
     "Tapi PROSES menulislah yang mengajarkan\n"
     "berpikir — bukan hasilnya."),
    ("🌍", "KETIMPANGAN", GOLD,
     "Negara kaya → infrastruktur AI kuat → semakin\n"
     "maju. Negara berkembang berisiko tertinggal.\n"
     "Peluang ini TIDAK demokratis."),
    ("📖", "LITERASI BARU", GREEN,
     "Melek huruf era AI = mampu membedakan\n"
     "narasi dari manusia vs algoritma.\n"
     "Tanpa ini, manusia mudah dimanipulasi."),
]

for i, (icon, label, col, body) in enumerate(cards):
    x = 0.4 + i * 3.15
    rect(sl, x, 1.0, 2.95, 5.8, CARD_BG, line=col)
    rect(sl, x, 1.0, 2.95, 0.55, rgb("0D0D15"))
    txt(sl, f"{icon}  {label}", x+0.1, 1.03, 2.75, 0.5,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, body, x+0.15, 1.65, 2.65, 4.8,
         size=12.5, color=WHITE)

txt(sl,
    "💬  Pertanyaan untuk peserta: \"Pernahkah kamu merasa lebih nyaman bertanya ke AI daripada ke guru/dosen?\"",
    0.4, 6.85, 12.5, 0.55, size=12, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HARARI: AI DUNIA KERJA
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, HARARI_R)
rect(sl, 0, 7.44, 13.33, 0.06, HARARI_R)
rect(sl, 0, 0.06, 0.25, 7.38, HARARI_R)

txt(sl, "🔴  HARARI  —  AI dalam Dunia Kerja",
    0.4, 0.18, 12.5, 0.6, size=22, bold=True, color=HARARI_R)

# big stat
rect(sl, 0.4, 0.9, 4.1, 2.5, rgb("2A0810"), line=HARARI_R)
htxt(sl, "800 Juta\npekerjaan", 0.5, 1.0, 3.9, 1.1,
     size=36, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)
htxt(sl, "berisiko terotomasi\noleh AI pada 2030\n(McKinsey Global Institute)",
     0.5, 2.1, 3.9, 1.2, size=11, italic=True,
     color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 4.7, 0.9, 4.1, 2.5, rgb("1A0A0A"), line=rgb("FF6B6B"))
htxt(sl, "\"Useless\nClass\"", 4.8, 1.0, 3.9, 1.1,
     size=36, bold=True, color=rgb("FF6B6B"), align=PP_ALIGN.CENTER)
htxt(sl, "Harari: Revolusi AI bisa\nmenciptakan kelas manusia\nyang tidak dibutuhkan ekonomi",
     4.8, 2.1, 3.9, 1.2, size=11, italic=True,
     color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 9.0, 0.9, 3.9, 2.5, rgb("1A150A"), line=GOLD)
htxt(sl, "Ketimpangan\nGlobal", 9.1, 1.0, 3.7, 1.1,
     size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
htxt(sl, "AI menguntungkan negara\nkaya. Negara berkembang\nberisiko tertinggal.",
     9.1, 2.1, 3.7, 1.2, size=11, italic=True,
     color=GRAY, align=PP_ALIGN.CENTER)

# bottom 3 points
points = [
    ("🔁", "Perubahan STRUKTURAL", "Bukan hanya otomasi fisik — AI menggantikan KOGNITIF:\nhukum, akuntansi, diagnosis, penulisan."),
    ("⚖️", "Tanggung Jawab SIAPA?", "\"Reskilling\" seolah mengalihkan tanggung jawab ke\nindividu. Padahal ini gagal sistem, bukan gagal pekerja."),
    ("🏛️", "Butuh REGULASI", "Kita butuh regulasi global AI yang demokratis —\nbukan hanya optimisme tanpa tata kelola."),
]
for i, (icon, title, body) in enumerate(points):
    x = 0.4 + i * 4.3
    rect(sl, x, 3.6, 4.1, 2.8, CARD_BG, line=HARARI_R)
    txt(sl, f"{icon}  {title}", x+0.1, 3.68, 3.9, 0.45,
        size=13, bold=True, color=HARARI_R)
    rect(sl, x+0.1, 4.15, 3.7, 0.04, rgb("3A0A10"))
    htxt(sl, body, x+0.1, 4.25, 3.85, 2.0, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ANDREW NG QUOTE CARD
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, rgb("030D12"))
rect(sl, 0, 0, 13.33, 0.06, NG_B)
rect(sl, 0, 7.44, 13.33, 0.06, NG_B)
rect(sl, 13.08, 0.06, 0.25, 7.38, NG_B)

# big quote
txt(sl, "\u201c", 9.5, 0.3, 3.0, 2.5, size=160, bold=True,
    color=rgb("031A25"), align=PP_ALIGN.LEFT)

htxt(sl,
     "AI adalah listrik baru.\nSeperti listrik mengubah segalanya,\nAI akan mengubah setiap industri.",
     1.0, 1.0, 8.5, 2.5, size=34, bold=True,
     color=WHITE, align=PP_ALIGN.LEFT)

rect(sl, 1.0, 3.65, 6.0, 0.05, NG_B)
htxt(sl, "— Andrew Ng\nPendiri DeepLearning.AI, Coursera, Landing AI",
     1.0, 3.8, 9.0, 0.9, size=16, italic=True,
     color=NG_B, align=PP_ALIGN.LEFT)

# left info box
rect(sl, 0.6, 1.2, 4.5, 4.2, rgb("03141E"), line=NG_B)
htxt(sl,
     "🔵  ANDREW NG\n\n"
     "🏫  Stanford University\n"
     "🧠  Pendiri Google Brain\n"
     "🎓  Chief Scientist — Baidu\n"
     "📱  Co-founder Coursera\n"
     "🚀  Founder DeepLearning.AI\n\n"
     "👨‍🎓  8 Juta+ siswa diajar\n"
     "📋  Time100 AI 2023",
     0.8, 1.35, 4.1, 3.9, size=12, color=WHITE)

txt(sl,
    "\"AI jobs are growing, not shrinking.\nPertanyaannya bukan apakah ada pekerjaan — tapi apakah kamu siap mengambilnya.\"",
    1.0, 5.2, 12.0, 1.1, size=14, italic=True,
    color=GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NG: AI PENDIDIKAN
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, NG_B)
rect(sl, 0, 7.44, 13.33, 0.06, NG_B)
rect(sl, 13.08, 0.06, 0.25, 7.38, NG_B)

txt(sl, "🔵  ANDREW NG  —  AI dalam Dunia Pendidikan",
    0.4, 0.18, 12.5, 0.6, size=22, bold=True, color=NG_B)

cards_ng = [
    ("🎓", "DEMOKRATISASI", NG_B,
     "Setiap siswa di desa terpencil bisa\npunya tutor pribadi 24/7 berkualitas\ntinggi — gratis. Ini belum pernah\nterjadi sepanjang sejarah manusia."),
    ("🤖", "AI AGENT GURU", GREEN,
     "Kira Learning (2025): AI agents\nmembantu guru fokus pada hal\nterpenting — relasi, motivasi,\n& bimbingan personal siswa."),
    ("🔓", "BARRIER TURUN", GOLD,
     "8 juta+ orang sudah belajar AI\nlewat Coursera & DeepLearning.AI.\n\"AI for Everyone\" = bukan slogan,\nini misi nyata yang terbukti."),
    ("📐", "REDESIGN PENILAIAN", rgb("B5E3FF"),
     "Masalah plagiarisme? Bukan larang\nAI — tapi rancang ulang cara kita\nmenilai. Uji pemahaman nyata,\nbukan kemampuan menghafal."),
]

for i, (icon, label, col, body) in enumerate(cards_ng):
    x = 0.4 + i * 3.15
    rect(sl, x, 1.0, 2.95, 5.6, CARD_BG, line=col)
    rect(sl, x, 1.0, 2.95, 0.55, rgb("050D10"))
    txt(sl, f"{icon}  {label}", x+0.1, 1.03, 2.75, 0.5,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, body, x+0.15, 1.65, 2.65, 4.7, size=12.5, color=WHITE)

txt(sl,
    "💬  \"AI memungkinkan satu guru luar biasa menjangkau jutaan siswa — bukan menggantikan guru, tapi melipatgandakan dampaknya.\" — Andrew Ng, 2025",
    0.4, 6.7, 12.5, 0.65, size=11.5, italic=True,
    color=NG_B, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NG: AI DUNIA KERJA
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, NG_B)
rect(sl, 0, 7.44, 13.33, 0.06, NG_B)
rect(sl, 13.08, 0.06, 0.25, 7.38, NG_B)

txt(sl, "🔵  ANDREW NG  —  AI dalam Dunia Kerja",
    0.4, 0.18, 12.5, 0.6, size=22, bold=True, color=NG_B)

# 3 big stats
stats = [
    (NG_B,  "97 Juta",  "pekerjaan BARU\nakan tercipta\n(WEF 2025)"),
    (GREEN, "Superpower", "Satu orang + AI =\ndampak setara\nsatu tim besar"),
    (GOLD,  "Vibe Coding", "Non-programmer\nbisa bangun\naplikasi sendiri"),
]
for i, (col, val, sub) in enumerate(stats):
    x = 0.4 + i * 4.3
    rect(sl, x, 0.9, 4.1, 2.3, CARD_BG, line=col)
    txt(sl, val, x+0.1, 0.98, 3.9, 0.95,
        size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, sub, x+0.1, 1.95, 3.9, 1.1,
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

# 3 skill areas
skills = [
    ("⚡", "AGENTIC AI WORKFLOW",
     "AI bekerja iteratif seperti\nmanusia berpikir — membuka\nproduktivitas tak terbatas."),
    ("🎯", "3 KOMPETENSI KUNCI",
     "① AI tools dalam kerja harian\n② Berpikir sistematis\n③ Soft skills tak terotomasi"),
    ("🔁", "MULAI SEKARANG",
     "Jangan tunggu perusahaan\nberubah. Mulai dari dirimu\n— pelajari 1 AI tool hari ini."),
]
for i, (icon, title, body) in enumerate(skills):
    x = 0.4 + i * 4.3
    rect(sl, x, 3.4, 4.1, 2.9, CARD_BG, line=NG_B)
    txt(sl, f"{icon}  {title}", x+0.1, 3.48, 3.9, 0.45,
        size=13, bold=True, color=NG_B)
    rect(sl, x+0.1, 3.95, 3.7, 0.04, rgb("051420"))
    htxt(sl, body, x+0.1, 4.1, 3.85, 2.0, size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DUEL PERSPEKTIF: PENDIDIKAN (Split-screen)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)

# left bg red tint, right blue tint
rect(sl, 0,    0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
# center divider
rect(sl, 6.5, 0, 0.33, 7.5, BG)
txt(sl, "VS", 6.5, 3.35, 0.33, 0.7,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# headers
rect(sl, 0, 0, 6.5, 0.65, rgb("2A0810"))
txt(sl, "🔴  HARARI", 0, 0.1, 6.5, 0.55,
    size=20, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)
rect(sl, 6.83, 0, 6.5, 0.65, rgb("031420"))
txt(sl, "🔵  ANDREW NG", 6.83, 0.1, 6.5, 0.55,
    size=20, bold=True, color=NG_B, align=PP_ALIGN.CENTER)

txt(sl, "🎓  AI & PENDIDIKAN", 0, 0.7, 13.33, 0.45,
    size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

duels = [
    ("Dampak pada\nsiswa",
     "Pola pikir anak dibentuk\nalgorithma tanpa disadari.\nBahaya tersembunyi.",
     "Tutor pribadi untuk\nsemua — mempercepat\npemahaman siswa."),
    ("Berpikir\nkritis",
     "AI mengancam kemampuan\nberpikir mandiri. Proses\nbelajar yang hilang.",
     "Redesign penilaian!\nUji pemahaman nyata,\nbukan hafalan."),
    ("Akses\npendidikan",
     "Peluang tidak merata.\nNegara kaya makin\nmaju, miskin tertinggal.",
     "Demokratisasi!\nDesa terpencil bisa\nakses kualitas terbaik."),
]
for i, (topic, left, right) in enumerate(duels):
    y = 1.25 + i * 1.9
    # topic label
    rect(sl, 0.15, y, 6.2, 1.7, CARD_BG, line=rgb("3A0A10"))
    rect(sl, 6.98, y, 6.2, 1.7, CARD_BG, line=rgb("031E2A"))
    htxt(sl, left,  0.3,  y+0.1, 5.9, 1.5, size=12.5, color=WHITE)
    htxt(sl, right, 7.1,  y+0.1, 5.9, 1.5, size=12.5, color=WHITE)
    txt(sl, topic, 4.7, y+0.55, 3.9, 0.6,
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DUEL PERSPEKTIF: DUNIA KERJA (Split-screen)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0,    0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
rect(sl, 6.5, 0, 0.33, 7.5, BG)
txt(sl, "VS", 6.5, 3.35, 0.33, 0.7,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

rect(sl, 0, 0, 6.5, 0.65, rgb("2A0810"))
txt(sl, "🔴  HARARI", 0, 0.1, 6.5, 0.55,
    size=20, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)
rect(sl, 6.83, 0, 6.5, 0.65, rgb("031420"))
txt(sl, "🔵  ANDREW NG", 6.83, 0.1, 6.5, 0.55,
    size=20, bold=True, color=NG_B, align=PP_ALIGN.CENTER)

txt(sl, "💼  AI & DUNIA KERJA", 0, 0.7, 13.33, 0.45,
    size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

duels_work = [
    ("Nasib\npekerjaan",
     "800 juta pekerjaan berisiko.\nRevolusi AI bisa ciptakan\n\"Useless Class\" baru.",
     "97 juta pekerjaan baru!\nSejarah terbukti: teknologi\nselalu ciptakan lebih banyak."),
    ("Reskilling\n& Adaptasi",
     "\"Reskilling\" = alibi korporasi.\nTanggung jawab bukan di\nindividu, tapi di sistem.",
     "Mulai dari dirimu sekarang!\nJangan tunggu. 1 AI tool\nper hari = karir masa depan."),
    ("Solusi\nyang dibutuhkan",
     "Regulasi global. Tata kelola\nAI demokratis. Jaminan\nsosial bagi yang terdampak.",
     "Pendidikan masif & cepat.\nAI for Everyone. Semua\nbisa, semua harus belajar."),
]
for i, (topic, left, right) in enumerate(duels_work):
    y = 1.25 + i * 1.9
    rect(sl, 0.15, y, 6.2, 1.7, CARD_BG, line=rgb("3A0A10"))
    rect(sl, 6.98, y, 6.2, 1.7, CARD_BG, line=rgb("031E2A"))
    htxt(sl, left,  0.3,  y+0.1, 5.9, 1.5, size=12.5, color=WHITE)
    htxt(sl, right, 7.1,  y+0.1, 5.9, 1.5, size=12.5, color=WHITE)
    txt(sl, topic, 4.7, y+0.55, 3.9, 0.6,
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — POLL INTERAKTIF: SIAPA YANG KAMU SETUJUI?
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, GOLD)
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)

rect(sl, 3.8, 0.18, 5.73, 0.55, rgb("1A1A2E"), line=GOLD)
txt(sl, "🗳️  POLL INTERAKTIF", 3.8, 0.18, 5.73, 0.55,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

htxt(sl, "Setelah mendengar kedua perspektif,\nkamu lebih setuju dengan siapa?",
     0.5, 0.9, 12.33, 1.3, size=32, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "🎯  Ketik di chat: HARARI / NG / KEDUANYA + alasan singkatmu!",
    0.5, 2.3, 12.33, 0.5, size=15, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

# 2 big vote boxes
rect(sl, 0.5, 2.95, 5.7, 3.1, rgb("2A0810"), line=HARARI_R)
txt(sl, "🔴", 0.5, 3.05, 5.7, 0.9,
    size=40, align=PP_ALIGN.CENTER, color=HARARI_R)
txt(sl, "TEAM HARARI", 0.5, 3.95, 5.7, 0.55,
    size=22, bold=True, color=HARARI_R, align=PP_ALIGN.CENTER)
htxt(sl, "AI perlu diawasi ketat.\nManusia harus tetap mengontrol\nnarasi & sistem.", 0.6, 4.55, 5.5, 1.3,
     size=13, color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 7.13, 2.95, 5.7, 3.1, rgb("031420"), line=NG_B)
txt(sl, "🔵", 7.13, 3.05, 5.7, 0.9,
    size=40, align=PP_ALIGN.CENTER, color=NG_B)
txt(sl, "TEAM ANDREW NG", 7.13, 3.95, 5.7, 0.55,
    size=22, bold=True, color=NG_B, align=PP_ALIGN.CENTER)
htxt(sl, "AI adalah peluang besar.\nBelajar & adaptasi adalah\nkunci — mulai sekarang!", 7.23, 4.55, 5.5, 1.3,
     size=13, color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 5.4, 3.15, 2.53, 2.7, CARD_BG, line=GOLD)
txt(sl, "⚡", 5.4, 3.3, 2.53, 0.8,
    size=36, align=PP_ALIGN.CENTER, color=GOLD)
txt(sl, "KEDUANYA", 5.4, 4.15, 2.53, 0.5,
    size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
htxt(sl, "Bisa ancaman\nsekaligus peluang\n— bergantung kita", 5.4, 4.7, 2.53, 1.1,
     size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TITIK TEMU: HARARI + NG SEPAKAT
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, GREEN)
rect(sl, 0, 7.44, 13.33, 0.06, GREEN)

txt(sl, "🤝  TITIK TEMU  —  Harari & Andrew Ng Sepakat",
    0.5, 0.15, 12.33, 0.65, size=24, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)
txt(sl, "Di balik perbedaan tajam, ada hal-hal yang mereka SEPAKATI bersama:",
    0.5, 0.85, 12.33, 0.45, size=14, color=GRAY,
    align=PP_ALIGN.CENTER)

agreements = [
    ("🧠", "Literasi AI adalah WAJIB",
     "Semua orang — guru, dokter, petani, pengusaha — harus\nmemahami AI. Bukan pilihan, ini keharusan era digital."),
    ("⚠️", "AI bisa disalahgunakan",
     "Tanpa regulasi & etika yang kuat, AI bisa menjadi alat\nmanipulasi, ketimpangan, dan hilangnya privasi."),
    ("🌱", "Soft skills makin vital",
     "Kreativitas, empati, kepemimpinan & berpikir kritis adalah\nkompetensi yang TIDAK bisa diotomasi oleh AI."),
    ("🔬", "Pendidikan harus berevolusi",
     "Sistem pendidikan yang mengajarkan hafalan tanpa\npemahaman kritis akan gagal di era AI."),
    ("⚖️", "Manusia tetap harus memimpin",
     "AI adalah alat. Manusia yang harus menentukan\ntujuan, nilai, dan arah penggunaannya."),
]

for i, (icon, title, body) in enumerate(agreements):
    y = 1.4 + i * 1.12
    rect(sl, 0.5, y, 12.33, 1.0, CARD_BG, line=GREEN)
    txt(sl, icon, 0.6, y+0.2, 0.6, 0.6, size=22, align=PP_ALIGN.CENTER, color=GREEN)
    txt(sl, title, 1.35, y+0.08, 4.5, 0.45, size=14, bold=True, color=WHITE)
    htxt(sl, body, 1.35, y+0.52, 11.0, 0.45, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ACTION PLAN: LANGKAH KONKRET UNTUK KAMU
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.06, GOLD)
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)

txt(sl, "🚀  ACTION PLAN  —  Mulai Dari Kamu Hari Ini",
    0.5, 0.15, 12.33, 0.65, size=24, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)

actions = [
    (NG_B,    "1", "HARI INI",
     "Coba satu AI tool dalam kerja/belajarmu:\nChatGPT, Gemini, atau Perplexity."),
    (GREEN,   "2", "MINGGU INI",
     "Daftar 1 kursus AI gratis:\nDeepLearning.AI / Coursera / Google AI."),
    (GOLD,    "3", "BULAN INI",
     "Identifikasi 3 tugas rutinmu\nyang bisa dibantu atau diotomasi AI."),
    (HARARI_R,"4", "SELALU",
     "Latih berpikir kritis: tanyakan\n\"Siapa yang buat ini? Apa tujuannya?\""),
]

for i, (col, num, when, action) in enumerate(actions):
    x = 0.4 + i * 3.15
    rect(sl, x, 1.0, 2.95, 5.6, CARD_BG, line=col)
    rect(sl, x, 1.0, 2.95, 1.1, rgb("0D0D15"))
    txt(sl, num, x+0.1, 1.05, 2.75, 0.55,
        size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, when, x+0.1, 1.6, 2.75, 0.45,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, action, x+0.15, 2.2, 2.65, 4.2, size=13, color=WHITE)

# bottom quote
rect(sl, 0.5, 6.7, 12.33, 0.7, CARD_BG, line=GOLD)
htxt(sl,
     "⚡  \"AI won't replace you. A person using AI will.\"  —  Prinsip yang Harari & Ng sama-sama akui kebenarannya.",
     0.65, 6.78, 12.0, 0.55, size=13, italic=True,
     color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING EPIC
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
rect(sl, 6.5, 0, 0.33, 7.5, BG)
rect(sl, 0, 0, 13.33, 0.06, GOLD)
rect(sl, 0, 7.44, 13.33, 0.06, GOLD)

# Harari final quote
rect(sl, 0.3, 0.3, 5.9, 2.4, rgb("1A0508"), line=HARARI_R)
txt(sl, "\u201c", 0.5, 0.1, 1.5, 1.5, size=60, bold=True,
    color=rgb("3A0A10"))
htxt(sl,
     "Siapa yang mengontrol AI,\nmengontrol masa depan.",
     0.5, 0.45, 5.5, 1.3, size=19, bold=True,
     color=WHITE)
txt(sl, "— Yuval Noah Harari", 0.5, 1.85, 5.5, 0.5,
    size=12, italic=True, color=HARARI_R)

# Ng final quote
rect(sl, 7.13, 0.3, 5.9, 2.4, rgb("031420"), line=NG_B)
txt(sl, "\u201c", 7.3, 0.1, 1.5, 1.5, size=60, bold=True,
    color=rgb("031E2A"))
htxt(sl,
     "Siapa yang belajar AI,\nakan memimpin masa depan.",
     7.3, 0.45, 5.5, 1.3, size=19, bold=True,
     color=WHITE)
txt(sl, "— Andrew Ng", 7.3, 1.85, 5.5, 0.5,
    size=12, italic=True, color=NG_B)

# center closing message
htxt(sl, "PESAN UNTUK KAMU", 0.5, 2.9, 12.33, 0.65,
     size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 3.6, 6.33, 0.05, LINE_COL)

htxt(sl,
     "Jangan hanya jadi penonton AI.\nJadi pemain yang cerdas, kritis, dan terus belajar.",
     0.5, 3.75, 12.33, 1.2, size=20, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)

# 3 closing icons
closes = [
    (HARARI_R, "🧠", "Berpikir\nKritis"),
    (GOLD,     "⚡", "Terus\nBelajar"),
    (NG_B,     "🚀", "Adaptasi\nAktif"),
]
for i, (col, icon, label) in enumerate(closes):
    x = 2.0 + i * 3.3
    rect(sl, x, 5.1, 2.8, 1.8, CARD_BG, line=col)
    txt(sl, icon, x, 5.2, 2.8, 0.8,
        size=32, align=PP_ALIGN.CENTER, color=col)
    txt(sl, label, x, 6.0, 2.8, 0.6,
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(sl, "#WebinarNasionalAI  |  Inovasi · Adaptasi · Kolaborasi  |  2025",
    0.5, 7.1, 12.33, 0.4, size=11, color=GRAY,
    align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save("/projects/sandbox/sonet-product/Webinar_Interaktif_AI.pptx")
print("✅  PPT saved: Webinar_Interaktif_AI.pptx")
print(f"   Total slides: {len(prs.slides)}")

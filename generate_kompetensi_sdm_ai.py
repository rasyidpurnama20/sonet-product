from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── COLOR PALETTE ──────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG      = rgb("0A0A0F")
BLUE    = rgb("00B4D8")
PURPLE  = rgb("7B2FBE")
GOLD    = rgb("FFD166")
WHITE   = rgb("FFFFFF")
GRAY    = rgb("AAAAAA")
GREEN   = rgb("06D6A0")
ORANGE  = rgb("FF6B35")
PINK    = rgb("EF476F")
CARD_BG = rgb("14141E")
LINE    = rgb("2A2A3E")

# ── HELPERS ────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, color, line=None, lw=1.5):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def htxt(slide, text, l, t, w, h, size=18, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)

# gradient accent strips
rect(sl, 0,   0, 13.33, 0.08, BLUE)
rect(sl, 0, 7.42, 13.33, 0.08, BLUE)
rect(sl, 0,   0, 0.12, 7.5, PURPLE)

# glow box center
rect(sl, 1.5, 0.8, 10.33, 5.6, rgb("0D0D1A"), line=BLUE)

# badge
rect(sl, 4.5, 1.1, 4.33, 0.52, rgb("1A1A2E"), line=GOLD)
txt(sl, "🎓  KOMPETENSI SDM ERA AI", 4.5, 1.1, 4.33, 0.52,
    size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# title
htxt(sl, "SDM UNGGUL\nDI ERA KECERDASAN BUATAN", 1.7, 1.75, 9.93, 2.2,
     size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# subtitle
txt(sl, "5 Pilar Kompetensi & Tantangan Pendidikan Indonesia",
    1.7, 4.05, 9.93, 0.55, size=17, color=GOLD, align=PP_ALIGN.CENTER)

# 5 icons for pilar
pillars = [("🔤","Foundational"), ("💻","Digital-AI"), ("🧠","Cognitive"), ("🤝","Human"), ("⚖️","Adaptive")]
for i, (icon, label) in enumerate(pillars):
    x = 2.0 + i * 1.9
    rect(sl, x, 4.8, 1.6, 1.1, rgb("1A1A2E"), line=LINE)
    txt(sl, icon, x, 4.85, 1.6, 0.5, size=20, align=PP_ALIGN.CENTER, color=WHITE)
    txt(sl, label, x, 5.38, 1.6, 0.45, size=9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

txt(sl, "2025  |  Inovasi · Adaptasi · Kolaborasi",
    0, 7.1, 13.33, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 5 PILAR KOMPETENSI (Overview)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, GOLD)
rect(sl, 0, 7.42, 13.33, 0.08, GOLD)

txt(sl, "5 PILAR KOMPETENSI SDM ERA AI", 0.4, 0.14, 12.5, 0.65,
    size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "Kompetensi yang harus dimiliki setiap individu untuk bertahan & berkembang di era AI",
    0.4, 0.8, 12.5, 0.4, size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

pilar_data = [
    (BLUE,   "01", "🔤 FOUNDATIONAL",   "Literasi · Numerasi\nSains · Bahasa",          "Fondasi belajar semua hal"),
    (GREEN,  "02", "💻 DIGITAL-AI",     "AI Literacy · Data\nPrompting · Cybersecurity", "Navigasi dunia digital"),
    (PURPLE, "03", "🧠 COGNITIVE",      "Kritis · Analitis\nKreatif · Problem Solving",  "Berpikir di level tinggi"),
    (ORANGE, "04", "🤝 HUMAN",          "Komunikasi · Empati\nKolaborasi · Kepemimpinan","Kelebihan manusia atas AI"),
    (PINK,   "05", "⚖️ ADAPTIVE-ETHIC", "Lifelong Learning\nEtika AI · Resiliensi",      "Bertumbuh & bertanggung jawab"),
]

for i, (col, num, title, skills, tagline) in enumerate(pilar_data):
    x = 0.3 + i * 2.58
    rect(sl, x, 1.35, 2.42, 5.7, CARD_BG, line=col)
    rect(sl, x, 1.35, 2.42, 0.65, rgb("0D0D15"))
    txt(sl, num, x+0.08, 1.38, 0.5, 0.58, size=11, bold=True, color=col)
    txt(sl, title, x+0.12, 1.38, 2.2, 0.58, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, skills, x+0.12, 2.15, 2.18, 1.5, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, x+0.15, 3.75, 2.12, 0.04, LINE)
    htxt(sl, tagline, x+0.12, 3.9, 2.18, 0.55, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    # bottom accent bar
    rect(sl, x, 6.95, 2.42, 0.1, col)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DIGITAL-AI SKILLS (Detail)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, BLUE)
rect(sl, 0, 7.42, 13.33, 0.08, BLUE)
rect(sl, 0, 0.08, 0.12, 7.34, BLUE)

txt(sl, "💻  DIGITAL-AI SKILLS  —  Fondasi Era Baru",
    0.3, 0.15, 12.5, 0.6, size=22, bold=True, color=BLUE)
rect(sl, 0.3, 0.8, 12.5, 0.04, LINE)

skills_ai = [
    (BLUE,   "🤖", "AI LITERACY",
     "Siswa memahami cara kerja,\nbatasan, dan risiko AI.\nBukan sekadar pengguna pasif."),
    (GREEN,  "📊", "DATA LITERACY",
     "Mampu membaca, menafsirkan,\ndan memvalidasi data.\nTidak mudah dibohongi statistik."),
    (GOLD,   "💬", "PROMPTING",
     "Seni berkomunikasi dengan AI\nuntuk hasil optimal.\nKualitas output = kualitas input."),
    (PINK,   "🔒", "CYBERSECURITY",
     "Menjaga identitas, privasi,\ndan keamanan digital.\nSadar ancaman siber sehari-hari."),
    (ORANGE, "⚙️", "COMPUTATIONAL\nTHINKING",
     "Memecah masalah kompleks\nmenjadi langkah logis.\nBerpikir seperti programmer."),
]

for i, (col, icon, title, body) in enumerate(skills_ai):
    x = 0.3 + i * 2.58
    rect(sl, x, 0.98, 2.42, 5.9, CARD_BG, line=col)
    txt(sl, icon, x+0.05, 1.05, 2.32, 0.75, size=28, align=PP_ALIGN.CENTER, color=col)
    htxt(sl, title, x+0.1, 1.85, 2.22, 0.7, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    rect(sl, x+0.15, 2.6, 2.12, 0.04, LINE)
    htxt(sl, body, x+0.12, 2.75, 2.18, 3.9, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom insight box
rect(sl, 0.3, 7.0, 12.5, 0.38, rgb("0D1A20"), line=BLUE)
txt(sl, "💡  AI bisa menghasilkan banyak jawaban → Manusia harus memilih mana yang BENAR, BERGUNA, ETIS, dan sesuai KONTEKS",
    0.45, 7.03, 12.2, 0.35, size=11, italic=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — COGNITIVE SKILLS (Detail)
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, PURPLE)
rect(sl, 0, 7.42, 13.33, 0.08, PURPLE)
rect(sl, 0, 0.08, 0.12, 7.34, PURPLE)

txt(sl, "🧠  COGNITIVE SKILLS  —  Otak yang Tak Bisa Digantikan AI",
    0.3, 0.15, 12.5, 0.6, size=22, bold=True, color=PURPLE)

# Left: 5 cognitive skills cards (2 columns)
cog_skills = [
    (PURPLE, "🔍", "Critical Thinking", "Mengevaluasi kualitas\njawaban & sumber"),
    (BLUE,   "📈", "Analytical Thinking", "Membaca pola\ndan sebab-akibat"),
    (GREEN,  "💡", "Creative Thinking", "Menghasilkan ide baru\nyang AI belum bisa"),
    (GOLD,   "🔧", "Problem Solving", "Merancang solusi,\nbukan sekadar menjawab"),
    (ORANGE, "🌐", "Systems Thinking", "Memahami dampak\nkeputusan secara luas"),
]

for i, (col, icon, title, body) in enumerate(cog_skills):
    col_x = 0.35 if i % 2 == 0 else 3.05
    row_y = 0.95 + (i // 2) * 2.0
    if i == 4:
        col_x = 1.7
        row_y = 4.95
    rect(sl, col_x, row_y, 2.55, 1.75, CARD_BG, line=col)
    txt(sl, f"{icon}  {title}", col_x+0.12, row_y+0.12, 2.3, 0.45, size=13, bold=True, color=col)
    htxt(sl, body, col_x+0.12, row_y+0.62, 2.3, 1.0, size=11.5, color=WHITE)

# Right: why it matters
rect(sl, 5.85, 0.88, 7.1, 6.0, rgb("0F0F1E"), line=PURPLE)
txt(sl, "⚡  MENGAPA INI PENTING?", 6.0, 0.98, 6.8, 0.5,
    size=15, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
rect(sl, 6.1, 1.52, 6.6, 0.04, LINE)

insights = [
    "🤖  AI bisa jawab soal ujian lebih cepat dari siswamu.",
    "📌  Tapi AI tidak bisa MEMAHAMI MAKNA di balik jawaban.",
    "🎯  Semakin mudah mencari jawaban lewat AI...",
    "🧩  Semakin penting kemampuan merumuskan pertanyaan.",
    "📚  Pendidikan kita masih terlalu fokus pada HAFALAN.",
    "🚨  Di era AI: hafalan tidak bernilai. Pemahaman segalanya.",
    "✅  Cognitive skills = senjata utama manusia vs AI.",
]
for i, line in enumerate(insights):
    y = 1.65 + i * 0.71
    rect(sl, 6.0, y, 6.7, 0.62, rgb("14141E"), line=LINE)
    txt(sl, line, 6.12, y+0.1, 6.45, 0.45, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HUMAN SKILLS: Keunggulan Manusia
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, ORANGE)
rect(sl, 0, 7.42, 13.33, 0.08, ORANGE)
rect(sl, 13.21, 0.08, 0.12, 7.34, ORANGE)

txt(sl, "🤝  HUMAN SKILLS  —  Yang Tak Bisa Diotomasi",
    0.3, 0.15, 12.5, 0.6, size=22, bold=True, color=ORANGE)

# Flow: semakin kuat AI → semakin penting human skills
rect(sl, 0.3, 0.88, 6.0, 5.9, rgb("0D0D15"), line=ORANGE)
txt(sl, "📢  TREN GLOBAL", 0.45, 0.97, 5.7, 0.45, size=14, bold=True, color=ORANGE)
rect(sl, 0.45, 1.45, 5.7, 0.04, LINE)

flow = [
    (ORANGE, "Semakin kuat AI dalam tugas teknis"),
    (GOLD,   "           ↓"),
    (GOLD,   "Semakin PENTING keunggulan manusia"),
    (GRAY,   "           ↓"),
    (WHITE,  "Empati · Komunikasi · Nilai · Kerja sama"),
    (GREEN,  "           ↓"),
    (GREEN,  "Menjadi PEMBEDA UTAMA di pasar kerja"),
]
for i, (col, line) in enumerate(flow):
    y = 1.6 + i * 0.63
    htxt(sl, line, 0.45, y, 5.7, 0.58, size=13, bold=(col != GRAY), color=col, align=PP_ALIGN.CENTER)

# Right: 4 human skill cards
human_skills = [
    (ORANGE, "💬", "KOMUNIKASI",    "Menyampaikan ide kompleks\nsecara jelas & persuasif."),
    (PINK,   "❤️", "EMPATI",        "Memahami perasaan orang lain.\nAI tidak bisa benar-benar peduli."),
    (BLUE,   "👥", "KOLABORASI",    "Bekerja lintas tim, budaya,\ndan disiplin ilmu."),
    (GREEN,  "🌟", "KEPEMIMPINAN",  "Menginspirasi, memotivasi,\ndan membuat keputusan sulit."),
]
for i, (col, icon, title, body) in enumerate(human_skills):
    cx = 6.6 + (i % 2) * 3.25
    ry = 0.9 + (i // 2) * 2.95
    rect(sl, cx, ry, 3.0, 2.65, CARD_BG, line=col)
    txt(sl, icon, cx+0.1, ry+0.15, 2.8, 0.75, size=30, align=PP_ALIGN.CENTER, color=col)
    txt(sl, title, cx+0.1, ry+0.95, 2.8, 0.45, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    rect(sl, cx+0.2, ry+1.42, 2.6, 0.04, LINE)
    htxt(sl, body, cx+0.12, ry+1.55, 2.76, 0.9, size=11.5, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ADAPTIVE-ETHICAL SKILLS + Lifelong Learning
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, PINK)
rect(sl, 0, 7.42, 13.33, 0.08, PINK)
rect(sl, 0, 0.08, 0.12, 7.34, PINK)

txt(sl, "⚖️  ADAPTIVE-ETHICAL SKILLS  —  Bertumbuh & Bertanggung Jawab",
    0.3, 0.15, 12.5, 0.6, size=20, bold=True, color=PINK)

# Left: Lifelong Learning traits
rect(sl, 0.3, 0.9, 5.6, 5.95, rgb("100810"), line=PINK)
txt(sl, "🌱  LIFELONG LEARNING TRAITS", 0.45, 0.98, 5.3, 0.48,
    size=13, bold=True, color=PINK, align=PP_ALIGN.CENTER)
rect(sl, 0.45, 1.5, 5.3, 0.04, LINE)

traits = [
    ("🔍", "Rasa ingin tahu yang terus tumbuh"),
    ("📖", "Kemampuan belajar mandiri"),
    ("🧭", "Mencari sumber yang kredibel"),
    ("🔄", "Menerima umpan balik dengan terbuka"),
    ("🎯", "Kemampuan adaptif terhadap perubahan"),
    ("💡", "Menyadari bahwa skill kerja berubah cepat"),
    ("🔐", "Memahami privasi data & nilai etika"),
    ("🤖", "Memahami bias algoritma"),
]
for i, (icon, trait) in enumerate(traits):
    y = 1.63 + i * 0.67
    rect(sl, 0.45, y, 5.3, 0.58, rgb("160D1A"), line=rgb("3A1A3E"))
    txt(sl, f"{icon}  {trait}", 0.6, y+0.09, 5.0, 0.42, size=12, color=WHITE)

# Right: Etika AI flow
rect(sl, 6.15, 0.9, 6.85, 5.95, rgb("0D0A10"), line=PINK)
txt(sl, "🚨  MENGAPA ETIKA AI KRUSIAL?", 6.3, 0.98, 6.55, 0.48,
    size=13, bold=True, color=PINK, align=PP_ALIGN.CENTER)
rect(sl, 6.3, 1.5, 6.55, 0.04, LINE)

htxt(sl, "AI adalah alat kuat\n         ↓\nAlat kuat memerlukan tanggung jawab\n         ↓\nTanpa etika, AI bisa:",
     6.3, 1.6, 6.55, 2.0, size=12.5, color=WHITE, align=PP_ALIGN.CENTER)

dangers = [
    (PINK,   "⚠️", "Memperkuat BIAS"),
    (ORANGE, "🎭", "Alat MANIPULASI"),
    (PURPLE, "📉", "Perkuat KETIMPANGAN"),
    (GRAY,   "🔗", "Ciptakan KETERGANTUNGAN"),
]
for i, (col, icon, label) in enumerate(dangers):
    cx = 6.3 + (i % 2) * 3.3
    ry = 3.75 + (i // 2) * 1.5
    rect(sl, cx, ry, 3.1, 1.25, CARD_BG, line=col)
    txt(sl, f"{icon}  {label}", cx+0.1, ry+0.35, 2.9, 0.5, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TANTANGAN PENDIDIKAN INDONESIA
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, GOLD)
rect(sl, 0, 7.42, 13.33, 0.08, GOLD)

txt(sl, "🚨  TANTANGAN PENDIDIKAN KITA  —  Gap yang Harus Dijembatani",
    0.3, 0.14, 12.5, 0.62, size=21, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

challenges = [
    (PINK,   "📚", "KURIKULUM\nKETINGGALAN",
     "Kurikulum masih berbasis hafalan.\nBelum sepenuhnya melatih berpikir kritis,\nkreatif, dan digital-AI skills."),
    (ORANGE, "🌍", "KESENJANGAN\nAKSES",
     "Infrastruktur digital tidak merata.\nSekolah di desa vs kota punya\nkualitas akses AI yang jauh berbeda."),
    (PURPLE, "👩‍🏫", "KOMPETENSI\nGURU",
     "Banyak guru belum terpapar AI literacy.\nSulit mengajarkan sesuatu yang\nbelum dipahami sendiri."),
    (BLUE,   "📏", "SISTEM\nPENILAIAN",
     "Ujian masih mengukur ingatan,\nbukan pemahaman & penerapan.\nMudah diakali oleh AI generatif."),
    (GREEN,  "⚠️", "KETERGANTUNGAN\nAI TANPA KRITIS",
     "Siswa pakai AI untuk menjawab\ntugas tanpa berpikir.\nProses belajar yang hilang."),
    (GOLD,   "🔗", "KONEKSI\nINDUSTRI",
     "Dunia pendidikan & industri\nbelum terkoneksi baik.\nSkill yang diajarkan vs dibutuhkan tidak sync."),
]

for i, (col, icon, title, body) in enumerate(challenges):
    cx = 0.3 + (i % 3) * 4.35
    ry = 0.95 + (i // 3) * 3.05
    rect(sl, cx, ry, 4.1, 2.75, CARD_BG, line=col)
    rect(sl, cx, ry, 4.1, 0.6, rgb("0D0D15"))
    txt(sl, f"{icon}  {title}", cx+0.12, ry+0.08, 3.86, 0.5, size=12, bold=True, color=col)
    htxt(sl, body, cx+0.15, ry+0.72, 3.8, 1.9, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SOLUSI & CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 13.33, 0.08, GREEN)
rect(sl, 0, 7.42, 13.33, 0.08, GREEN)

txt(sl, "🚀  SOLUSI  —  Langkah Konkret Transformasi Pendidikan",
    0.3, 0.14, 12.5, 0.62, size=21, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

solutions = [
    (GREEN,  "🔄", "REDESIGN KURIKULUM",
     "Integrasikan AI literacy & digital skills\ndalam setiap mata pelajaran.\nBukan kelas baru, tapi cara baru mengajar."),
    (BLUE,   "👨‍🏫", "UPSKILLING GURU",
     "Program pelatihan AI masif untuk guru.\nGuru yang paham AI = siswa yang\nsiap masa depan."),
    (PURPLE, "📐", "UBAH CARA PENILAIAN",
     "Nilai proses, bukan hanya hasil.\nUji pemahaman nyata & kemampuan\nmemecahkan masalah baru."),
    (GOLD,   "🌐", "PEMERATAAN AKSES",
     "Infrastruktur digital di seluruh Indonesia.\nSetiap siswa berhak mendapat\nkualitas pendidikan yang sama."),
]

for i, (col, icon, title, body) in enumerate(solutions):
    cx = 0.3 + i * 3.18
    rect(sl, cx, 0.92, 3.02, 4.6, CARD_BG, line=col)
    txt(sl, icon, cx+0.1, 1.0, 2.82, 0.75, size=30, align=PP_ALIGN.CENTER, color=col)
    txt(sl, title, cx+0.1, 1.82, 2.82, 0.55, size=12.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    rect(sl, cx+0.2, 2.42, 2.62, 0.04, LINE)
    htxt(sl, body, cx+0.12, 2.55, 2.78, 2.8, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Key insight quote
rect(sl, 0.3, 5.72, 12.5, 1.05, rgb("050F0A"), line=GREEN)
htxt(sl,
     "💡  \"AI menguntungkan orang yang mau belajar.\n"
     "Yang tertinggal bukan yang pekerjaannya diganti — tapi yang berhenti belajar.\"",
     0.5, 5.8, 12.1, 0.9, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING: PESAN AKHIR
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)

# Split color halves
rect(sl, 0,    0, 6.66, 7.5, rgb("050510"))
rect(sl, 6.66, 0, 6.67, 7.5, rgb("05100A"))

rect(sl, 0, 0, 13.33, 0.08, GOLD)
rect(sl, 0, 7.42, 13.33, 0.08, GOLD)
rect(sl, 6.6, 0, 0.13, 7.5, LINE)

# Left: tantangan header
txt(sl, "⚠️  TANTANGAN", 0.4, 0.28, 6.0, 0.55,
    size=20, bold=True, color=PINK, align=PP_ALIGN.CENTER)
challenge_pts = [
    "Pendidikan masih berbasis hafalan",
    "Guru belum siap mengajar literasi AI",
    "Akses digital tidak merata",
    "Sistem penilaian ketinggalan zaman",
    "Siswa gunakan AI tanpa berpikir kritis",
]
for i, pt in enumerate(challenge_pts):
    rect(sl, 0.4, 0.98 + i * 0.9, 6.0, 0.75, rgb("150508"), line=rgb("3A0A10"))
    txt(sl, f"🔴  {pt}", 0.55, 1.06 + i * 0.9, 5.7, 0.6, size=12.5, color=WHITE)

# Right: peluang header
txt(sl, "✅  PELUANG & SOLUSI", 6.88, 0.28, 6.0, 0.55,
    size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
solution_pts = [
    "Redesign kurikulum berbasis kompetensi",
    "Upskilling masif untuk guru Indonesia",
    "Demokratisasi akses teknologi AI",
    "Ubah penilaian: proses & pemahaman",
    "Bangun generasi kritis, adaptif, etis",
]
for i, pt in enumerate(solution_pts):
    rect(sl, 6.88, 0.98 + i * 0.9, 6.0, 0.75, rgb("05150A"), line=rgb("0A3A14"))
    txt(sl, f"🟢  {pt}", 7.03, 1.06 + i * 0.9, 5.7, 0.6, size=12.5, color=WHITE)

# Center closing message
rect(sl, 1.5, 5.55, 10.33, 1.55, rgb("0D0D0D"), line=GOLD)
htxt(sl,
     "🌟  SDM unggul era AI bukan yang paling pintar menggunakan AI\n"
     "— tapi yang paling bijak MEMILIH, MEMILAH, dan BERTANGGUNG JAWAB atas penggunaannya.",
     1.65, 5.65, 10.03, 1.35, size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save("/projects/sandbox/sonet-product/Kompetensi_SDM_Era_AI.pptx")
print("✅  PPT saved: Kompetensi_SDM_Era_AI.pptx")

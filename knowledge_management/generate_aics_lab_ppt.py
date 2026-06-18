from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── COLOUR PALETTE ───────────────────────────────────────────────────────────
NAVY       = RGBColor(0x05, 0x10, 0x28)   # deep navy background
DARK_BLUE  = RGBColor(0x0A, 0x1F, 0x44)   # card / panel
MID_BLUE   = RGBColor(0x0D, 0x2E, 0x6B)   # secondary panel
CYAN       = RGBColor(0x00, 0xC8, 0xFF)   # primary accent
TEAL       = RGBColor(0x00, 0xE5, 0xB0)   # secondary accent
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)   # warning / highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE8)
DARK_GREY  = RGBColor(0x1E, 0x35, 0x5A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ─── HELPERS ──────────────────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, alpha=None):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_line(slide, t, color=CYAN, l=Inches(0.5), w=None):
    """Horizontal accent bar."""
    if w is None:
        w = W - Inches(1)
    r = rect(slide, l, t, w, Inches(0.04), color)
    return r

def divider_line(slide, t, color=DARK_GREY):
    accent_line(slide, t, color=color, l=Inches(0.5), w=W - Inches(1))

def section_tag(slide, label, l=Inches(0.5), t=Inches(0.18)):
    """Small coloured pill label at top-left."""
    pill = rect(slide, l, t, Inches(2.6), Inches(0.32), CYAN)
    txt(slide, label, l + Inches(0.08), t + Inches(0.02),
        Inches(2.4), Inches(0.3), size=9, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h, fill=DARK_BLUE):
    """Dark panel card."""
    r = rect(slide, l, t, w, h, fill)
    return r

def bullet_items(slide, items, l, t, w, h, size=13, color=LIGHT_GREY,
                 bullet="▸", spacing=0.38):
    """Render a list of bullet strings in a textbox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
bg(s1, NAVY)

# Left dark panel
rect(s1, 0, 0, Inches(7.2), H, DARK_BLUE)

# Top cyan stripe
rect(s1, 0, 0, Inches(7.2), Inches(0.12), CYAN)

# Bottom teal stripe
rect(s1, 0, H - Inches(0.12), Inches(7.2), Inches(0.12), TEAL)

# Decorative grid dots (right side atmosphere)
for row in range(8):
    for col in range(9):
        dot = rect(s1,
                   Inches(7.5) + col * Inches(0.62),
                   Inches(0.6) + row * Inches(0.85),
                   Inches(0.1), Inches(0.1),
                   CYAN if (row + col) % 3 == 0 else DARK_GREY)

# Shield icon placeholder
shield = rect(s1, Inches(8.5), Inches(2.0), Inches(2.2), Inches(2.6), MID_BLUE)
txt(s1, "🛡", Inches(8.5), Inches(2.15), Inches(2.2), Inches(2.4),
    size=72, align=PP_ALIGN.CENTER, color=CYAN)

# Tag line on right
txt(s1, "AI & CYBERSECURITY", Inches(7.4), Inches(4.9), Inches(5.5), Inches(0.5),
    size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s1, "LABORATORY", Inches(7.4), Inches(5.3), Inches(5.5), Inches(0.5),
    size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Main title
txt(s1, "PARTNERSHIP PROPOSAL", Inches(0.4), Inches(1.0), Inches(6.5), Inches(0.6),
    size=13, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
txt(s1, "Development of the\nArtificial Intelligence &\nCybersecurity Laboratory",
    Inches(0.4), Inches(1.55), Inches(6.5), Inches(2.0),
    size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s1, "(AICS-Lab UNDIP)", Inches(0.4), Inches(3.45), Inches(6.5), Inches(0.5),
    size=18, bold=False, color=TEAL, align=PP_ALIGN.LEFT)

# Divider
accent_line(s1, Inches(4.05), color=CYAN, l=Inches(0.4), w=Inches(6.4))

# Partners
txt(s1, "DIPONEGORO UNIVERSITY  ×  POSITIVE TECHNOLOGIES",
    Inches(0.4), Inches(4.2), Inches(6.5), Inches(0.45),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s1, "Faculty of Science and Mathematics  |  Department of Informatics",
    Inches(0.4), Inches(4.6), Inches(6.5), Inches(0.35),
    size=10, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
txt(s1, "7th Floor, FSM Central Laboratory Building, Semarang  |  2026",
    Inches(0.4), Inches(4.92), Inches(6.5), Inches(0.35),
    size=10, color=LIGHT_GREY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
bg(s2, NAVY)
rect(s2, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s2, Inches(1.1), color=CYAN)

txt(s2, "AGENDA", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)
txt(s2, "AI & Cybersecurity Laboratory – UNDIP × Positive Technologies",
    Inches(0.5), Inches(0.72), Inches(10), Inches(0.35),
    size=11, color=CYAN)

agenda_items = [
    ("01", "Executive Summary"),
    ("02", "Background & Rationale"),
    ("03", "The Parties & Partnership Scheme"),
    ("04", "Laboratory Concept & Functional Zones"),
    ("05", "Technical Specifications & Equipment"),
    ("06", "Academic Programmes & Human Capital"),
    ("07", "Budget & Investment"),
    ("08", "Governance & KPIs"),
    ("09", "Implementation Timeline"),
    ("10", "Risk Management"),
    ("11", "Closing & Call to Action"),
]

cols = [agenda_items[:6], agenda_items[6:]]
col_x = [Inches(0.5), Inches(6.85)]
for ci, col in enumerate(cols):
    for ri, (num, label) in enumerate(col):
        cy = Inches(1.35) + ri * Inches(0.97)
        cx = col_x[ci]
        card(s2, cx, cy, Inches(6.0), Inches(0.82), fill=DARK_BLUE)
        rect(s2, cx, cy, Inches(0.08), Inches(0.82), CYAN)
        txt(s2, num, cx + Inches(0.18), cy + Inches(0.12),
            Inches(0.55), Inches(0.6), size=20, bold=True, color=CYAN)
        txt(s2, label, cx + Inches(0.78), cy + Inches(0.18),
            Inches(5.0), Inches(0.5), size=13, bold=False, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
bg(s3, NAVY)
rect(s3, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s3, Inches(1.1), color=CYAN)
section_tag(s3, "01  EXECUTIVE SUMMARY")
txt(s3, "Executive Summary", Inches(0.5), Inches(0.22), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

kpi_boxes = [
    ("IDR 45.1B", "Total CAPEX",        CYAN),
    ("IDR 5.8B/yr","Annual OPEX",       TEAL),
    ("± 977 m²",   "Lab Floor Area",    ORANGE),
    ("Q1 2028",    "Official Launch",   CYAN),
    ("S2 / 39 SKS","Master's Programme",TEAL),
    ("5 Zones",    "Functional Areas",  ORANGE),
]
for i, (val, lbl, col) in enumerate(kpi_boxes):
    cx = Inches(0.45) + i * Inches(2.12)
    card(s3, cx, Inches(1.22), Inches(2.0), Inches(1.0), fill=DARK_BLUE)
    rect(s3, cx, Inches(1.22), Inches(2.0), Inches(0.07), col)
    txt(s3, val, cx, Inches(1.32), Inches(2.0), Inches(0.5),
        size=17, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s3, lbl, cx, Inches(1.78), Inches(2.0), Inches(0.38),
        size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Summary text
summary = (
    "The AICS-Lab UNDIP will be established at the 7th Floor of the FSM Central Laboratory "
    "Building, Diponegoro University, as a national center of excellence integrating AI and "
    "cybersecurity education, research, innovation, and professional services.\n\n"
    "The initiative addresses Indonesia's structural digital-talent gap (9M talents needed by 2030) "
    "and 403M+ cyberattacks recorded by BSSN in 2023. A strategic partnership with "
    "Positive Technologies (Russia) brings world-class tools: MaxPatrol SIEM/VM, PT NAD, "
    "PT Sandbox, PT AF, and the Standoff 365 global cyber-range.\n\n"
    "Construction begins December 2026; Official Launch Q1 2028."
)
card(s3, Inches(0.45), Inches(2.38), Inches(12.4), Inches(2.7), fill=DARK_BLUE)
rect(s3, Inches(0.45), Inches(2.38), Inches(0.07), Inches(2.7), TEAL)
txt(s3, summary, Inches(0.65), Inches(2.5), Inches(12.0), Inches(2.5),
    size=12, color=LIGHT_GREY, wrap=True)

highlights = [
    "Cooperation staged: LoI → MoU → MoA → Implementation Arrangement",
    "Funding: UNDIP ±55% | Positive Technologies ±20% (in-kind) | Grants/CSR ±25%",
    "DED completed June 2026  |  Site preparation commenced 17 June 2026",
]
bullet_items(s3, highlights, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.5),
             size=11, color=CYAN, bullet="✔")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – BACKGROUND & RATIONALE
# ══════════════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
bg(s4, NAVY)
rect(s4, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s4, Inches(1.1), color=CYAN)
section_tag(s4, "02  BACKGROUND & RATIONALE")
txt(s4, "Why This Lab — Why Now?", Inches(0.5), Inches(0.22), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

stats = [
    ("USD 366B", "Indonesia's projected\ndigital economy by 2030", CYAN, "📈"),
    ("403M+",    "Cyber-attack traffic anomalies\nrecorded by BSSN in 2023",    ORANGE, "⚠"),
    ("9 Million","Digital talents needed\nby 2030 — supply far behind",         TEAL,   "👥"),
]
for i, (val, lbl, col, icon) in enumerate(stats):
    cx = Inches(0.45) + i * Inches(4.25)
    card(s4, cx, Inches(1.28), Inches(4.0), Inches(1.7), fill=DARK_BLUE)
    rect(s4, cx, Inches(1.28), Inches(4.0), Inches(0.08), col)
    txt(s4, icon, cx + Inches(0.15), Inches(1.4), Inches(0.7), Inches(0.7),
        size=22, color=col, align=PP_ALIGN.CENTER)
    txt(s4, val, cx + Inches(0.9), Inches(1.38), Inches(2.9), Inches(0.55),
        size=24, bold=True, color=col)
    txt(s4, lbl, cx + Inches(0.15), Inches(1.88), Inches(3.7), Inches(0.9),
        size=11, color=LIGHT_GREY, wrap=True)

# Two-column rationale
left_points = [
    "AI & cybersecurity are decisive determinants of national competitiveness & sovereignty",
    "UNDIP as PTN-BH has the mandate and autonomy to respond strategically",
    "7th floor of FSM Central Lab Building is ideally suited for AI/cyber infrastructure",
]
right_points = [
    "Persistent shortage of certified AI, data science & cybersecurity professionals",
    "Escalating threats to vital information infrastructure require immediate action",
    "Strategic moment: Physical infrastructure ready, partner aligned, DED complete",
]
card(s4, Inches(0.45), Inches(3.15), Inches(6.0), Inches(3.5), fill=DARK_BLUE)
rect(s4, Inches(0.45), Inches(3.15), Inches(0.07), Inches(3.5), CYAN)
txt(s4, "National Imperative", Inches(0.65), Inches(3.22), Inches(5.6), Inches(0.4),
    size=12, bold=True, color=CYAN)
bullet_items(s4, left_points, Inches(0.65), Inches(3.62), Inches(5.6), Inches(2.8), size=11)

card(s4, Inches(6.85), Inches(3.15), Inches(6.0), Inches(3.5), fill=DARK_BLUE)
rect(s4, Inches(6.85), Inches(3.15), Inches(0.07), Inches(3.5), TEAL)
txt(s4, "Structural Opportunity", Inches(7.05), Inches(3.22), Inches(5.6), Inches(0.4),
    size=12, bold=True, color=TEAL)
bullet_items(s4, right_points, Inches(7.05), Inches(3.62), Inches(5.6), Inches(2.8), size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – PARTIES & PARTNERSHIP SCHEME
# ══════════════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
bg(s5, NAVY)
rect(s5, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s5, Inches(1.1), color=CYAN)
section_tag(s5, "03  THE PARTIES & PARTNERSHIP SCHEME")
txt(s5, "Strategic Collaboration", Inches(0.5), Inches(0.22), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

# UNDIP card
card(s5, Inches(0.45), Inches(1.25), Inches(5.8), Inches(2.7), fill=DARK_BLUE)
rect(s5, Inches(0.45), Inches(1.25), Inches(5.8), Inches(0.08), CYAN)
txt(s5, "🏛  DIPONEGORO UNIVERSITY (UNDIP)",
    Inches(0.6), Inches(1.38), Inches(5.5), Inches(0.45),
    size=13, bold=True, color=CYAN)
undip_pts = [
    "Leading PTN-BH state university, Semarang, Central Java",
    "Faculty of Science & Mathematics (FSM) — Dept. of Informatics",
    "Academic foundation: intelligent computing, data analytics & cybersecurity",
    "Land & building (7th floor), civil works, HR & operating costs  →  ±55% CAPEX",
]
bullet_items(s5, undip_pts, Inches(0.6), Inches(1.85), Inches(5.5), Inches(1.9), size=11)

# PosTech card
card(s5, Inches(6.55), Inches(1.25), Inches(6.3), Inches(2.7), fill=DARK_BLUE)
rect(s5, Inches(6.55), Inches(1.25), Inches(6.3), Inches(0.08), TEAL)
txt(s5, "🔐  POSITIVE TECHNOLOGIES (Russian Federation)",
    Inches(6.7), Inches(1.38), Inches(6.0), Inches(0.45),
    size=13, bold=True, color=TEAL)
pt_pts = [
    "World-class cybersecurity company",
    "Portfolio: MaxPatrol SIEM, MaxPatrol VM, PT NAD, PT Sandbox, PT AF, PT ISIM",
    "Standoff 365 — global cyber-range platform & active education initiatives",
    "Solution licenses, curriculum, ToT & certification  →  ±20% CAPEX (in-kind)",
]
bullet_items(s5, pt_pts, Inches(6.7), Inches(1.85), Inches(6.0), Inches(1.9), size=11)

# Cooperation scheme flow
txt(s5, "Cooperation Scheme  (staged, equal-partnership)", Inches(0.5), Inches(4.1),
    Inches(12), Inches(0.4), size=12, bold=True, color=WHITE)

stages = [("LoI", "Letter\nof Intent", DARK_GREY),
          ("MoU", "Memorandum\nof Understanding", MID_BLUE),
          ("MoA", "Memorandum\nof Agreement", DARK_BLUE),
          ("IA",  "Implementation\nArrangement", DARK_BLUE)]
arrow_col = CYAN
for i, (abbr, desc, fill_c) in enumerate(stages):
    cx = Inches(0.45) + i * Inches(3.15)
    card(s5, cx, Inches(4.6), Inches(2.7), Inches(1.55), fill=fill_c)
    rect(s5, cx, Inches(4.6), Inches(2.7), Inches(0.07),
         CYAN if i < 2 else TEAL)
    txt(s5, abbr, cx, Inches(4.72), Inches(2.7), Inches(0.55),
        size=20, bold=True, color=CYAN if i < 2 else TEAL, align=PP_ALIGN.CENTER)
    txt(s5, desc, cx, Inches(5.22), Inches(2.7), Inches(0.85),
        size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s5, "→", Inches(0.45) + i * Inches(3.15) + Inches(2.7),
            Inches(5.0), Inches(0.38), Inches(0.4),
            size=18, bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – LAB CONCEPT & FUNCTIONAL ZONES
# ══════════════════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
bg(s6, NAVY)
rect(s6, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s6, Inches(1.1), color=CYAN)
section_tag(s6, "04  LABORATORY CONCEPT")
txt(s6, "Lab Vision & 5 Functional Zones — 7th Floor (±977 m²)",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

# Vision strip
card(s6, Inches(0.45), Inches(1.18), Inches(12.4), Inches(0.68), fill=MID_BLUE)
rect(s6, Inches(0.45), Inches(1.18), Inches(0.07), Inches(0.68), CYAN)
txt(s6,
    "Vision: A national & regional center of excellence in AI and cybersecurity — "
    "delivering excellence across four pillars: Education | Research & Innovation | "
    "Down-streaming & Professional Services | Community Outreach",
    Inches(0.65), Inches(1.23), Inches(12.0), Inches(0.58),
    size=11, color=LIGHT_GREY, italic=True)

zones = [
    ("Zone 1", "AI Compute\nCenter",           "±170 m²",
     "GPU/HPC cluster, storage,\nprivate cloud, MLOps", CYAN),
    ("Zone 2", "Cyber-Range\n& SOC",            "±150 m²",
     "Red/Blue/Purple team sim,\nStandoff 365, video wall", TEAL),
    ("Zone 3", "Practicum Labs &\nClassrooms",  "±250 m²",
     "AI, data science & cybersecurity\npracticums, 50-52 seats/room", ORANGE),
    ("Zone 4", "Research &\nIncubation",        "±150 m²",
     "Faculty-student research,\nco-working, startup incubation", CYAN),
    ("Zone 5", "Support\nFacilities",           "±157 m²",
     "Technician room, store,\nmeeting room, lounge", TEAL),
]
for i, (znum, zname, area, desc, col) in enumerate(zones):
    cx = Inches(0.45) + i * Inches(2.54)
    card(s6, cx, Inches(2.0), Inches(2.4), Inches(4.65), fill=DARK_BLUE)
    rect(s6, cx, Inches(2.0), Inches(2.4), Inches(0.1), col)
    # Zone label badge
    rect(s6, cx + Inches(0.15), Inches(2.18), Inches(0.92), Inches(0.32), MID_BLUE)
    txt(s6, znum, cx + Inches(0.15), Inches(2.2), Inches(0.92), Inches(0.28),
        size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s6, zname, cx + Inches(0.1), Inches(2.56), Inches(2.2), Inches(0.72),
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s6, area, cx + Inches(0.1), Inches(3.32), Inches(2.2), Inches(0.38),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s6, desc, cx + Inches(0.1), Inches(3.72), Inches(2.2), Inches(1.3),
        size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Total label
rect(s6, Inches(0.45), Inches(6.72), Inches(12.4), Inches(0.5), MID_BLUE)
txt(s6, "Total Floor Area: ±977 m²   |   Functional Zones Subtotal: ±877 m²   |   "
        "Circulation & Utilities: ±100 m²",
    Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.42),
    size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – TECHNICAL SPECIFICATIONS & EQUIPMENT
# ══════════════════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
bg(s7, NAVY)
rect(s7, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s7, Inches(1.1), color=CYAN)
section_tag(s7, "05  TECHNICAL SPECIFICATIONS")
txt(s7, "Technology Architecture & Key Equipment",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

# Architecture layers
arch = [
    ("Applications Layer",   "AI/ML Frameworks · MLOps · MaxPatrol SIEM/VM · PT Sandbox · Standoff 365", TEAL),
    ("Security Layer",       "PT NAD · PT AF · NGFW · Micro-segmentation · SOC/SOAR · Physical access control", CYAN),
    ("Network Layer",        "L3 Core Switch · Access Switches · Isolated Cyber-Range Network · Redundant UPS", ORANGE),
    ("Compute & Storage",    "GPU/HPC Cluster (≥8 accel/unit) · ≥200 TB All-flash Storage · Private Cloud", CYAN),
]
for i, (layer, desc, col) in enumerate(arch):
    cy = Inches(1.25) + i * Inches(0.78)
    rect(s7, Inches(0.45), cy, Inches(0.07), Inches(0.65), col)
    card(s7, Inches(0.52), cy, Inches(12.33), Inches(0.65), fill=DARK_BLUE)
    txt(s7, layer, Inches(0.68), cy + Inches(0.05), Inches(2.6), Inches(0.55),
        size=11, bold=True, color=col)
    txt(s7, desc,  Inches(3.3),  cy + Inches(0.1), Inches(9.3), Inches(0.5),
        size=10, color=LIGHT_GREY)

# Equipment table header
txt(s7, "Key Equipment (Indicative)", Inches(0.45), Inches(4.5), Inches(12), Inches(0.38),
    size=12, bold=True, color=WHITE)
rect(s7, Inches(0.45), Inches(4.88), Inches(12.4), Inches(0.38), MID_BLUE)
for hdr, lx, lw in [("Category", Inches(0.55), Inches(3.0)),
                     ("Components & Specs",  Inches(3.65), Inches(7.5)),
                     ("Qty",   Inches(11.25), Inches(1.4))]:
    txt(s7, hdr, lx, Inches(4.9), lw, Inches(0.34),
        size=10, bold=True, color=CYAN)

equip_rows = [
    ("AI Compute",       "Data-center GPU servers (≥8 accelerators/unit), NVMe storage", "3–4 units"),
    ("Storage",          "High-capacity all-flash / NAS-SAN                              ", "≥200 TB"),
    ("Cyber-Range/SOC",  "Cyber-range platform, simulation servers, SIEM/SOAR console, video wall", "1 system"),
    ("Workstations",     "High-spec practicum PCs with workstation-class GPU             ", "±100 units"),
    ("Security SW",      "MaxPatrol SIEM/VM, PT NAD, PT Sandbox, PT AF, Standoff 365    ", "Academic lic."),
    ("Power & Cooling",  "Redundant UPS, PDU, precision cooling (Zones 1–2)             ", "1 pkg"),
]
for ri, (cat, comp, qty) in enumerate(equip_rows):
    cy = Inches(5.28) + ri * Inches(0.36)
    row_bg = DARK_BLUE if ri % 2 == 0 else MID_BLUE
    card(s7, Inches(0.45), cy, Inches(12.4), Inches(0.35), fill=row_bg)
    txt(s7, cat,  Inches(0.55), cy + Inches(0.04), Inches(3.0),  Inches(0.3), size=9, color=TEAL)
    txt(s7, comp, Inches(3.65), cy + Inches(0.04), Inches(7.5),  Inches(0.3), size=9, color=LIGHT_GREY)
    txt(s7, qty,  Inches(11.25),cy + Inches(0.04), Inches(1.4),  Inches(0.3), size=9, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – ACADEMIC PROGRAMMES & HUMAN CAPITAL
# ══════════════════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
bg(s8, NAVY)
rect(s8, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s8, Inches(1.1), color=CYAN)
section_tag(s8, "06  ACADEMIC PROGRAMMES")
txt(s8, "Master's Programme (S2) & Human Capital Development",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

# Left: S2 programme details
card(s8, Inches(0.45), Inches(1.25), Inches(6.2), Inches(5.8), fill=DARK_BLUE)
rect(s8, Inches(0.45), Inches(1.25), Inches(6.2), Inches(0.08), CYAN)
txt(s8, "Master of Computer Science (M.Kom.)",
    Inches(0.6), Inches(1.38), Inches(5.9), Inches(0.42),
    size=13, bold=True, color=CYAN)
txt(s8, "Artificial Intelligence & Cybersecurity  |  4 Semesters  |  ±39 SKS",
    Inches(0.6), Inches(1.78), Inches(5.9), Inches(0.35),
    size=10, color=LIGHT_GREY)
txt(s8, 'Bridged by "AI for Security" theme  |  Target 1st intake: AY 2027/2028',
    Inches(0.6), Inches(2.1), Inches(5.9), Inches(0.35),
    size=10, italic=True, color=TEAL)

sem_data = [
    ("I",   "Adv. Machine Learning · Foundations of Cybersecurity · Big Data Analytics · Research Methodology", "11 SKS"),
    ("II",  "Adversarial ML & AI for Security · Concentration Electives I–III", "12 SKS"),
    ("III", "Cyber-Range Capstone (Standoff 365) · Elective IV · Thesis Proposal · Certification/Internship", "10 SKS"),
    ("IV",  "Thesis & Scientific Publication", "6 SKS"),
]
for ri, (sem, courses, sks) in enumerate(sem_data):
    cy = Inches(2.55) + ri * Inches(1.08)
    row_col = MID_BLUE if ri % 2 == 0 else DARK_BLUE
    card(s8, Inches(0.55), cy, Inches(6.0), Inches(0.98), fill=row_col)
    rect(s8, Inches(0.55), cy, Inches(0.55), Inches(0.98), CYAN if ri % 2 == 0 else TEAL)
    txt(s8, f"Sem {sem}", Inches(0.55), cy + Inches(0.04), Inches(0.55), Inches(0.9),
        size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s8, courses, Inches(1.18), cy + Inches(0.08), Inches(4.8), Inches(0.72),
        size=9, color=LIGHT_GREY, wrap=True)
    txt(s8, sks, Inches(5.85), cy + Inches(0.3), Inches(0.65), Inches(0.38),
        size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Right: concentrations + other programmes
card(s8, Inches(6.9), Inches(1.25), Inches(5.9), Inches(2.65), fill=DARK_BLUE)
rect(s8, Inches(6.9), Inches(1.25), Inches(5.9), Inches(0.08), TEAL)
txt(s8, "Concentration Tracks", Inches(7.05), Inches(1.38), Inches(5.6), Inches(0.42),
    size=13, bold=True, color=TEAL)
ai_elec = ["Natural Language Processing", "Computer Vision",
           "Generative AI & LLMs", "MLOps & Scalable AI", "Reinforcement Learning"]
cs_elec = ["Offensive Security & Pen Testing", "Digital Forensics & IR",
           "Network & Infrastructure Security", "Applied Cryptography", "SOC & Threat Intelligence"]
txt(s8, "🤖  Artificial Intelligence", Inches(7.05), Inches(1.82), Inches(5.6), Inches(0.35),
    size=10, bold=True, color=CYAN)
bullet_items(s8, ai_elec, Inches(7.15), Inches(2.15), Inches(5.55), Inches(1.25),
             size=9, bullet="·")
txt(s8, "🔐  Cybersecurity", Inches(7.05), Inches(3.08), Inches(5.6), Inches(0.35),
    size=10, bold=True, color=TEAL)
bullet_items(s8, cs_elec, Inches(7.15), Inches(3.4), Inches(5.55), Inches(1.3),
             size=9, bullet="·")

card(s8, Inches(6.9), Inches(4.05), Inches(5.9), Inches(2.98), fill=DARK_BLUE)
rect(s8, Inches(6.9), Inches(4.05), Inches(5.9), Inches(0.08), ORANGE)
txt(s8, "Other Human Capital Programmes", Inches(7.05), Inches(4.18), Inches(5.6), Inches(0.42),
    size=13, bold=True, color=ORANGE)
other_progs = [
    "Professional Certification tracks (MaxPatrol, PT NAD, Standoff 365)",
    "Training of Trainers (ToT) — jointly led with Positive Technologies",
    "Intensive bootcamps and short courses for industry & government",
    "CTF competitions — host & participate (national / international)",
    "Security assessment & penetration-testing services",
    "Startup incubation and community cybersecurity literacy",
]
bullet_items(s8, other_progs, Inches(7.05), Inches(4.62), Inches(5.65), Inches(2.2),
             size=10, bullet="▸", color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – BUDGET & INVESTMENT
# ══════════════════════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
bg(s9, NAVY)
rect(s9, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s9, Inches(1.1), color=CYAN)
section_tag(s9, "07  BUDGET & INVESTMENT")
txt(s9, "Budget & Investment Overview", Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)

# Summary KPI row
for val, lbl, col, cx in [
    ("IDR 45.1B", "Total CAPEX",            CYAN,   Inches(0.45)),
    ("IDR 5.8B",  "Annual OPEX",            TEAL,   Inches(3.45)),
    ("55 / 20 / 25%", "UNDIP / PT / Ext.", ORANGE, Inches(6.45)),
    ("Multi-source","Funding Strategy",     CYAN,   Inches(9.7)),
]:
    card(s9, cx, Inches(1.18), Inches(2.7), Inches(0.9), fill=DARK_BLUE)
    rect(s9, cx, Inches(1.18), Inches(2.7), Inches(0.07), col)
    txt(s9, val, cx, Inches(1.28), Inches(2.7), Inches(0.48),
        size=17, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s9, lbl, cx, Inches(1.73), Inches(2.7), Inches(0.32),
        size=9, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# CAPEX table
txt(s9, "Capital Expenditure (CAPEX) Breakdown", Inches(0.45), Inches(2.25),
    Inches(8.1), Inches(0.38), size=11, bold=True, color=WHITE)
rect(s9, Inches(0.45), Inches(2.63), Inches(8.1), Inches(0.36), MID_BLUE)
txt(s9, "Investment Component", Inches(0.55), Inches(2.66), Inches(5.1), Inches(0.3),
    size=9, bold=True, color=CYAN)
txt(s9, "IDR (Billion)", Inches(5.75), Inches(2.66), Inches(1.5), Inches(0.3),
    size=9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s9, "Source", Inches(7.35), Inches(2.66), Inches(1.15), Inches(0.3),
    size=9, bold=True, color=CYAN)

capex_rows = [
    ("Interior & MEP fit-out (±977 m²)",                "7.8",  "RKAT / Grant"),
    ("AI Compute (GPU/HPC cluster, storage)",            "12.0", "Grant / RKAT"),
    ("Cyber-Range & SOC infrastructure",                 "6.5",  "RKAT / Partner"),
    ("Workstations & practicum devices (±100 units)",    "4.0",  "RKAT"),
    ("Positive Technologies licenses & platform access", "6.0",  "Partner in-kind"),
    ("Network, UPS, precision cooling & physical sec.",  "3.2",  "RKAT"),
    ("Furniture, AV systems & signage",                  "1.5",  "RKAT"),
    ("Project management & contingency (±10%)",          "4.1",  "RKAT"),
]
for ri, (comp, idr, src) in enumerate(capex_rows):
    cy = Inches(2.99) + ri * Inches(0.37)
    row_bg = DARK_BLUE if ri % 2 == 0 else MID_BLUE
    card(s9, Inches(0.45), cy, Inches(8.1), Inches(0.36), fill=row_bg)
    txt(s9, comp, Inches(0.55), cy + Inches(0.05), Inches(5.1), Inches(0.28), size=9, color=LIGHT_GREY)
    txt(s9, idr,  Inches(5.75), cy + Inches(0.05), Inches(1.5), Inches(0.28),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s9, src,  Inches(7.35), cy + Inches(0.05), Inches(1.15), Inches(0.28), size=8, color=TEAL)

# TOTAL row
card(s9, Inches(0.45), Inches(5.96), Inches(8.1), Inches(0.42), fill=MID_BLUE)
rect(s9, Inches(0.45), Inches(5.96), Inches(8.1), Inches(0.06), CYAN)
txt(s9, "TOTAL INDICATIVE CAPEX", Inches(0.55), Inches(6.0), Inches(5.1), Inches(0.35),
    size=11, bold=True, color=WHITE)
txt(s9, "± 45.1B", Inches(5.75), Inches(6.0), Inches(2.7), Inches(0.35),
    size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# OPEX panel (right)
card(s9, Inches(8.75), Inches(2.25), Inches(4.1), Inches(4.2), fill=DARK_BLUE)
rect(s9, Inches(8.75), Inches(2.25), Inches(4.1), Inches(0.08), TEAL)
txt(s9, "Annual OPEX ≈ IDR 5.8B/year", Inches(8.9), Inches(2.38),
    Inches(3.8), Inches(0.42), size=11, bold=True, color=TEAL)
opex_items = [
    ("Licenses & SW maintenance",  "1.8B"),
    ("Infrastructure & utilities",  "1.5B"),
    ("Management staff",            "1.2B"),
    ("Training & certification",    "0.8B"),
    ("Operations & admin",          "0.5B"),
]
for ri, (item, val) in enumerate(opex_items):
    cy = Inches(2.88) + ri * Inches(0.65)
    pct = float(val[:-1]) / 5.8
    bar_w = Inches(3.3) * pct
    txt(s9, item, Inches(8.9), cy, Inches(3.8), Inches(0.3), size=9, color=LIGHT_GREY)
    rect(s9, Inches(8.9), cy + Inches(0.3), Inches(3.3), Inches(0.2), DARK_GREY)
    rect(s9, Inches(8.9), cy + Inches(0.3), bar_w,       Inches(0.2), TEAL)
    txt(s9, val, Inches(8.9) + bar_w + Inches(0.05), cy + Inches(0.27),
        Inches(0.6), Inches(0.26), size=9, bold=True, color=TEAL)

# Funding split
card(s9, Inches(8.75), Inches(6.53), Inches(4.1), Inches(0.72), fill=MID_BLUE)
txt(s9, "Funding Split:   UNDIP 55%   |   Positive Technologies 20% (in-kind)   |   Grants/CSR 25%",
    Inches(8.78), Inches(6.6), Inches(4.0), Inches(0.58),
    size=9, color=LIGHT_GREY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – GOVERNANCE & KPIs
# ══════════════════════════════════════════════════════════════════════════════
s10 = blank_slide(prs)
bg(s10, NAVY)
rect(s10, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s10, Inches(1.1), color=CYAN)
section_tag(s10, "08  GOVERNANCE & KPIs")
txt(s10, "Governance Structure & Key Performance Indicators",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

# Governance org
card(s10, Inches(0.45), Inches(1.25), Inches(5.5), Inches(5.9), fill=DARK_BLUE)
rect(s10, Inches(0.45), Inches(1.25), Inches(5.5), Inches(0.08), CYAN)
txt(s10, "Governance Structure", Inches(0.6), Inches(1.38),
    Inches(5.2), Inches(0.4), size=13, bold=True, color=CYAN)

gov_levels = [
    ("Dean of FSM, Diponegoro University",       CYAN,  "Strategic oversight"),
    ("Joint Steering Committee (UNDIP–Partner)", TEAL,  "Strategic direction"),
    ("Head of Laboratory",                       WHITE, "Laboratory leadership"),
    ("Joint Technical Team",                     TEAL,  "Implementation"),
    ("Operations Manager",                       LIGHT_GREY, "Day-to-day operations"),
    ("Research Coordinator",                     LIGHT_GREY, "Research management"),
    ("Education & Certification Coord.",         LIGHT_GREY, "Academic programmes"),
    ("Infrastructure Engineers & SOC Analysts",  LIGHT_GREY, "Technical operations"),
]
for ri, (role, col, desc) in enumerate(gov_levels):
    cy = Inches(1.85) + ri * Inches(0.62)
    indent = Inches(0.15) * min(ri, 3)
    card(s10, Inches(0.6) + indent, cy, Inches(5.15) - indent, Inches(0.52), fill=MID_BLUE)
    rect(s10, Inches(0.6) + indent, cy, Inches(0.07), Inches(0.52), col)
    txt(s10, role, Inches(0.78) + indent, cy + Inches(0.04),
        Inches(3.5) - indent, Inches(0.28), size=9, bold=(ri < 4), color=col)
    txt(s10, desc, Inches(0.78) + indent, cy + Inches(0.3),
        Inches(3.5) - indent, Inches(0.22), size=8, color=DARK_GREY if col == LIGHT_GREY else LIGHT_GREY,
        italic=True)

# KPI table
txt(s10, "Key Performance Indicators", Inches(6.15), Inches(1.25),
    Inches(6.7), Inches(0.38), size=12, bold=True, color=WHITE)
rect(s10, Inches(6.15), Inches(1.63), Inches(6.7), Inches(0.36), MID_BLUE)
for hdr, lx, lw in [("KPI", Inches(6.25), Inches(3.2)),
                     ("2027",Inches(9.55), Inches(0.95)),
                     ("2028",Inches(10.55),Inches(0.95)),
                     ("2029",Inches(11.55),Inches(1.1))]:
    txt(s10, hdr, lx, Inches(1.66), lw, Inches(0.3),
        size=9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

kpi_rows = [
    ("Students in lab practicums/programmes (persons/yr)",   "300",  "800",   "1,200"),
    ("Professionally certified participants (persons/yr)",   "40",   "120",   "250"),
    ("Lecturers/instructors completing ToT (cumulative)",    "10",   "20",    "30"),
    ("Scientific publications (indexed) (articles/yr)",     "5",    "15",    "25"),
    ("Research & grants won (titles/yr)",                   "3",    "6",     "10"),
    ("Competitions hosted/participated (Standoff/CTF/yr)",  "1",    "2",     "4"),
    ("Services to industry/gov't (contracts/yr)",           "2",    "6",     "12"),
    ("Incubated startups/products (units/yr)",              "1",    "2",     "4"),
    ("Service revenue / PNBP (IDR billion/yr)",             "1.0",  "3.0",   "6.0"),
]
for ri, (kpi, y27, y28, y29) in enumerate(kpi_rows):
    cy = Inches(1.99) + ri * Inches(0.58)
    row_bg = DARK_BLUE if ri % 2 == 0 else MID_BLUE
    card(s10, Inches(6.15), cy, Inches(6.7), Inches(0.56), fill=row_bg)
    txt(s10, kpi, Inches(6.25), cy + Inches(0.08), Inches(3.2), Inches(0.42), size=9, color=LIGHT_GREY)
    for val, col, lx in [(y27, TEAL,   Inches(9.55)),
                          (y28, CYAN,   Inches(10.55)),
                          (y29, ORANGE, Inches(11.55))]:
        txt(s10, val, lx, cy + Inches(0.1), Inches(0.95), Inches(0.38),
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – IMPLEMENTATION TIMELINE
# ══════════════════════════════════════════════════════════════════════════════
s11 = blank_slide(prs)
bg(s11, NAVY)
rect(s11, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s11, Inches(1.1), color=CYAN)
section_tag(s11, "09  IMPLEMENTATION TIMELINE")
txt(s11, "Implementation Roadmap — 2026 to 2028",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)

phases = [
    ("Phase 0", "Jun–Oct 2026",  "Initiation,\nDED & Site Prep",
     ["LoI & MoU", "DED completed (Jun 2026)", "Site prep begun (17 Jun 2026)",
      "S2 curriculum drafting", "Procurement planning"],
     CYAN),
    ("Phase 1", "Oct–Dec 2026", "Procurement &\nMobilization",
     ["MoA & Impl. Arrangement", "Final budget approval",
      "Tender & procurement", "Lab construction begins (Dec 2026)"],
     TEAL),
    ("Phase 2", "Jan–Jun 2027", "Construction\n& Fit-out",
     ["Interior works & MEP", "Server room & raised floor",
      "Precision cooling installed", "Lab space ready (Jun 2027)"],
     ORANGE),
    ("Phase 3", "Apr–Sep 2027", "Equipment\nInstallation",
     ["HPC/GPU & storage install", "Partner solutions deployed",
      "Cyber-range configured", "Systems tested (Sep 2027)"],
     CYAN),
    ("Phase 4", "Sep–Dec 2027", "Integration\n& ToT",
     ["System integration & testing", "Training of Trainers",
      "SOPs finalized", "Soft launch / pilot (Dec 2027)"],
     TEAL),
    ("Phase 5", "Jan 2028+",    "Full Operation\n& Launch",
     ["Official launch", "First S2 (Master's) intake",
      "Certification programmes", "Research & industry services"],
     ORANGE),
]

# Timeline connector bar
rect(s11, Inches(0.45), Inches(2.52), Inches(12.4), Inches(0.06), DARK_GREY)

for i, (phase, period, title, pts, col) in enumerate(phases):
    cx = Inches(0.45) + i * Inches(2.12)
    # dot on timeline
    dot = rect(s11, cx + Inches(0.72), Inches(2.35), Inches(0.4), Inches(0.4), col)
    card(s11, cx, Inches(2.85), Inches(2.0), Inches(4.3), fill=DARK_BLUE)
    rect(s11, cx, Inches(2.85), Inches(2.0), Inches(0.08), col)
    txt(s11, phase,  cx, Inches(2.96), Inches(2.0), Inches(0.35),
        size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s11, period, cx, Inches(3.28), Inches(2.0), Inches(0.32),
        size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    txt(s11, title,  cx, Inches(3.6), Inches(2.0), Inches(0.62),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_items(s11, pts, cx + Inches(0.08), Inches(4.28), Inches(1.88),
                 Inches(2.65), size=9, color=LIGHT_GREY, bullet="·")

# Milestone highlights row
rect(s11, Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.06), CYAN)
milestones = [
    ("Jun 2026", "DED Complete"),
    ("17 Jun 2026", "Site Prep"),
    ("Oct 2026", "MoU Signed"),
    ("Dec 2026", "Build Starts"),
    ("Jun 2027", "Space Ready"),
    ("Sep 2027", "Equip. Done"),
    ("Dec 2027", "Soft Launch"),
    ("Q1 2028", "🚀 Official Launch"),
]
for i, (dt, lbl) in enumerate(milestones):
    cx = Inches(0.45) + i * Inches(1.58)
    txt(s11, dt,  cx, Inches(7.18), Inches(1.55), Inches(0.2),
        size=7, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    txt(s11, lbl, cx, Inches(7.35), Inches(1.55), Inches(0.2),
        size=7, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – RISK MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════
s12 = blank_slide(prs)
bg(s12, NAVY)
rect(s12, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s12, Inches(1.1), color=CYAN)
section_tag(s12, "10  RISK MANAGEMENT")
txt(s12, "Risk Register & Mitigation Strategies",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)

rect(s12, Inches(0.45), Inches(1.18), Inches(12.4), Inches(0.4), MID_BLUE)
for hdr, lx, lw in [
    ("#",              Inches(0.55), Inches(0.5)),
    ("Risk",           Inches(1.1),  Inches(2.8)),
    ("Category",       Inches(3.95), Inches(1.5)),
    ("Impact",         Inches(5.5),  Inches(1.1)),
    ("Likelihood",     Inches(6.65), Inches(1.1)),
    ("Mitigation",     Inches(7.82), Inches(5.0)),
]:
    txt(s12, hdr, lx, Inches(1.22), lw, Inches(0.32),
        size=9, bold=True, color=CYAN)

risks = [
    ("1", "Funding shortfall / delay",           "Financial",   "High",   "Med",
     "Multi-source funding; phased investment; firm RKAT commitment"),
    ("2", "Construction / fit-out delay",         "Operational", "Med",    "Med",
     "Realistic scheduling; competent contractor; critical-path monitoring"),
    ("3", "Shortage of qualified HR",             "Human Capital","High",  "Med",
     "Early ToT; partner-supported recruitment; certification programmes"),
    ("4", "Cybersecurity incident in the lab",    "Security",    "High",   "Med",
     "Strict network isolation; hardening; SOPs; regular audits"),
    ("5", "Technology dependence on partner",     "Strategic",   "Med",    "Med",
     "Technology-transfer clauses; internal capacity; multi-vendor where feasible"),
    ("6", "Operational sustainability",           "Financial",   "Med",    "Med",
     "Service/research revenue (PNBP); periodic OPEX review"),
    ("7", "Regulatory / compliance issues",       "Compliance",  "Med",    "Low",
     "Legal review; alignment with national regulations; data protection"),
    ("8", "Geopolitical / sanctions exposure",    "Compliance",  "Med",    "Med",
     "Due diligence; academic/in-kind scope; legal & inter-ministerial consultation"),
]
for ri, (num, risk, cat, imp, lik, mit) in enumerate(risks):
    cy = Inches(1.6) + ri * Inches(0.71)
    row_bg = DARK_BLUE if ri % 2 == 0 else MID_BLUE
    card(s12, Inches(0.45), cy, Inches(12.4), Inches(0.68), fill=row_bg)
    imp_col  = ORANGE if imp == "High" else (TEAL if imp == "Med" else CYAN)
    lik_col  = ORANGE if lik == "Med"  else TEAL
    txt(s12, num, Inches(0.55), cy + Inches(0.18), Inches(0.5), Inches(0.32),
        size=10, bold=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    txt(s12, risk, Inches(1.1), cy + Inches(0.1), Inches(2.8), Inches(0.5),
        size=9, color=WHITE, wrap=True)
    txt(s12, cat, Inches(3.95), cy + Inches(0.18), Inches(1.5), Inches(0.32),
        size=9, color=LIGHT_GREY)
    txt(s12, imp, Inches(5.5), cy + Inches(0.18), Inches(1.0), Inches(0.32),
        size=9, bold=True, color=imp_col)
    txt(s12, lik, Inches(6.65), cy + Inches(0.18), Inches(1.0), Inches(0.32),
        size=9, bold=True, color=lik_col)
    txt(s12, mit, Inches(7.82), cy + Inches(0.1), Inches(5.0), Inches(0.5),
        size=9, color=LIGHT_GREY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – IMPACT & BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
s13 = blank_slide(prs)
bg(s13, NAVY)
rect(s13, 0, 0, W, Inches(1.1), DARK_BLUE)
accent_line(s13, Inches(1.1), color=CYAN)
section_tag(s13, "EXPECTED IMPACT")
txt(s13, "Who Benefits — and How?",
    Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)

stakeholders = [
    ("🏛", "UNDIP / FSM",
     ["Enhanced academic & international reputation",
      "World-class computing & security facilities",
      "Stronger IKU (Key Performance Indicators)",
      "New revenue streams via PNBP services"],
     CYAN),
    ("🎓", "Students &\nLecturers",
     ["Practice-based, cyber-range learning",
      "Embedded professional certification",
      "Research & publication opportunities",
      "CTF competitions (national & international)"],
     TEAL),
    ("🏢", "Positive\nTechnologies",
     ["Academic engagement & talent pipeline",
      "Regional reference deployment",
      "Collaborative research & publications",
      "Curriculum & Training of Trainers"],
     ORANGE),
    ("🇮🇩", "Nation &\nRegion",
     ["Certified digital talent at scale",
      "Strengthened national cyber resilience",
      "Digital sovereignty protection",
      "Applied AI & cybersecurity research"],
     CYAN),
]
for i, (icon, name, pts, col) in enumerate(stakeholders):
    cx = Inches(0.45) + i * Inches(3.18)
    card(s13, cx, Inches(1.25), Inches(3.0), Inches(5.35), fill=DARK_BLUE)
    rect(s13, cx, Inches(1.25), Inches(3.0), Inches(0.1), col)
    txt(s13, icon, cx, Inches(1.42), Inches(3.0), Inches(0.62),
        size=28, align=PP_ALIGN.CENTER, color=col)
    txt(s13, name, cx, Inches(2.05), Inches(3.0), Inches(0.65),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    accent_line(s13, Inches(2.75), color=col, l=cx + Inches(0.3), w=Inches(2.4))
    bullet_items(s13, pts, cx + Inches(0.1), Inches(2.88), Inches(2.8),
                 Inches(3.5), size=11, color=LIGHT_GREY, bullet="✓")

# Bottom tagline
card(s13, Inches(0.45), Inches(6.72), Inches(12.4), Inches(0.52), fill=MID_BLUE)
rect(s13, Inches(0.45), Inches(6.72), Inches(12.4), Inches(0.07), CYAN)
txt(s13,
    '"A strategic, timely, and high-impact investment in Indonesia\'s '
    'digital future — accelerating certified talent, research excellence, '
    'and national cyber resilience."',
    Inches(0.55), Inches(6.78), Inches(12.2), Inches(0.42),
    size=11, italic=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CLOSING / CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
s14 = blank_slide(prs)
bg(s14, NAVY)

# Large background accent shapes
rect(s14, 0, 0, Inches(7.0), H, DARK_BLUE)
rect(s14, 0, 0, Inches(7.0), Inches(0.12), CYAN)
rect(s14, 0, H - Inches(0.12), Inches(7.0), Inches(0.12), TEAL)

# Decorative dots right panel
for row in range(8):
    for col in range(9):
        rect(s14,
             Inches(7.4) + col * Inches(0.64),
             Inches(0.5) + row * Inches(0.9),
             Inches(0.1), Inches(0.1),
             TEAL if (row + col) % 3 == 0 else DARK_GREY)

txt(s14, "🚀", Inches(8.5), Inches(1.8), Inches(2.5), Inches(2.0),
    size=80, align=PP_ALIGN.CENTER, color=TEAL)
txt(s14, "OFFICIAL LAUNCH", Inches(7.3), Inches(4.7), Inches(5.7), Inches(0.55),
    size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s14, "Q1 2028", Inches(7.3), Inches(5.2), Inches(5.7), Inches(0.55),
    size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left call-to-action
txt(s14, "11  CLOSING", Inches(0.45), Inches(0.8), Inches(6.2), Inches(0.38),
    size=11, bold=True, color=CYAN)
accent_line(s14, Inches(1.28), color=CYAN, l=Inches(0.45), w=Inches(6.2))
txt(s14, "Call to Action",
    Inches(0.45), Inches(1.38), Inches(6.2), Inches(0.65),
    size=28, bold=True, color=WHITE)

closing_msg = (
    "The AICS-Lab UNDIP represents a strategic and timely investment "
    "in Indonesia's digital future — combining UNDIP's academic strength "
    "with Positive Technologies' world-class cybersecurity tools."
)
txt(s14, closing_msg, Inches(0.45), Inches(2.1), Inches(6.2), Inches(1.2),
    size=12, color=LIGHT_GREY, wrap=True)

actions = [
    "✅  Approve and endorse this Partnership Proposal",
    "✅  Authorize LoI signing → progress to MoU (Oct 2026)",
    "✅  Confirm budget allocation in RKAT PTN-BH",
    "✅  Support S2 study-programme permit processing",
    "✅  Enable legal & compliance consultation (MoFA/BSSN)",
]
for i, action in enumerate(actions):
    cy = Inches(3.45) + i * Inches(0.62)
    card(s14, Inches(0.45), cy, Inches(6.2), Inches(0.52), fill=MID_BLUE)
    rect(s14, Inches(0.45), cy, Inches(0.07), Inches(0.52), TEAL)
    txt(s14, action, Inches(0.6), cy + Inches(0.1), Inches(5.9), Inches(0.35),
        size=11, bold=False, color=WHITE)

# Contact footer
rect(s14, Inches(0.45), Inches(6.65), Inches(6.2), Inches(0.62), DARK_GREY)
txt(s14,
    "Faculty of Science and Mathematics, Diponegoro University\n"
    "Department of Informatics  |  7th Floor FSM Central Laboratory Building, Semarang  |  2026",
    Inches(0.55), Inches(6.7), Inches(6.0), Inches(0.55),
    size=9, color=LIGHT_GREY, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output_path = "knowledge_management/AICS_Lab_UNDIP_Proposal_Presentation.pptx"
prs.save(output_path)
print(f"✅  Presentation saved to: {output_path}")
print(f"    Slides: {len(prs.slides)}")

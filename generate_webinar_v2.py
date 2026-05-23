"""
Webinar PPT v2 — Battle of Minds: Harari vs Andrew Ng
Versi dengan foto tokoh + grafik statistik embed
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

ASSETS = Path("/projects/sandbox/sonet-product/assets")

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG      = rgb("0A0A0F")
RED     = rgb("E63946")
BLUE    = rgb("00B4D8")
GOLD    = rgb("FFD166")
GREEN   = rgb("06D6A0")
WHITE   = rgb("FFFFFF")
GRAY    = rgb("AAAAAA")
CARD    = rgb("14141E")
MIDBLUE = rgb("1A3A5C")
PURPLE  = rgb("9B5DE5")

def bg(slide, color=None):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color or BG

def rect(slide, l, t, w, h, color, line=None, lw=1.5):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False,
        color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or WHITE
    return tb

def htxt(slide, text, l, t, w, h, size=14, bold=False,
         color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color or WHITE
    return tb

def img(slide, path, l, t, w, h=None):
    """Add image; h=None → auto aspect ratio"""
    p = ASSETS / path
    if not p.exists():
        return None
    if h:
        return slide.shapes.add_picture(
            str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        return slide.shapes.add_picture(
            str(p), Inches(l), Inches(t), Inches(w))

def top_bar(slide, color):
    rect(slide, 0, 0,     13.33, 0.07, color)
    rect(slide, 0, 7.43,  13.33, 0.07, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — EPIC COVER  (split red/blue + foto kedua tokoh)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)

rect(sl, 0,    0, 6.66, 7.5, rgb("0D0508"))
rect(sl, 6.66, 0, 6.67, 7.5, rgb("030D12"))
rect(sl, 6.5,  0, 0.07, 7.5, rgb("222235"))
top_bar(sl, GOLD)

# Badge
rect(sl, 4.3, 0.2, 4.73, 0.52, rgb("1A1A2E"), line=GOLD)
txt(sl, "⚡  WEBINAR NASIONAL  ⚡", 4.3, 0.2, 4.73, 0.52,
    size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Title
htxt(sl, "BATTLE OF MINDS", 1.1, 1.0, 11.13, 1.35,
     size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
htxt(sl, "AI: ANCAMAN atau PELUANG?", 1.1, 2.45, 11.13, 0.75,
     size=26, color=GOLD, align=PP_ALIGN.CENTER)

# VS badge
rect(sl, 5.9, 3.3, 1.53, 0.68, rgb("1A1A2E"), line=GOLD)
txt(sl, "VS", 5.9, 3.3, 1.53, 0.68,
    size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Foto Harari (kiri)
img(sl, "harari.jpg", 0.3, 3.0, 2.2, 3.8)
rect(sl, 0.3, 3.0, 2.2, 0.55, rgb("2A0810"))
txt(sl, "YUVAL NOAH HARARI", 0.3, 3.04, 2.2, 0.46,
    size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)
htxt(sl, "Sejarawan · Penulis Sapiens\nPerspektif Kritis & Historis",
     0.25, 6.85, 2.3, 0.6, size=9, italic=True,
     color=GRAY, align=PP_ALIGN.CENTER)

# Foto Andrew Ng (kanan)
img(sl, "andrew.jpg", 10.83, 3.0, 2.2, 3.8)
rect(sl, 10.83, 3.0, 2.2, 0.55, rgb("031420"))
txt(sl, "ANDREW NG", 10.83, 3.04, 2.2, 0.46,
    size=9, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
htxt(sl, "Pendiri DeepLearning.AI\nPerspektif Optimis & Pragmatis",
     10.78, 6.85, 2.3, 0.6, size=9, italic=True,
     color=GRAY, align=PP_ALIGN.CENTER)

# Left info box
rect(sl, 2.65, 3.2, 3.1, 3.5, rgb("180408"), line=RED)
htxt(sl, "🔴  SISI KRITIS\n\n"
         "\"AI bukan sekadar alat.\n"
         "IA bisa memanipulasi\n"
         "dan mengontrol narasi.\"",
     2.78, 3.35, 2.85, 3.2,
     size=12, color=WHITE)

# Right info box
rect(sl, 7.58, 3.2, 3.1, 3.5, rgb("03141E"), line=BLUE)
htxt(sl, "🔵  SISI OPTIMIS\n\n"
         "\"AI adalah listrik baru.\n"
         "Siapa yang belajar AI\n"
         "akan memimpin dunia.\"",
     7.72, 3.35, 2.85, 3.2,
     size=12, color=WHITE)

# Topics footer
htxt(sl, "Topik 1: AI dalam Dunia Pendidikan   |   Topik 2: AI dalam Dunia Kerja",
     1.0, 7.0, 11.33, 0.38, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA (rundown)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, GOLD)

txt(sl, "📋  AGENDA WEBINAR", 0.5, 0.18, 12.33, 0.62,
    size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

items = [
    ("🔥","01","OPENING — Pertanyaan Pemantik",
     "Apakah AI ancaman atau peluang bagimu?"),
    ("🔴","02","PERSPEKTIF HARARI",
     "AI dalam Pendidikan & Karir — Sisi Kritis + Foto Tokoh"),
    ("🔵","03","PERSPEKTIF ANDREW NG",
     "AI dalam Pendidikan & Karir — Sisi Optimis + Foto Tokoh"),
    ("📊","04","GRAFIK & DATA GLOBAL",
     "WEF 2025 · McKinsey · Statistik Pendidikan Real"),
    ("⚡","05","DUEL PERSPEKTIF",
     "Radar chart + face-to-face comparison"),
    ("🚀","06","ACTION PLAN",
     "Langkah konkret untuk kamu mulai hari ini"),
]
for i, (icon, num, title, sub) in enumerate(items):
    y = 1.0 + i * 1.02
    rect(sl, 0.5, y, 0.68, 0.78, CARD, line=rgb("2A2A3E"))
    txt(sl, num, 0.5, y, 0.68, 0.78,
        size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, 1.32, y, 11.5, 0.78, CARD, line=rgb("2A2A3E"))
    txt(sl, f"{icon}  {title}",
        1.46, y+0.06, 7.5, 0.38, size=14, bold=True, color=WHITE)
    txt(sl, sub, 1.46, y+0.43, 11.1, 0.3,
        size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — POLL PEMBUKA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, GOLD)

rect(sl, 4.0, 0.18, 5.33, 0.52, rgb("1A1A2E"), line=GOLD)
txt(sl, "⚡  PERTANYAAN PEMBUKA", 4.0, 0.18, 5.33, 0.52,
    size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

htxt(sl, "Menurutmu, AI itu\nANCAMAN atau PELUANG?",
     0.6, 0.88, 12.13, 1.55, size=40, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "🗳️  Ketik jawabanmu di kolom chat sekarang!",
    0.6, 2.55, 12.13, 0.48, size=16, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

rect(sl, 0.55, 3.15, 3.95, 2.05, rgb("2A0810"), line=RED)
txt(sl, "🔴  A", 0.65, 3.22, 3.75, 0.5,
    size=24, bold=True, color=RED, align=PP_ALIGN.CENTER)
htxt(sl, "ANCAMAN\nAI akan mengambil\npekerjaan kita",
     0.65, 3.76, 3.75, 1.3, size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 4.7, 3.15, 3.93, 2.05, rgb("0A0A10"), line=GOLD)
txt(sl, "⚡  C", 4.8, 3.22, 3.73, 0.5,
    size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
htxt(sl, "KEDUANYA\nTergantung bagaimana\nkita menyikapinya",
     4.8, 3.76, 3.73, 1.3, size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 8.83, 3.15, 3.95, 2.05, rgb("031420"), line=BLUE)
txt(sl, "🔵  B", 8.93, 3.22, 3.75, 0.5,
    size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
htxt(sl, "PELUANG\nAI membuka lapangan\nkerja & kemajuan",
     8.93, 3.76, 3.75, 1.3, size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "💡  Tidak ada jawaban salah — semua perspektif valid!",
    0.6, 5.35, 12.13, 0.48, size=13, italic=True,
    color=GRAY, align=PP_ALIGN.CENTER)

# mini chart preview
img(sl, "chart_key_numbers.png", 0.5, 5.9, 12.33, 1.4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — HARARI PROFILE  (foto besar + bio + quote)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, rgb("0D0508"))
top_bar(sl, RED)
rect(sl, 0, 0.07, 0.25, 7.36, RED)

# Foto besar Harari kiri
img(sl, "harari.jpg", 0.35, 0.35, 3.5, 6.8)
rect(sl, 0.35, 0.35, 3.5, 0.55, rgb("2A0810"))
txt(sl, "YUVAL NOAH HARARI", 0.35, 0.38, 3.5, 0.48,
    size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Big quote
txt(sl, "\u201c", 4.1, 0.2, 1.5, 1.8,
    size=130, bold=True, color=rgb("3A0A10"))
htxt(sl, "AI bukan sekadar alat.\nIA adalah kekuatan yang bisa\nbelajar, berubah, dan memanipulasi.",
     4.1, 0.9, 8.8, 2.3, size=32, bold=True, color=WHITE)
rect(sl, 4.1, 3.3, 6.0, 0.05, RED)
htxt(sl, "— Yuval Noah Harari\nSejarawan & Penulis Nexus (2024)",
     4.1, 3.45, 8.5, 0.85, size=15, italic=True, color=RED)

# Bio cards
bio = [
    ("📚 Sapiens (2011) · Homo Deus (2015)\n"
     "   21 Lessons (2018) · Nexus (2024)"),
    ("🎓 Hebrew University Jerusalem\n"
     "   Cambridge — Existential Risk"),
    ("🌍 45 juta+ buku terjual, 65 bahasa\n"
     "   Pembicara Davos WEF 2018, 2020, 2026"),
]
for i, b in enumerate(bio):
    y = 4.4 + i * 0.98
    rect(sl, 4.1, y, 8.85, 0.88, CARD, line=rgb("3A0A10"))
    htxt(sl, b, 4.25, y+0.1, 8.55, 0.7, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — ANDREW NG PROFILE  (foto besar + bio + quote)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl, rgb("030D12"))
top_bar(sl, BLUE)
rect(sl, 13.08, 0.07, 0.25, 7.36, BLUE)

# Foto besar Andrew Ng kanan
img(sl, "andrew.jpg", 9.48, 0.35, 3.5, 6.8)
rect(sl, 9.48, 0.35, 3.5, 0.55, rgb("031420"))
txt(sl, "ANDREW NG", 9.48, 0.38, 3.5, 0.48,
    size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Big quote
txt(sl, "\u201c", 0.4, 0.2, 1.5, 1.8,
    size=130, bold=True, color=rgb("031A25"))
htxt(sl, "AI adalah listrik baru.\nSeperti listrik mengubah segalanya,\nAI akan mengubah setiap industri.",
     0.5, 0.9, 8.8, 2.3, size=32, bold=True, color=WHITE)
rect(sl, 0.5, 3.3, 6.0, 0.05, BLUE)
htxt(sl, "— Andrew Ng\nPendiri DeepLearning.AI, Coursera, Landing AI",
     0.5, 3.45, 8.5, 0.85, size=15, italic=True, color=BLUE)

bio_ng = [
    ("🏫 Stanford University · Google Brain Founder\n"
     "   Chief Scientist Baidu · Co-founder Coursera"),
    ("🚀 DeepLearning.AI · Landing AI · AI Fund\n"
     "   Kira Learning (K-12 AI Education, 2025)"),
    ("👨‍🎓 8 Juta+ siswa diajar · Time100 AI 2023\n"
     "   Board of Directors Amazon 2024"),
]
for i, b in enumerate(bio_ng):
    y = 4.4 + i * 0.98
    rect(sl, 0.5, y, 8.85, 0.88, CARD, line=rgb("031E2A"))
    htxt(sl, b, 0.65, y+0.1, 8.55, 0.7, size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — HARARI: AI & PENDIDIKAN  (4 cards + chart)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, RED)
rect(sl, 0, 0.07, 0.25, 7.36, RED)

txt(sl, "🔴  HARARI  —  AI dalam Dunia Pendidikan",
    0.4, 0.14, 12.5, 0.58, size=22, bold=True, color=RED)

cards_h = [
    ("⚠️","DAMPAK",RED,
     "Generasi pertama tumbuh\nbersama AI — pola pikir\ndibentuk algoritma."),
    ("🧠","TANTANGAN",rgb("FF6B6B"),
     "AI menulis esai lebih baik.\nTapi PROSES berpikir\nyang hilang."),
    ("🌍","KETIMPANGAN",GOLD,
     "Peluang ini tidak merata:\nnegara kaya makin maju,\nyang lain tertinggal."),
    ("📖","LITERASI BARU",GREEN,
     "Melek AI = bisa bedakan\nnarasi manusia vs\nalgoritma."),
]
for i, (icon, label, col, body) in enumerate(cards_h):
    x = 0.35 + i * 3.2
    rect(sl, x, 0.85, 3.05, 3.5, CARD, line=col)
    rect(sl, x, 0.85, 3.05, 0.52, rgb("0D0D15"))
    txt(sl, f"{icon}  {label}", x+0.1, 0.88, 2.85, 0.46,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, body, x+0.12, 1.46, 2.82, 2.8, size=12.5, color=WHITE)

# Education stats chart bawah
img(sl, "chart_edu_stats.png", 0.35, 4.5, 7.8, 2.85)

# Education market chart kanan bawah
img(sl, "chart_edu_market.png", 8.3, 4.5, 4.85, 2.85)

txt(sl, "💬  Pertanyaan: \"Pernahkah kamu lebih nyaman bertanya ke AI daripada ke guru?\"",
    0.4, 7.15, 12.5, 0.3, size=11, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — HARARI: AI & DUNIA KERJA  (stats + chart WEF)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, RED)
rect(sl, 0, 0.07, 0.25, 7.36, RED)

txt(sl, "🔴  HARARI  —  AI dalam Dunia Kerja",
    0.4, 0.14, 12.5, 0.58, size=22, bold=True, color=RED)

# Big WEF chart
img(sl, "chart_wef_jobs.png", 0.35, 0.85, 6.5, 3.6)

# Disruption+skills chart right
img(sl, "chart_disruption_skills.png", 7.0, 0.85, 6.08, 3.6)

# 3 key points bottom
points = [
    ("🔁","Perubahan KOGNITIF",
     "Bukan hanya fisik — AI gantikan\nhukum, akuntansi, penulisan."),
    ("⚖️","Tanggung Jawab SISTEM",
     "\"Reskilling\" bukan jawaban cukup.\nIni kegagalan sistem, bukan individu."),
    ("🏛️","Regulasi GLOBAL",
     "Butuh tata kelola AI demokratis &\njaminan sosial bagi yang terdampak."),
]
for i, (icon, title, body) in enumerate(points):
    x = 0.35 + i * 4.35
    rect(sl, x, 4.6, 4.15, 2.75, CARD, line=RED)
    txt(sl, f"{icon}  {title}", x+0.12, 4.68, 3.95, 0.44,
        size=13, bold=True, color=RED)
    rect(sl, x+0.12, 5.14, 3.7, 0.04, rgb("3A0A10"))
    htxt(sl, body, x+0.12, 5.22, 3.85, 2.0, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — ANDREW NG: AI & PENDIDIKAN  (4 cards + chart)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, BLUE)
rect(sl, 13.08, 0.07, 0.25, 7.36, BLUE)

txt(sl, "🔵  ANDREW NG  —  AI dalam Dunia Pendidikan",
    0.4, 0.14, 12.5, 0.58, size=22, bold=True, color=BLUE)

cards_ng = [
    ("🎓","DEMOKRATISASI",BLUE,
     "Setiap siswa punya tutor\npribadi 24/7 berkualitas\ntinggi — gratis."),
    ("🤖","AI AGENT GURU",GREEN,
     "Kira Learning (2025):\nAI agents bantu guru\nfokus ke relasi siswa."),
    ("🔓","BARRIER TURUN",GOLD,
     "8M+ orang belajar AI\nlewat Coursera &\nDeepLearning.AI."),
    ("📐","REDESIGN",rgb("B5E3FF"),
     "Plagiarisme? Redesign\npenilaian — uji\npemahaman nyata."),
]
for i, (icon, label, col, body) in enumerate(cards_ng):
    x = 0.35 + i * 3.2
    rect(sl, x, 0.85, 3.05, 3.5, CARD, line=col)
    rect(sl, x, 0.85, 3.05, 0.52, rgb("050D10"))
    txt(sl, f"{icon}  {label}", x+0.1, 0.88, 2.85, 0.46,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, body, x+0.12, 1.46, 2.82, 2.8, size=12.5, color=WHITE)

# Edu stats chart
img(sl, "chart_edu_stats.png", 0.35, 4.5, 7.8, 2.85)
img(sl, "chart_edu_market.png", 8.3, 4.5, 4.85, 2.85)

txt(sl, "💬  \"AI memungkinkan satu guru menjangkau jutaan siswa.\" — Andrew Ng, 2025",
    0.4, 7.15, 12.5, 0.3, size=11, italic=True,
    color=BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — ANDREW NG: AI & DUNIA KERJA  (stats + chart)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, BLUE)
rect(sl, 13.08, 0.07, 0.25, 7.36, BLUE)

txt(sl, "🔵  ANDREW NG  —  AI dalam Dunia Kerja",
    0.4, 0.14, 12.5, 0.58, size=22, bold=True, color=BLUE)

# McKinsey adoption chart kiri
img(sl, "chart_mckinsey_adoption.png", 0.35, 0.85, 6.5, 3.6)

# WEF jobs chart kanan
img(sl, "chart_wef_jobs.png", 7.0, 0.85, 6.08, 3.6)

# 3 skill cards bottom
skills = [
    (BLUE, "97 Juta\npekerjaan BARU",
     "WEF 2025: teknologi ciptakan\nlebih banyak dari yang hilang."),
    (GREEN, "Superpower\nIndividu",
     "Satu orang + AI =\ndampak setara satu tim besar."),
    (GOLD, "3 Kompetensi\nKunci",
     "① AI tools\n② Berpikir sistematis\n③ Soft skills"),
]
for i, (col, title, body) in enumerate(skills):
    x = 0.35 + i * 4.35
    rect(sl, x, 4.6, 4.15, 2.75, CARD, line=col)
    htxt(sl, title, x+0.12, 4.68, 3.9, 0.9,
         size=16, bold=True, color=col)
    rect(sl, x+0.12, 5.6, 3.7, 0.04, rgb("051420"))
    htxt(sl, body, x+0.12, 5.7, 3.9, 1.55, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — DUEL SPLIT-SCREEN: PENDIDIKAN
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0,    0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
rect(sl, 6.5, 0, 0.33, 7.5, BG)

rect(sl, 0,    0, 6.5,  0.62, rgb("2A0810"))
rect(sl, 6.83, 0, 6.5,  0.62, rgb("031420"))
txt(sl, "🔴  HARARI", 0, 0.1, 6.5, 0.5,
    size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(sl, "🔵  ANDREW NG", 6.83, 0.1, 6.5, 0.5,
    size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(sl, "VS", 6.5, 3.3, 0.33, 0.65,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "🎓  AI & PENDIDIKAN", 0, 0.67, 13.33, 0.42,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

duels = [
    ("Dampak\nSiswa",
     "Pola pikir anak dibentuk\nalgoritma tanpa disadari.\nBahaya tersembunyi.",
     "Tutor pribadi untuk semua.\nMempercepat pemahaman.\nBelum pernah ada sebelumnya."),
    ("Berpikir\nKritis",
     "AI mengancam kemampuan\nberpikir mandiri. Proses\nbelajar yang hilang.",
     "Redesign penilaian!\nUji pemahaman nyata,\nbukan kemampuan hafalan."),
    ("Akses\nPendidikan",
     "Peluang tidak merata.\nNegara miskin makin\ntertinggal oleh AI.",
     "Demokratisasi total!\nDesa terpencil bisa akses\nkualitas pendidikan terbaik."),
]
for i, (topic, left, right) in enumerate(duels):
    y = 1.18 + i * 1.98
    rect(sl, 0.18, y, 6.18, 1.82, CARD, line=rgb("3A0A10"))
    rect(sl, 6.98, y, 6.18, 1.82, CARD, line=rgb("031E2A"))
    htxt(sl, left,  0.32, y+0.12, 5.9, 1.58, size=12.5, color=WHITE)
    htxt(sl, right, 7.12, y+0.12, 5.9, 1.58, size=12.5, color=WHITE)
    txt(sl, topic, 4.75, y+0.6, 3.83, 0.62,
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# mini foto di divider
img(sl, "harari.jpg",  0.0,  0.62, 1.2, 2.2)
img(sl, "andrew.jpg", 12.13, 0.62, 1.2, 2.2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — DUEL SPLIT-SCREEN: DUNIA KERJA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0,    0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
rect(sl, 6.5, 0, 0.33, 7.5, BG)

rect(sl, 0,    0, 6.5,  0.62, rgb("2A0810"))
rect(sl, 6.83, 0, 6.5,  0.62, rgb("031420"))
txt(sl, "🔴  HARARI", 0, 0.1, 6.5, 0.5,
    size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(sl, "🔵  ANDREW NG", 6.83, 0.1, 6.5, 0.5,
    size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(sl, "VS", 6.5, 3.3, 0.33, 0.65,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "💼  AI & DUNIA KERJA", 0, 0.67, 13.33, 0.42,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

duels_w = [
    ("Nasib\nPekerjaan",
     "92 juta pekerjaan terhapus.\nRevolusi AI ciptakan\n\"Useless Class\" baru.",
     "170 juta pekerjaan baru!\nSejarah terbukti: teknologi\nselalu ciptakan lebih banyak."),
    ("Siapa\nBertanggung\nJawab?",
     "\"Reskilling\" = alibi korporasi.\nTanggung jawab bukan di\nindividu, tapi sistem.",
     "Mulai dari dirimu!\nJangan tunggu. 1 AI tool/hari\n= karir masa depan."),
    ("Solusi\nYang\nDibutuhkan",
     "Regulasi global demokratis.\nTata kelola AI yang adil.\nJaminan sosial terdampak.",
     "Pendidikan masif & cepat.\nAI for Everyone.\nSemua bisa belajar AI."),
]
for i, (topic, left, right) in enumerate(duels_w):
    y = 1.18 + i * 1.98
    rect(sl, 0.18, y, 6.18, 1.82, CARD, line=rgb("3A0A10"))
    rect(sl, 6.98, y, 6.18, 1.82, CARD, line=rgb("031E2A"))
    htxt(sl, left,  0.32, y+0.12, 5.9, 1.58, size=12.5, color=WHITE)
    htxt(sl, right, 7.12, y+0.12, 5.9, 1.58, size=12.5, color=WHITE)
    txt(sl, topic, 4.75, y+0.5, 3.83, 0.8,
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

img(sl, "harari.jpg",  0.0,  0.62, 1.2, 2.2)
img(sl, "andrew.jpg", 12.13, 0.62, 1.2, 2.2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — RADAR CHART: Peta Perspektif
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, PURPLE)

txt(sl, "📊  PETA PERSPEKTIF  —  Harari vs Andrew Ng",
    0.5, 0.14, 12.33, 0.62, size=24, bold=True,
    color=PURPLE, align=PP_ALIGN.CENTER)

# Radar chart besar tengah
img(sl, "chart_radar.png", 0.35, 0.9, 8.0, 5.5)

# Legend kanan
rect(sl, 8.6, 1.0, 4.48, 5.4, CARD, line=PURPLE)
txt(sl, "Penjelasan Dimensi",
    8.75, 1.1, 4.2, 0.45, size=14, bold=True, color=PURPLE)
rect(sl, 8.75, 1.58, 4.0, 0.04, rgb("2A0A3A"))

dims = [
    (RED,  "Optimisme AI",    "Seberapa positif melihat AI"),
    (RED,  "Urgensi Regulasi","Seberapa perlu aturan ketat"),
    (BLUE, "Peluang Edukasi", "AI sebagai alat pendidikan"),
    (RED,  "Risiko Sosial",   "Bahaya dampak sosial AI"),
    (BLUE, "Kesiapan Indiv.", "Peran individu dalam adaptasi"),
    (BLUE, "Demokratisasi",   "Aksesibilitas AI untuk semua"),
]
for i, (col, dim, desc) in enumerate(dims):
    y = 1.72 + i * 0.78
    rect(sl, 8.75, y, 0.35, 0.38, col)
    txt(sl, dim,  9.18, y+0.02, 3.7, 0.35, size=12, bold=True, color=WHITE)
    txt(sl, desc, 9.18, y+0.38, 3.7, 0.3,  size=10, italic=True, color=GRAY)

rect(sl, 8.75, 6.56, 1.8, 0.4, RED,  line=RED)
txt(sl, "= Harari",  8.85, 6.6, 1.6, 0.3, size=12, bold=True, color=WHITE)
rect(sl, 10.68, 6.56, 1.8, 0.4, BLUE, line=BLUE)
txt(sl, "= Andrew Ng", 10.78, 6.6, 1.6, 0.3, size=12, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — POLL INTERAKTIF: Team siapa?
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, GOLD)

rect(sl, 3.8, 0.18, 5.73, 0.52, rgb("1A1A2E"), line=GOLD)
txt(sl, "🗳️  POLL INTERAKTIF", 3.8, 0.18, 5.73, 0.52,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

htxt(sl, "Setelah melihat data & perspektif,\nkamu lebih setuju dengan siapa?",
     0.5, 0.88, 12.33, 1.28, size=32, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "🎯  Ketik di chat: HARARI / NG / KEDUANYA + alasan singkatmu!",
    0.5, 2.28, 12.33, 0.48, size=15, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

# Team Harari (dengan foto kecil)
rect(sl, 0.45, 2.9, 5.78, 4.0, rgb("2A0810"), line=RED)
img(sl, "harari.jpg", 0.55, 3.05, 1.55, 2.8)
txt(sl, "TEAM HARARI", 2.22, 3.05, 3.9, 0.55,
    size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
htxt(sl, "AI perlu diawasi ketat.\nManusia harus tetap\nmengontrol narasi & sistem.",
     2.22, 3.7, 3.9, 1.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(sl, "🔴  Ketik: HARARI", 2.22, 5.3, 3.9, 0.45,
    size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Team NG (dengan foto kecil)
rect(sl, 7.1, 2.9, 5.78, 4.0, rgb("031420"), line=BLUE)
img(sl, "andrew.jpg", 7.2, 3.05, 1.55, 2.8)
txt(sl, "TEAM ANDREW NG", 8.87, 3.05, 3.9, 0.55,
    size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
htxt(sl, "AI adalah peluang besar.\nBelajar & adaptasi adalah\nkunci — mulai sekarang!",
     8.87, 3.7, 3.9, 1.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(sl, "🔵  Ketik: NG", 8.87, 5.3, 3.9, 0.45,
    size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Tengah: Keduanya
rect(sl, 5.4, 3.15, 1.53, 3.6, CARD, line=GOLD)
txt(sl, "⚡", 5.4, 3.4, 1.53, 0.8,
    size=32, align=PP_ALIGN.CENTER, color=GOLD)
htxt(sl, "KE-\nDUA-\nNYA", 5.4, 4.3, 1.53, 1.5,
     size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "⚡", 5.4, 5.9, 1.53, 0.55,
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — TITIK TEMU  (hal yang keduanya sepakati)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, GREEN)

txt(sl, "🤝  TITIK TEMU  —  Harari & Andrew Ng Ternyata Sepakat!",
    0.5, 0.14, 12.33, 0.62, size=22, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)

# Foto kedua tokoh kecil di header
img(sl, "harari.jpg", 0.3,  0.85, 1.1, 2.0)
img(sl, "andrew.jpg", 11.93, 0.85, 1.1, 2.0)

rect(sl, 1.55, 0.9, 10.23, 1.92, CARD, line=GREEN)
htxt(sl, "Di balik perbedaan yang sangat tajam, Harari & Ng memiliki\nkesepakatan mendasar yang sering tidak disadari publik.",
     1.7, 1.05, 9.9, 1.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)

agrees = [
    ("🧠","Literasi AI adalah WAJIB",
     "Semua orang harus memahami AI. Bukan pilihan, ini keharusan era digital."),
    ("⚠️","AI bisa disalahgunakan",
     "Tanpa etika & regulasi kuat, AI menjadi alat manipulasi & ketimpangan."),
    ("🌱","Soft skills makin VITAL",
     "Kreativitas, empati, kepemimpinan — kompetensi yang tidak bisa diotomasi."),
    ("🔬","Pendidikan harus BEREVOLUSI",
     "Sistem hafalan tanpa pemahaman kritis akan gagal total di era AI."),
    ("⚖️","Manusia tetap harus MEMIMPIN",
     "AI adalah alat. Manusia yang menentukan tujuan, nilai & arahnya."),
]
for i, (icon, title, body) in enumerate(agrees):
    y = 3.0 + i * 0.87
    rect(sl, 0.45, y, 12.43, 0.78, CARD, line=GREEN)
    txt(sl, icon, 0.58, y+0.18, 0.55, 0.42,
        size=20, align=PP_ALIGN.CENTER, color=GREEN)
    txt(sl, title, 1.3, y+0.08, 5.5, 0.38,
        size=13, bold=True, color=WHITE)
    txt(sl, body,  1.3, y+0.46, 11.3, 0.3,
        size=11, italic=True, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — ACTION PLAN  (4 langkah konkret)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
top_bar(sl, GOLD)

txt(sl, "🚀  ACTION PLAN  —  Mulai Dari Kamu Hari Ini",
    0.5, 0.14, 12.33, 0.62, size=24, bold=True,
    color=GOLD, align=PP_ALIGN.CENTER)

actions = [
    (BLUE,  "1", "HARI INI",
     "Coba satu AI tool:\nChatGPT, Gemini,\natau Perplexity."),
    (GREEN, "2", "MINGGU INI",
     "Daftar kursus AI gratis:\nDeepLearning.AI /\nCoursera / Google AI."),
    (GOLD,  "3", "BULAN INI",
     "Identifikasi 3 tugas\nrutinmu yang bisa\ndibantu AI."),
    (RED,   "4", "SELALU",
     "Latih berpikir kritis:\n\"Siapa yang buat?\nApa tujuannya?\""),
]
for i, (col, num, when, act) in enumerate(actions):
    x = 0.35 + i * 3.2
    rect(sl, x, 0.85, 3.05, 5.2, CARD, line=col)
    rect(sl, x, 0.85, 3.05, 1.05, rgb("0D0D15"))
    txt(sl, num,  x+0.1, 0.9, 2.85, 0.54,
        size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, when, x+0.1, 1.44, 2.85, 0.42,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    htxt(sl, act, x+0.15, 2.02, 2.75, 3.8, size=13.5, color=WHITE)

# Key numbers chart
img(sl, "chart_key_numbers.png", 0.35, 6.12, 12.63, 1.25)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING EPIC  (foto kedua + quote penutup)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0,    0, 6.5, 7.5, rgb("0D0508"))
rect(sl, 6.83, 0, 6.5, 7.5, rgb("030D12"))
rect(sl, 6.5, 0, 0.33, 7.5, BG)
top_bar(sl, GOLD)

# Foto closing
img(sl, "harari.jpg",  0.2, 0.28, 3.2, 5.8)
img(sl, "andrew.jpg", 9.93, 0.28, 3.2, 5.8)

# Harari quote kiri
rect(sl, 3.6, 0.28, 2.7, 2.1, rgb("1A0508"), line=RED)
htxt(sl, "\"Siapa yang\nmengontrol AI,\nmengontrol masa depan.\"",
     3.72, 0.4, 2.46, 1.7, size=13, bold=True, color=WHITE)
txt(sl, "— Yuval Noah Harari",
    3.72, 2.0, 2.46, 0.38, size=10, italic=True, color=RED)

# Ng quote kanan
rect(sl, 7.03, 0.28, 2.7, 2.1, rgb("031420"), line=BLUE)
htxt(sl, "\"Siapa yang belajar AI\nakan memimpin\nmasa depan.\"",
     7.15, 0.4, 2.46, 1.7, size=13, bold=True, color=WHITE)
txt(sl, "— Andrew Ng",
    7.15, 2.0, 2.46, 0.38, size=10, italic=True, color=BLUE)

# Central closing message
htxt(sl, "PESAN UNTUK KAMU",
     0.5, 2.6, 12.33, 0.65, size=28, bold=True,
     color=GOLD, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 3.3, 6.33, 0.05, rgb("2A2A3E"))
htxt(sl, "Jangan hanya jadi penonton AI.\nJadi pemain yang cerdas, kritis, dan terus belajar.",
     0.5, 3.45, 12.33, 1.15, size=20, bold=True,
     color=WHITE, align=PP_ALIGN.CENTER)

# 3 closing cards
closes = [(RED,"🧠","Berpikir\nKritis"), (GOLD,"⚡","Terus\nBelajar"), (BLUE,"🚀","Adaptasi\nAktif")]
for i, (col, icon, label) in enumerate(closes):
    x = 2.1 + i * 3.3
    rect(sl, x, 4.75, 2.75, 1.75, CARD, line=col)
    txt(sl, icon, x, 4.88, 2.75, 0.75, size=30, align=PP_ALIGN.CENTER, color=col)
    txt(sl, label, x, 5.65, 2.75, 0.6, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(sl, "\"AI won't replace you — a person using AI will.\"",
    0.5, 6.6, 12.33, 0.45, size=14, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "#WebinarNasionalAI  |  Battle of Minds  |  2025",
    0.5, 7.1, 12.33, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ─── SAVE ──────────────────────────────────────────────────────
out = "/projects/sandbox/sonet-product/Webinar_Interaktif_AI_v2.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Total slides: {len(prs.slides)}")
import os
print(f"   File size   : {os.path.getsize(out)//1024} KB")

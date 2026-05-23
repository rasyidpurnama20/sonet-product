"""
generate_history_ppt.py
Sejarah Kerja, Belajar, Teknologi, dan AI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
ASSETS = Path("/projects/sandbox/sonet-product/assets")

# ── palette ───────────────────────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG     = rgb("0D1117"); CARD  = rgb("161B22"); LINE  = rgb("30363D")
BLUE   = rgb("58A6FF"); RED   = rgb("F78166"); GREEN = rgb("56D364")
GOLD   = rgb("E3B341"); PURP  = rgb("BC8CFF"); WHITE = rgb("E6EDF3")
GRAY   = rgb("8B949E"); DKRED = rgb("FF6B6B")

# ── helpers ───────────────────────────────────────────────────────────────────
def bg(sl, color=None):
    f = sl.background.fill; f.solid()
    f.fore_color.rgb = color or BG

def rect(sl, l, t, w, h, fill, line=None, lw=1.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h, size=15, bold=False,
        color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color or WHITE
    return tb

def htxt(sl, text, l, t, w, h, size=15, bold=False,
         color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color or WHITE
    return tb

def img(sl, fname, l, t, w, h=None):
    p = ASSETS / fname
    if not p.exists(): return
    if h: sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    else: sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w))

def topbar(sl, col): rect(sl, 0, 0, 13.33, 0.08, col); rect(sl, 0, 7.42, 13.33, 0.08, col)
def sidebar(sl, col): rect(sl, 0, 0.08, 0.18, 7.34, col)

def slide_header(sl, title, sub=None, col=BLUE):
    topbar(sl, col); sidebar(sl, col)
    txt(sl, title, 0.35, 0.14, 12.8, 0.72, size=28, bold=True, color=col)
    if sub:
        txt(sl, sub, 0.35, 0.82, 12.8, 0.45, size=14, color=GRAY, italic=True)
    rect(sl, 0.35, 0.82 if not sub else 1.25, 12.6, 0.04, col)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
rect(sl, 0, 0, 0.55, 7.5, BLUE)
rect(sl, 0.55, 0, 12.78, 7.5, rgb("0A0E14"))
topbar(sl, GOLD)

htxt(sl, "Sejarah Kerja,\nBelajar, Teknologi,\ndan AI",
     0.8, 0.6, 9.0, 3.5, size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(sl, "Kaitan Mesin Uap, Listrik, Komputer, Internet, dan AI",
    0.8, 4.2, 10.0, 0.6, size=18, italic=True, color=GOLD)
txt(sl, "dengan Pendapat Para Ahli Dunia",
    0.8, 4.8, 10.0, 0.6, size=18, italic=True, color=GOLD)

rect(sl, 0.8, 5.5, 11.0, 0.05, BLUE)
items = ["Kerja & Kebutuhan Hidup", "Revolusi Teknologi",
         "Dampak pada Pendidikan", "Pekerjaan di Era AI"]
for i, item in enumerate(items):
    x = 0.8 + i * 2.9
    rect(sl, x, 5.7, 2.6, 0.55, rgb("161B22"), line=BLUE)
    txt(sl, item, x+0.1, 5.75, 2.4, 0.45, size=11, bold=True,
        color=BLUE, align=PP_ALIGN.CENTER)
txt(sl, "Webinar Nasional  ·  2025", 0.8, 6.9, 12.0, 0.45,
    size=12, color=GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ALUR BESAR (bagan + penjelasan)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Alur Besar: Mengapa AI Ada?",
             "Perjalanan dari kebutuhan manusia hingga lahirnya AI", col=BLUE)

img(sl, "hist_alur_besar.png", 0.3, 1.4, 4.5, 5.9)

points = [
    (GREEN, "Kebutuhan Hidup",
     "Manusia memiliki kebutuhan dasar: makan, tempat tinggal, keamanan, dan makna sosial."),
    (BLUE,  "Kerja sebagai Jawaban",
     "Kerja lahir sebagai cara manusia memenuhi kebutuhan dan berperan dalam masyarakat."),
    (RED,   "Keterbatasan Manual",
     "Tenaga dan waktu manusia terbatas, produksi tidak bisa terus ditingkatkan secara manual."),
    (GOLD,  "Teknologi sebagai Solusi",
     "Setiap era menciptakan alat baru: mesin uap → listrik → komputer → internet → AI."),
    (PURP,  "AI: Tahap Terbaru",
     "AI berbeda karena bukan hanya membantu otot atau data, tetapi membantu kerja pikiran."),
]
for i, (col, title, body) in enumerate(points):
    y = 1.42 + i * 1.15
    rect(sl, 5.05, y, 8.05, 1.02, CARD, line=col)
    rect(sl, 5.05, y, 0.18, 1.02, col)
    txt(sl, title, 5.32, y+0.05, 7.6, 0.38, size=13, bold=True, color=col)
    htxt(sl, body,  5.32, y+0.44, 7.6, 0.52, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MENGAPA MANUSIA BEKERJA?
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Mengapa Manusia Bekerja?",
             "Kerja bukan hanya soal uang — ada makna yang lebih dalam", col=GREEN)

# Left: 4 kebutuhan boxes
needs = [("🍚 Pangan & Air", GREEN), ("🏠 Tempat Tinggal", BLUE),
         ("👕 Pakaian & Kesehatan", GOLD), ("🛡️ Keamanan & Sosial", PURP)]
for i, (label, col) in enumerate(needs):
    x = 0.35 + (i % 2) * 3.0
    y = 1.45 + (i // 2) * 1.5
    rect(sl, x, y, 2.75, 1.25, CARD, line=col)
    txt(sl, label, x+0.12, y+0.35, 2.5, 0.55,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)

# Right: 4 meaning dimensions
rect(sl, 6.5, 1.3, 6.65, 5.0, CARD, line=GREEN)
txt(sl, "Dimensi Makna Bekerja", 6.7, 1.38, 6.2, 0.5,
    size=14, bold=True, color=GREEN)
rect(sl, 6.7, 1.88, 6.0, 0.04, rgb("30363D"))
dims = [
    ("💰 Ekonomi",   "Menghasilkan uang untuk memenuhi kebutuhan hidup"),
    ("🏅 Identitas",  "Kerja membentuk siapa diri kita di masyarakat"),
    ("🤝 Kontribusi", "Manusia butuh merasa berguna dan dibutuhkan orang lain"),
    ("🌟 Makna",      "Kerja memberi tujuan, arah, dan rasa pencapaian dalam hidup"),
]
for i, (icon, body) in enumerate(dims):
    y = 2.05 + i * 1.05
    rect(sl, 6.65, y, 6.25, 0.9, rgb("0D1117"), line=rgb("30363D"))
    txt(sl, icon, 6.75, y+0.18, 1.4, 0.55, size=13, bold=True, color=GREEN)
    htxt(sl, body, 8.05, y+0.16, 4.7, 0.6, size=11.5, color=WHITE)

# Bottom quote box
rect(sl, 0.35, 5.0, 12.65, 0.95, rgb("161B22"), line=GREEN)
txt(sl, "💡  Inti: Kerja bukan hanya soal uang — kerja berkaitan dengan identitas, "
        "status sosial, kontribusi, dan makna hidup manusia.",
    0.55, 5.1, 12.3, 0.75, size=12.5, italic=True, color=GREEN)

txt(sl, "Kebutuhan Dasar Manusia", 1.2, 1.32, 5.5, 0.45,
    size=13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TIMELINE REVOLUSI TEKNOLOGI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Timeline Revolusi Teknologi",
             "Dari otot manusia hingga kecerdasan buatan", col=GOLD)

img(sl, "hist_timeline.png", 0.25, 1.25, 12.85, 4.5)

# bottom strip
rect(sl, 0.25, 5.9, 12.85, 1.35, CARD, line=GOLD)
txt(sl, "🔑  Pola Sejarah:", 0.45, 5.98, 2.5, 0.4, size=12, bold=True, color=GOLD)
txt(sl, "Setiap era teknologi baru MENGGANTIKAN sebagian pekerjaan lama "
        "dan sekaligus MENCIPTAKAN jenis pekerjaan baru yang sebelumnya tidak ada. "
        "Pola ini berlaku dari mesin uap hingga AI.",
    0.45, 6.38, 12.4, 0.75, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — APA YANG DIBANTU TIAP TEKNOLOGI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Apa yang Dibantu Setiap Teknologi?",
             "AI adalah satu-satunya teknologi yang membantu kerja PIKIRAN", col=PURP)

img(sl, "hist_ai_vs_prev.png", 0.25, 1.3, 12.85, 4.2)

# highlight AI row
rect(sl, 0.25, 5.6, 12.85, 1.65, rgb("1E0F2E"), line=PURP)
txt(sl, "⚡  Mengapa AI Berbeda?", 0.45, 5.68, 5.0, 0.45,
    size=13, bold=True, color=PURP)
htxt(sl, "Mesin uap, listrik, komputer, dan internet membantu manusia "
         "bekerja LEBIH CEPAT atau mengakses informasi LEBIH MUDAH.\n"
         "AI adalah pertama kalinya teknologi membantu manusia BERPIKIR, "
         "MENULIS, MENGANALISIS, dan MENGAMBIL KEPUTUSAN — kerja kognitif yang "
         "selama ini dianggap eksklusif milik manusia.",
    0.45, 6.15, 12.4, 1.0, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — JENIS PEKERJAAN YANG TERPENGARUH (bar chart)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Dampak Teknologi pada Jenis Pekerjaan",
             "Dari kerja fisik ke kerja kognitif — AI masuk paling dalam", col=RED)

img(sl, "hist_replaced.png", 0.25, 1.3, 12.85, 4.55)

rect(sl, 0.25, 5.98, 12.85, 1.3, CARD, line=RED)
htxt(sl, "📌  Komputer mulai masuk ke kerja kognitif (65%), internet memperdalam (80%), "
         "dan AI menembus hampir semua lapisan kerja berpikir (90%).\n"
         "Pola sejarah ini menunjukkan: setiap teknologi baru selalu menggeser "
         "batas kemampuan mesin lebih jauh ke dalam wilayah yang tadinya hanya bisa dilakukan manusia.",
    0.45, 6.06, 12.4, 1.1, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MESIN UAP & LISTRIK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Mesin Uap & Listrik: Revolusi Tenaga Fisik",
             "Era pertama di mana mesin menggantikan otot manusia", col=RED)

for i, (era, col, era_label, points_list) in enumerate([
    ("Mesin Uap\n(1760s)", RED, "Revolusi Industri Pertama", [
        "Tenaga otot mulai digantikan mesin",
        "Pabrik dapat produksi lebih cepat",
        "Kereta & kapal uap percepat transportasi",
        "Lahir: operator mesin, mekanik, insinyur pabrik",
    ]),
    ("Listrik\n(1880s)", GOLD, "Revolusi Industri Kedua", [
        "Lampu: aktivitas melampaui siang hari",
        "Motor listrik gantikan mesin uap",
        "Telegraf & telepon percepat komunikasi",
        "Lahir: teknisi listrik, operator telepon, akuntan",
    ]),
]):
    x = 0.3 + i * 6.5
    rect(sl, x, 1.38, 6.25, 5.85, CARD, line=col)
    rect(sl, x, 1.38, 6.25, 0.72, rgb("0D1117"))
    txt(sl, era, x+0.18, 1.44, 4.0, 0.6, size=18, bold=True, color=col)
    txt(sl, era_label, x+0.18, 1.9, 5.8, 0.45, size=12,
        italic=True, color=GRAY)
    rect(sl, x+0.18, 2.32, 5.75, 0.04, rgb("30363D"))
    for j, pt in enumerate(points_list):
        y = 2.5 + j * 1.12
        rect(sl, x+0.22, y, 5.72, 0.92, rgb("0D1117"), line=LINE)
        txt(sl, f"→  {pt}", x+0.38, y+0.18, 5.4, 0.55,
            size=12.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KOMPUTER & INTERNET
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Komputer & Internet: Revolusi Informasi",
             "Mesin mulai masuk ke wilayah kerja pikiran manusia", col=BLUE)

for i, (era, col, sub, pts) in enumerate([
    ("Komputer\n(1950s–)", BLUE, "Mesin Bantu Pikiran", [
        "Pekerjaan hitung & ketik jauh lebih cepat",
        "Pekerjaan kantor berubah besar-besaran",
        "Yang bisa pakai komputer jadi lebih unggul",
        "Pola: tidak bisa komputer = TERTINGGAL",
    ]),
    ("Internet\n(1990s–)", GREEN, "Informasi Terbuka untuk Semua", [
        "Info tidak lagi hanya di sekolah & buku",
        "Guru bukan satu-satunya sumber pengetahuan",
        "Manusia bisa belajar dari mana saja",
        "Peran guru bergeser → pembimbing cara berpikir",
    ]),
]):
    x = 0.3 + i * 6.5
    rect(sl, x, 1.38, 6.25, 5.85, CARD, line=col)
    rect(sl, x, 1.38, 6.25, 0.72, rgb("0D1117"))
    txt(sl, era, x+0.18, 1.44, 4.0, 0.6, size=18, bold=True, color=col)
    txt(sl, sub, x+0.18, 1.9, 5.8, 0.45, size=12, italic=True, color=GRAY)
    rect(sl, x+0.18, 2.32, 5.75, 0.04, LINE)
    for j, pt in enumerate(pts):
        y = 2.5 + j * 1.12
        rect(sl, x+0.22, y, 5.72, 0.92, rgb("0D1117"), line=LINE)
        txt(sl, f"→  {pt}", x+0.38, y+0.18, 5.4, 0.55, size=12.5, color=WHITE)

# bottom pola sejarah
rect(sl, 0.3, 7.1, 12.7, 0.28, BLUE)
txt(sl, "Pola Sejarah:  tidak bisa komputer → tertinggal  |  "
        "sekarang: tidak pakai AI → berisiko tertinggal",
    0.5, 7.13, 12.3, 0.22, size=11, bold=True, color=rgb("0D1117"),
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI: TAHAP TERBARU
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "AI: Tahap Terbaru dalam Sejarah Teknologi",
             "Yang pertama kali membantu kerja KOGNITIF manusia", col=PURP)

# 5 era pills (top row)
eras = [("Mesin Uap","Tenaga\nFisik",RED),
        ("Listrik","Energi\nCepat",GOLD),
        ("Komputer","Data &\nAdmin",BLUE),
        ("Internet","Akses\nInfo",GREEN),
        ("AI","Berpikir\n& Belajar",PURP)]
for i,(name,sub,col) in enumerate(eras):
    x = 0.35 + i*2.55
    rect(sl, x, 1.3, 2.3, 1.5, CARD, line=col)
    txt(sl, name, x+0.12, 1.38, 2.05, 0.5, size=13, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    txt(sl, sub,  x+0.12, 1.88, 2.05, 0.78, size=11,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", 2.58+i*2.55, 1.98, 0.28, 0.38,
            size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# AI special highlight
rect(sl, 0.35, 3.0, 12.65, 3.0, rgb("1E1030"), line=PURP, lw=2.5)
txt(sl, "Mengapa AI Berbeda dari Semua Teknologi Sebelumnya?",
    0.6, 3.08, 12.0, 0.55, size=16, bold=True, color=PURP)

diffs = [
    (PURP, "Membantu BERPIKIR",
     "AI bisa menganalisis data kompleks, membuat koneksi, dan menarik kesimpulan"),
    (BLUE, "Membantu MENULIS",
     "AI menghasilkan teks, laporan, esai, kode program, dan konten kreatif"),
    (GREEN,"Membantu BELAJAR",
     "AI menjelaskan konsep secara personal, interaktif, dan sesuai kecepatan belajar"),
    (GOLD, "Membantu MEMUTUSKAN",
     "AI memberikan rekomendasi berbasis data untuk pengambilan keputusan lebih baik"),
]
for i,(col,title,body) in enumerate(diffs):
    x = 0.55 + (i%2)*6.25
    y = 3.72 + (i//2)*1.1
    rect(sl, x, y, 6.0, 0.95, rgb("0D1117"), line=col)
    rect(sl, x, y, 0.14, 0.95, col)
    txt(sl, title, x+0.25, y+0.06, 2.8, 0.38, size=12, bold=True, color=col)
    txt(sl, body,  x+0.25, y+0.45, 5.6, 0.44, size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PENDIDIKAN SEBELUM & SESUDAH AI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Pendidikan: Sebelum AI vs Setelah AI",
             "Pergeseran dari pendidikan massal ke pembelajaran personal", col=GREEN)

img(sl, "hist_edu_compare.png", 0.25, 1.28, 12.85, 5.1)

rect(sl, 0.25, 6.48, 12.85, 0.82, rgb("0F1F15"), line=GREEN)
txt(sl, "💡  Internet membuka akses informasi, tetapi AI membuat informasi "
        "dapat dijelaskan secara personal dan interaktif — persis seperti "
        "memiliki tutor pribadi yang sabar, tersedia 24/7, dan tidak menghakimi.",
    0.45, 6.56, 12.4, 0.65, size=12, italic=True, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ANDREW NG
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
topbar(sl, BLUE); sidebar(sl, BLUE)

# profile strip
rect(sl, 0.25, 0.15, 12.85, 1.5, rgb("0A1628"))
img(sl, "andrew.jpg", 0.35, 0.2, 1.3, 1.35)
txt(sl, "Andrew Ng", 1.82, 0.25, 8.0, 0.58, size=26, bold=True, color=BLUE)
htxt(sl, "Pendiri DeepLearning.AI & Coursera  ·  Stanford University\n"
         "Co-founder Google Brain  ·  Mengajar 8 juta+ pelajar",
     1.82, 0.85, 10.0, 0.72, size=12, color=GRAY)

# quote cards
quotes = [
    (BLUE, "AI & Listrik — Pekerjaan",
     '"AI is the new electricity. Just as electricity transformed almost\n'
     'everything 100 years ago, today I actually have a hard time thinking\n'
     'of an industry that I don\'t think AI will transform."',
     "Stanford GSB / Live Mint, 2026"),
    (GREEN,"AI Redefine Guru — Pendidikan",
     '"AI is helping redefine what it means to be a great teacher."',
     "Business Insider, April 2025"),
]
for i,(col,title,quote,src) in enumerate(quotes):
    y = 1.82 + i*2.65
    rect(sl, 0.25, y, 12.85, 2.45, CARD, line=col, lw=2)
    rect(sl, 0.25, y, 0.22, 2.45, col)
    txt(sl, title, 0.58, y+0.1, 12.0, 0.45, size=14, bold=True, color=col)
    rect(sl, 0.58, y+0.58, 12.0, 0.03, LINE)
    htxt(sl, quote, 0.58, y+0.72, 12.0, 1.38, size=14, italic=True, color=WHITE)
    txt(sl, f"— {src}", 0.58, y+2.08, 11.8, 0.3,
        size=10, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

# key takeaway
rect(sl, 0.25, 7.15, 12.85, 0.25, BLUE)
txt(sl, "Inti: AI bukan pengganti manusia — AI adalah senjata bagi yang mau belajar dan beradaptasi.",
    0.45, 7.17, 12.4, 0.2, size=11, bold=True,
    color=rgb("0D1117"), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — JENSEN HUANG
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
topbar(sl, GOLD); sidebar(sl, GOLD)

rect(sl, 0.25, 0.15, 12.85, 1.5, rgb("1A1400"))
txt(sl, "Jensen Huang", 0.45, 0.25, 8.0, 0.58, size=26, bold=True, color=GOLD)
htxt(sl, "CEO NVIDIA  ·  Salah satu pemimpin industri AI paling berpengaruh\n"
         "NVIDIA: perusahaan chip AI senilai $3+ triliun (2025)",
     0.45, 0.85, 11.8, 0.72, size=12, color=GRAY)

quotes_j = [
    (GOLD, "AI & Pekerjaan — Peringatan Keras",
     '"Every job will be affected, and immediately. It is unquestionable.\n'
     'You\'re not going to lose your job to an AI, but you\'re going to\n'
     'lose your job to someone who uses AI."',
     "CNBC, 28 Mei 2025"),
    (GREEN,"AI Tutor — Pendidikan",
     '"If there\'s one thing I would encourage everybody to do,\n'
     'go get yourself an AI tutor right away."',
     "CNBC, Februari 2025"),
    (BLUE, "Setiap Lulusan Harus Kuasai AI",
     '"Every college student should graduate and be an expert in AI.\n'
     'This skill could determine whether someone secures a job or\n'
     'struggles in an AI-driven economy."',
     "Money Control, 2025"),
]
for i,(col,title,quote,src) in enumerate(quotes_j):
    y = 1.82 + i*1.8
    rect(sl, 0.25, y, 12.85, 1.65, CARD, line=col, lw=1.8)
    rect(sl, 0.25, y, 0.2,  1.65, col)
    txt(sl, title, 0.55, y+0.08, 12.1, 0.42, size=13, bold=True, color=col)
    htxt(sl, quote, 0.55, y+0.55, 12.0, 0.95, size=13, italic=True, color=WHITE)
    txt(sl, f"— {src}", 0.55, y+1.35, 12.1, 0.25,
        size=10, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

rect(sl, 0.25, 7.15, 12.85, 0.25, GOLD)
txt(sl, "Inti: Yang kalah bukan orang yang digantikan AI — melainkan orang yang TIDAK menggunakan AI.",
    0.45, 7.17, 12.4, 0.2, size=11, bold=True,
    color=rgb("0D1117"), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ELON MUSK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
topbar(sl, RED); sidebar(sl, RED)

rect(sl, 0.25, 0.15, 12.85, 1.5, rgb("1A0808"))
txt(sl, "Elon Musk", 0.45, 0.25, 8.0, 0.58, size=26, bold=True, color=RED)
htxt(sl, "CEO Tesla, SpaceX, X  ·  Salah satu pendiri OpenAI (kemudian keluar)\n"
         "Tokoh paling vokal soal dampak AI pada masa depan pekerjaan",
     0.45, 0.85, 11.8, 0.72, size=12, color=GRAY)

quotes_e = [
    (RED, "Pekerjaan Menjadi Opsional",
     '"In a benign scenario, probably none of us will have a job.\n'
     'But in that benign scenario there will be universal high income —\n'
     'not universal base income — there will be no shortage of goods or services."',
     "VivaTech 2024, Paris — Fortune, Mei 2024"),
    (GOLD,"Universal Basic Income (UBI)",
     '"I think we\'ll end up doing universal basic income. It\'s going to be\n'
     'necessary. There will be fewer and fewer jobs that a robot cannot do better.\n'
     'These are not things I wish will happen — these are things I think probably will happen."',
     "Business Insider, Juni 2024"),
    (PURP,"Kerja Jadi Pilihan",
     '"My prediction is that work will be optional."',
     "U.S.-Saudi Investment Forum, Washington — Fortune, Januari 2026"),
]
for i,(col,title,quote,src) in enumerate(quotes_e):
    y = 1.82 + i*1.8
    rect(sl, 0.25, y, 12.85, 1.65, CARD, line=col, lw=1.8)
    rect(sl, 0.25, y, 0.2,  1.65, col)
    txt(sl, title, 0.55, y+0.08, 12.1, 0.42, size=13, bold=True, color=col)
    htxt(sl, quote, 0.55, y+0.55, 12.0, 0.95, size=13, italic=True, color=WHITE)
    txt(sl, f"— {src}", 0.55, y+1.35, 12.1, 0.25,
        size=10, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

rect(sl, 0.25, 7.15, 12.85, 0.25, RED)
txt(sl, "Inti: Jika AI membuat manusia tidak wajib bekerja, sistem ekonomi perlu menyiapkan UBI.",
    0.45, 7.17, 12.4, 0.2, size=11, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — YUVAL HARARI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
topbar(sl, DKRED); sidebar(sl, DKRED)

rect(sl, 0.25, 0.15, 12.85, 1.5, rgb("1A0508"))
img(sl, "harari.jpg", 0.35, 0.2, 1.3, 1.35)
txt(sl, "Yuval Noah Harari", 1.82, 0.25, 9.0, 0.58, size=26, bold=True, color=DKRED)
htxt(sl, "Sejarawan & Penulis  ·  Hebrew University Jerusalem\n"
         "Buku: Sapiens (2011), Homo Deus (2016), 21 Lessons (2018), Nexus (2024)",
     1.82, 0.85, 10.8, 0.72, size=12, color=GRAY)

quotes_h = [
    (DKRED,"AI Masuk Sistem Berbasis Bahasa",
     '"If laws are made of words, then AI will take over the legal system.\n'
     'If books are just combinations of words, then AI will take over books.\n'
     'If religion is built from words, then AI will take over religion."',
     "Davos / WEF, dikutip Newsweek, Januari 2025"),
    (GOLD, "AI & Akselerasi Tanpa Henti",
     '"AI is a knife that can invent new kinds of knives as well as new kinds\n'
     'of music, medicine and money. The point is not just novelty. It is acceleration."',
     "Davos 2026, dikutip Forbes, Januari 2026"),
    (PURP, "Manfaatnya akan terasa dalam 200 tahun",
     '"AI\'s true impact will unfold over 200 years.\n'
     'Judging AI on short-term progress badly misunderstands its scale and risks."',
     "Business Insider, Januari 2026"),
]
for i,(col,title,quote,src) in enumerate(quotes_h):
    y = 1.82 + i*1.8
    rect(sl, 0.25, y, 12.85, 1.65, CARD, line=col, lw=1.8)
    rect(sl, 0.25, y, 0.2,  1.65, col)
    txt(sl, title, 0.55, y+0.08, 12.1, 0.42, size=13, bold=True, color=col)
    htxt(sl, quote, 0.55, y+0.55, 12.0, 0.95, size=13, italic=True, color=WHITE)
    txt(sl, f"— {src}", 0.55, y+1.35, 12.1, 0.25,
        size=10, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

rect(sl, 0.25, 7.15, 12.85, 0.25, DKRED)
txt(sl, "Inti: AI sangat kuat karena manusia membangun peradaban dengan BAHASA — dan AI menguasai bahasa.",
    0.45, 7.17, 12.4, 0.2, size=11, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SAL KHAN & BILL GATES & SAM ALTMAN
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Para Ahli tentang AI & Pendidikan",
             "Sepakat: AI adalah transformasi terbesar pendidikan sepanjang sejarah", col=GREEN)

experts = [
    (GREEN, "Sal Khan", "Pendiri Khan Academy",
     '"We\'re at the cusp of using AI for probably the biggest positive\n'
     'transformation that education has ever seen."\n\n'
     '"AI in the classroom will be like having four or five helpful\n'
     'grad students — acting as teachers\' assistants."',
     "TED Talk 2023 & Business Insider 2025"),
    (BLUE,  "Bill Gates", "Co-founder Microsoft",
     '"The AI will be like a great high school teacher who really\n'
     'marks your essay, and you go back and think:\n'
     '\'OK, I need to step up there.\'"',
     "CNBC, Agustus 2023"),
    (PURP,  "Sam Altman", "CEO OpenAI",
     '"Our children will have virtual tutors who can provide personalized\n'
     'instruction in any subject, in any language, and at whatever pace\n'
     'they need."',
     "samaltman.com — The Intelligence Age, September 2024"),
]
for i,(col,name,role,quote,src) in enumerate(experts):
    x = 0.25 + i * 4.37
    rect(sl, x, 1.4, 4.18, 5.85, CARD, line=col, lw=2)
    rect(sl, x, 1.4, 4.18, 0.9, rgb("0D1117"))
    txt(sl, name, x+0.15, 1.47, 3.85, 0.48, size=16, bold=True, color=col)
    txt(sl, role, x+0.15, 1.92, 3.85, 0.35, size=10, italic=True, color=GRAY)
    rect(sl, x+0.15, 2.3, 3.72, 0.04, LINE)
    htxt(sl, quote, x+0.15, 2.42, 3.85, 3.65, size=11.5,
         italic=True, color=WHITE)
    txt(sl, f"— {src}", x+0.1, 5.95, 3.9, 0.25,
        size=9, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

rect(sl, 0.25, 7.15, 12.85, 0.25, GREEN)
txt(sl, "Sepakat: AI tutor personal untuk semua siswa — demokratisasi pendidikan terbesar dalam sejarah.",
    0.45, 7.17, 12.4, 0.2, size=11, bold=True,
    color=rgb("0D1117"), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — BAGAN ALUR: PEKERJAAN MENURUT PARA AHLI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Alur Pekerjaan Menurut Para Ahli",
             "Dari sejarah teknologi hingga prediksi masa depan pekerjaan", col=RED)

img(sl, "hist_work_flow.png", 0.25, 1.3, 12.85, 3.7)

# Detail cards bawah
details = [
    (BLUE, "Andrew Ng", "AI adalah senjata bagi yang mau belajar"),
    (GOLD, "Jensen Huang", "Yang tertinggal = yang tidak pakai AI"),
    (DKRED,"Yuval Harari", "AI masuk sistem berbasis bahasa"),
    (RED,  "Elon Musk", "Pekerjaan bisa jadi opsional + UBI"),
]
for i,(col,name,inti) in enumerate(details):
    x = 0.25 + i * 3.27
    rect(sl, x, 5.22, 3.08, 1.95, CARD, line=col)
    rect(sl, x, 5.22, 3.08, 0.48, LINE)
    txt(sl, name, x+0.12, 5.27, 2.82, 0.38, size=13, bold=True, color=col)
    htxt(sl, inti, x+0.12, 5.74, 2.84, 1.2, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — BAGAN ALUR: PENDIDIKAN MENURUT PARA AHLI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Alur Pendidikan Menurut Para Ahli",
             "Konsensus: AI = transformasi personal terbesar dalam sejarah pendidikan", col=GREEN)

img(sl, "hist_edu_flow.png", 0.25, 1.3, 12.85, 3.7)

details_e = [
    (BLUE,  "Andrew Ng",   "AI definisikan ulang\narti guru hebat"),
    (GOLD,  "Jensen Huang","Segera gunakan\nAI tutor"),
    (GREEN, "Bill Gates",  "AI tutor = sebaik\nguru terbaik"),
    (GREEN, "Sal Khan",    "Transformasi\nterbesar pendidikan"),
    (PURP,  "Sam Altman",  "Generasi paling\nberuntung sejarah"),
]
box_w = 2.42
for i,(col,name,inti) in enumerate(details_e):
    x = 0.25 + i * 2.58
    rect(sl, x, 5.18, box_w, 2.08, CARD, line=col)
    rect(sl, x, 5.18, box_w, 0.48, LINE)
    txt(sl, name, x+0.1, 5.23, 2.2, 0.38, size=12, bold=True, color=col)
    htxt(sl, inti, x+0.1, 5.7, 2.24, 1.35, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — EXPERT CARDS (all 5)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Pendapat Para Ahli — Ringkasan",
             "5 tokoh global tentang AI, pekerjaan, dan pendidikan", col=PURP)

img(sl, "hist_expert_cards.png", 0.25, 1.3, 12.85, 3.7)

# table of key takeaways
headers = ["Tokoh","Tentang Pekerjaan","Tentang Pendidikan"]
col_x   = [0.3, 3.5, 8.3]
rect(sl, 0.25, 5.15, 12.85, 0.48, CARD)
for hx, htxt_s in zip(col_x, headers):
    txt(sl, htxt_s, hx+0.1, 5.2, 4.5, 0.38, size=12, bold=True, color=PURP)
rect(sl, 0.3, 5.63, 12.7, 0.03, PURP)

rows = [
    ("Andrew Ng",    "AI = senjata bagi pelajar",         "AI definisikan ulang guru"),
    ("Jensen Huang", "Tidak pakai AI = tertinggal",        "Segera pakai AI tutor"),
    ("Elon Musk",    "Kerja akan opsional, perlu UBI",     "—"),
    ("Yuval Harari", "AI masuk sistem berbasis bahasa",    "Literasi kritis sangat penting"),
    ("Sal Khan",     "—",                                  "Transformasi terbesar pendidikan"),
]
colors = [BLUE, GOLD, RED, DKRED, GREEN]
for i,(name,job,edu) in enumerate(rows):
    y = 5.7 + i * 0.37
    bg_c = CARD if i%2==0 else rgb("0D1117")
    rect(sl, 0.25, y, 12.85, 0.35, bg_c)
    col = colors[i]
    txt(sl, name, col_x[0]+0.1, y+0.04, 2.8, 0.28, size=11, bold=True, color=col)
    txt(sl, job,  col_x[1]+0.1, y+0.04, 4.4, 0.28, size=11, color=WHITE)
    txt(sl, edu,  col_x[2]+0.1, y+0.04, 4.4, 0.28, size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — BAGAN BESAR KESELURUHAN
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
slide_header(sl, "Bagan Besar: Manusia → Teknologi → Adaptasi",
             "Satu alur utuh dari kebutuhan hingga era AI", col=GOLD)

img(sl, "hist_mega_flow.png", 5.0, 1.3, 8.08, 6.0)

# Left annotation
annotations = [
    (GREEN, "Akar",
     "Kebutuhan hidup mendorong manusia untuk bekerja sejak zaman purba."),
    (BLUE,  "Pola",
     "Setiap era teknologi menciptakan pekerjaan baru yang lebih kompleks."),
    (GOLD,  "AI",
     "AI bukan awal baru — ia adalah kelanjutan dari pola panjang ini."),
    (PURP,  "Tugas kita",
     "Belajar ulang, beradaptasi, dan memastikan manusia tetap memimpin AI."),
]
for i,(col,title,body) in enumerate(annotations):
    y = 1.42 + i * 1.45
    rect(sl, 0.28, y, 4.5, 1.28, CARD, line=col)
    rect(sl, 0.28, y, 0.18, 1.28, col)
    txt(sl, title, 0.55, y+0.08, 3.9, 0.42, size=13, bold=True, color=col)
    htxt(sl, body,  0.55, y+0.52, 3.9, 0.68, size=11.5, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — KESIMPULAN UTAMA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); bg(sl)
topbar(sl, GOLD); rect(sl, 0, 7.42, 13.33, 0.08, GOLD)

# big title
txt(sl, "Kesimpulan Utama", 0.5, 0.15, 12.33, 0.72,
    size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
rect(sl, 1.5, 0.9, 10.33, 0.05, GOLD)

# 5 tech pills — reuse pattern
eras2 = [("Mesin Uap","Tenaga Fisik",RED),
         ("Listrik","Energi",GOLD),
         ("Komputer","Data",BLUE),
         ("Internet","Informasi",GREEN),
         ("AI","Kerja Pikiran",PURP)]
for i,(name,sub,col) in enumerate(eras2):
    x = 0.35 + i*2.55
    rect(sl, x, 1.05, 2.3, 1.3, CARD, line=col)
    txt(sl, name, x+0.12, 1.12, 2.05, 0.45, size=12, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    txt(sl, sub,  x+0.12, 1.57, 2.05, 0.65, size=10,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i<4:
        txt(sl, "→", 2.58+i*2.55, 1.58, 0.28, 0.35,
            size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# 4 conclusion boxes
conclusions = [
    (BLUE,  "AI adalah Kelanjutan Sejarah",
     "AI bukan revolusi yang datang tiba-tiba. Ia adalah bagian dari pola panjang "
     "teknologi yang selalu mengubah cara manusia bekerja dan belajar."),
    (GREEN, "Bukan Ancaman, Tapi Alat",
     "Menurut para ahli, AI adalah alat yang memperkuat manusia — "
     "terutama bagi mereka yang mau belajar dan beradaptasi."),
    (GOLD,  "Pendidikan Harus Berevolusi",
     "AI memungkinkan setiap siswa memiliki tutor personal. "
     "Sistem pendidikan perlu menyesuaikan diri dengan realitas ini."),
    (PURP,  "Manusia Harus Tetap Memimpin",
     "AI adalah alat. Manusia yang menentukan tujuan, nilai, dan arah "
     "penggunaannya. Berpikir kritis tetap menjadi kompetensi utama."),
]
for i,(col,title,body) in enumerate(conclusions):
    x = 0.3 + (i%2)*6.53
    y = 2.55 + (i//2)*2.35
    rect(sl, x, y, 6.25, 2.18, CARD, line=col, lw=2)
    rect(sl, x, y, 6.25, 0.56, LINE)
    rect(sl, x, y, 0.2,  2.18, col)
    txt(sl, title, x+0.3, y+0.1, 5.75, 0.45, size=14, bold=True, color=col)
    htxt(sl, body,  x+0.3, y+0.65, 5.75, 1.35, size=12.5, color=WHITE)

# final quote
rect(sl, 0.3, 7.08, 12.73, 0.3, rgb("1A1400"), line=GOLD)
txt(sl, "\"Menurut para ahli, AI bukan sekadar ancaman — AI adalah alat besar yang dapat memperkuat manusia "
        "bagi mereka yang mau belajar, memakai AI, dan beradaptasi.\"",
    0.5, 7.11, 12.3, 0.24, size=10.5, italic=True,
    color=GOLD, align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
OUT = "/projects/sandbox/sonet-product/Sejarah_Teknologi_AI.pptx"
prs.save(OUT)
import os
print(f"✅  Saved: {OUT}")
print(f"   Slides : {len(prs.slides)}")
print(f"   Size   : {os.path.getsize(OUT)//1024} KB")

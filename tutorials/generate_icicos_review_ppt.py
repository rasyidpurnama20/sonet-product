"""
Generate ICICOS Review Paper Tutorial - PowerPoint Presentation
Requires: python-pptx (pip install python-pptx)
Output: ICICOS_Review_Tutorial.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────
C_NAVY    = RGBColor(0x1F, 0x38, 0x64)   # dark navy - header bg
C_BLUE    = RGBColor(0x27, 0x6F, 0xBF)   # mid blue  - accents
C_LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)   # light blue - table alt rows
C_ORANGE  = RGBColor(0xE0, 0x7B, 0x39)   # orange    - highlights
C_GREEN   = RGBColor(0x2E, 0x86, 0x48)   # green     - positive items
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # red       - negative items
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)   # light gray - content bg
C_DGRAY   = RGBColor(0x44, 0x44, 0x44)   # dark gray  - body text
C_YELLOW  = RGBColor(0xFF, 0xF0, 0xCC)   # soft yellow - callout bg

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_DGRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=16, bold=False, color=C_DGRAY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False, level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=None):
    """Top navy bar with title and optional subtitle."""
    h = 1.15 if subtitle else 0.95
    add_rect(slide, 0, 0, 13.33, h, fill=C_NAVY)
    add_rect(slide, 0, h, 13.33, 0.06, fill=C_ORANGE)
    add_text(slide, title, 0.35, 0.08, 12.5, 0.65,
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.68, 12.5, 0.38,
                 size=15, color=C_LBLUE)


def bullet_box(slide, items, l, t, w, h, title=None,
               title_color=C_NAVY, bullet="▸", size=16,
               bg=None, border=None):
    """A box with optional title and bulleted items."""
    if bg:
        add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(1))
    if title:
        add_text(slide, title, l+0.12, t+0.08, w-0.2, 0.38,
                 size=17, bold=True, color=title_color)
        item_t = t + 0.46
    else:
        item_t = t + 0.14
    txb = slide.shapes.add_textbox(
        Inches(l+0.18), Inches(item_t),
        Inches(w-0.28), Inches(h - (0.46 if title else 0.2)))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = C_DGRAY
    return txb


def callout(slide, text, l, t, w, h, bg=C_YELLOW, border=C_ORANGE,
            size=16, bold=False, align=PP_ALIGN.LEFT, color=C_DGRAY):
    add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(2))
    add_text(slide, text, l+0.15, t+0.12, w-0.25, h-0.18,
             size=size, bold=bold, color=color, align=align)



# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 4.7, 13.33, 2.8, fill=C_BLUE)
add_rect(sl, 0, 4.62, 13.33, 0.1, fill=C_ORANGE)

add_text(sl, "Review Paper ICICOS", 0.7, 1.0, 11.9, 1.1,
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "dengan ChatGPT Plus", 0.7, 1.95, 11.9, 0.8,
         size=36, bold=False, color=C_LBLUE, align=PP_ALIGN.CENTER)
add_rect(sl, 3.5, 2.95, 6.33, 0.06, fill=C_ORANGE)
add_text(sl, "Panduan Praktis untuk Dosen Reviewer", 0.7, 3.15, 11.9, 0.55,
         size=20, color=C_LBLUE, align=PP_ALIGN.CENTER, italic=True)
add_text(sl, "Berbasis Keahlian · Terstruktur · Efisien", 0.7, 3.7, 11.9, 0.5,
         size=17, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "ICICOS 2026  ·  Juni 2026", 0.7, 5.1, 11.9, 0.5,
         size=16, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "30–45 menit per paper · Kualitas lebih tinggi · Lebih konsisten",
         0.7, 5.7, 11.9, 0.45,
         size=15, color=C_LBLUE, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "Agenda", "Dari keahlian kamu → review berkualitas dalam 30–45 menit")

agenda = [
    ("01", "Review = Cerminan Keahlian Kamu",    "Skenario & contoh nyata"),
    ("02", "Setup ChatGPT Plus + Project",         "Encode style kamu satu kali"),
    ("03", "Ekstraksi Paper: Langkah Pertama",      "Struktur sebelum analisis"),
    ("04", "AI = Pelaksana Tugas",                  "Pertanyaan vs. tugas yang jelas"),
    ("05", "Alur Lengkap 5 Langkah",               "Dari upload ke draft review"),
    ("06", "Bias, Variasi & Konsistensi",           "Kualitas terjaga di setiap paper"),
]
for i, (num, title, sub) in enumerate(agenda):
    row = i // 2; col = i % 2
    lx = 0.3 + col * 6.55
    ty = 1.35 + row * 1.85
    add_rect(sl, lx, ty, 6.2, 1.65, fill=C_WHITE, line=C_LBLUE, line_w=Pt(1))
    add_rect(sl, lx, ty, 0.7, 1.65, fill=C_NAVY)
    add_text(sl, num, lx+0.04, ty+0.5, 0.62, 0.65,
             size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, lx+0.82, ty+0.18, 5.2, 0.65,
             size=16, bold=True, color=C_NAVY)
    add_text(sl, sub,  lx+0.82, ty+0.82, 5.2, 0.55,
             size=13, color=C_DGRAY, italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — REVIEW = CERMINAN KEAHLIAN
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "01 · Review = Cerminan Keahlian Kamu",
           "Reviewer yang berbeda → komentar yang berbeda → semua valid")

cols = [
    ("🤖  Ahli ML / AI",    C_NAVY,  [
        "Baseline vs. SOTA terbaru",
        "Ablation study ada?",
        "Dataset leakage / bias?",
        "Hyperparameter justified?",
    ]),
    ("⚙️  Ahli Sistem",     C_BLUE,  [
        "Skalabilitas & performa",
        "Latensi, throughput, memory",
        "Reproducibility & open code",
        "Deployment realistic?",
    ]),
    ("📊  Ahli IS / Aplikasi", RGBColor(0x2E,0x86,0x48), [
        "Masalah nyata atau akademis?",
        "User study ada?",
        "Bukan sekadar prototype?",
        "Business value jelas?",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    lx = 0.28 + i * 4.32
    add_rect(sl, lx, 1.35, 4.1, 4.85, fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, lx, 1.35, 4.1, 0.62, fill=color)
    add_text(sl, title, lx+0.12, 1.4, 3.9, 0.52,
             size=16, bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        add_text(sl, f"▸  {item}", lx+0.18, 2.12+j*0.85, 3.82, 0.72,
                 size=15, color=C_DGRAY)

callout(sl,
    "Kunci: ChatGPT perlu tahu perspektif & keahlian KAMU sebelum mulai bekerja.",
    0.28, 6.3, 12.78, 0.72,
    size=16, bold=True, align=PP_ALIGN.CENTER, color=C_NAVY)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — SKENARIO NYATA
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "01 · Skenario: Satu Paper, Tiga Reviewer",
           'Paper: "Deep Learning-based Intrusion Detection System for IoT Networks"')

rows = [
    ("Ahli Security",      C_NAVY,
     '"Dataset hanya 4 tipe serangan — tidak representatif untuk IoT production"'),
    ("Ahli Deep Learning", C_BLUE,
     '"Tidak ada ablation study — kontribusi layer attention tidak dibuktikan terpisah"'),
    ("Ahli IoT Systems",   RGBColor(0x2E,0x86,0x48),
     '"Model 47MB tidak bisa di-deploy di microcontroller — tidak ada evaluasi resource overhead"'),
]
for i, (role, color, comment) in enumerate(rows):
    ty = 1.45 + i * 1.75
    add_rect(sl, 0.3, ty, 2.6, 1.5, fill=color)
    add_text(sl, role, 0.38, ty+0.42, 2.44, 0.65,
             size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.9, ty, 10.1, 1.5, fill=C_WHITE, line=color, line_w=Pt(1.5))
    add_text(sl, comment, 3.08, ty+0.28, 9.75, 0.95,
             size=16, color=C_DGRAY, italic=True)

callout(sl,
    "Semua komentar valid — muncul dari perspektif berbeda.  "
    "Itulah mengapa setup keahlian di awal sangat menentukan kualitas review.",
    0.3, 6.7, 12.73, 0.55,
    size=14, bold=False, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — SETUP PROJECT
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "02 · Setup ChatGPT Plus + Project",
           "Encode keahlian kamu sekali — aktif di semua sesi review")

# Left: comparison table
add_rect(sl, 0.28, 1.3, 6.0, 5.85, fill=C_WHITE, line=C_LBLUE, line_w=Pt(1))
add_rect(sl, 0.28, 1.3, 6.0, 0.52, fill=C_NAVY)
add_text(sl, "Chat Biasa  vs.  Project", 0.38, 1.34, 5.8, 0.44,
         size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

comparisons = [
    ("❌ Konteks hilang tiap sesi",      "✅ Konteks permanen"),
    ("❌ Jelaskan ulang setiap kali",    "✅ Setup sekali untuk semua paper"),
    ("❌ File hilang setelah sesi",      "✅ File tersimpan & bisa dirujuk"),
    ("❌ Tidak bisa dikustomisasi",      "✅ Custom instructions = DNA reviewer"),
    ("❌ Review tidak konsisten",        "✅ Standar sama di semua paper"),
]
for i, (bad, good) in enumerate(comparisons):
    ty = 1.96 + i * 0.84
    bg = C_LGRAY if i % 2 == 0 else C_WHITE
    add_rect(sl, 0.28, ty, 6.0, 0.84, fill=bg)
    add_text(sl, bad,  0.38, ty+0.14, 2.72, 0.56, size=13, color=C_RED)
    add_text(sl, good, 3.2,  ty+0.14, 2.9,  0.56, size=13, color=C_GREEN, bold=True)

# Right: 3 setup steps
add_rect(sl, 6.62, 1.3, 6.42, 5.85, fill=C_WHITE, line=C_BLUE, line_w=Pt(1))
add_rect(sl, 6.62, 1.3, 6.42, 0.52, fill=C_BLUE)
add_text(sl, "3 Langkah Setup", 6.72, 1.34, 6.2, 0.44,
         size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

steps = [
    ("1", "Buat Project",
     "Sidebar kiri → + Projects\nNama: ICICOS Reviewer — [Nama Kamu]"),
    ("2", "Isi Custom Instructions",
     "Bidang keahlian, prioritas review,\ngaya komentar, format output"),
    ("3", "Upload File Pendukung",
     "Form review ICICOS, template ekstraksi,\nauthor guidelines"),
]
for i, (num, title, desc) in enumerate(steps):
    ty = 1.96 + i * 1.72
    add_rect(sl, 6.78, ty+0.1, 0.62, 0.62, fill=C_BLUE)
    add_text(sl, num, 6.78, ty+0.12, 0.62, 0.58,
             size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 7.54, ty+0.1, 5.3, 0.45,
             size=15, bold=True, color=C_NAVY)
    add_text(sl, desc,  7.54, ty+0.54, 5.3, 0.82,
             size=13, color=C_DGRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — EKSTRAKSI PAPER
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "03 · Langkah Pertama: Ekstraksi Paper",
           "Jangan langsung minta review — struktur dulu, analisis kemudian")

# Wrong way
add_rect(sl, 0.3, 1.32, 5.9, 2.4, fill=RGBColor(0xFF,0xEB,0xEB),
         line=C_RED, line_w=Pt(2))
add_text(sl, "❌  Yang Sering Dilakukan", 0.45, 1.38, 5.6, 0.5,
         size=16, bold=True, color=C_RED)
add_text(sl,
    "Upload PDF\n→  \"Tolong review paper ini\"\n\nHasil: generik, tidak sesuai ekspektasi,\nperlu banyak revisi",
    0.48, 1.88, 5.65, 1.72, size=15, color=C_DGRAY)

# Right way
add_rect(sl, 0.3, 3.92, 5.9, 3.15, fill=RGBColor(0xE8,0xF5,0xE9),
         line=C_GREEN, line_w=Pt(2))
add_text(sl, "✅  Yang Seharusnya Dilakukan", 0.45, 3.98, 5.6, 0.5,
         size=16, bold=True, color=C_GREEN)
add_text(sl,
    "Upload PDF\n→  Ekstraksi terstruktur\n→  Koreksi hasil ekstraksi\n→  Beri tugas spesifik per area\n→  Edit & submit",
    0.48, 4.5, 5.65, 2.42, size=15, color=C_DGRAY)

# Right: why extraction matters
add_rect(sl, 6.55, 1.32, 6.5, 5.75, fill=C_WHITE, line=C_BLUE, line_w=Pt(1))
add_rect(sl, 6.55, 1.32, 6.5, 0.52, fill=C_BLUE)
add_text(sl, "Kenapa Ekstraksi Penting?", 6.65, 1.36, 6.3, 0.44,
         size=15, bold=True, color=C_WHITE)

reasons = [
    "Paper PDF panjang & berantakan →\nChatGPT bisa kehilangan fokus",
    "Ekstraksi memaksa struktur →\nKamu tahu persis info yang ditangkap AI",
    "Hasil bisa dikoreksi dulu →\nSebelum lanjut ke analisis",
    "Mudah temukan gap →\nKlaim vs. isi paper yang sebenarnya",
    "4 template tersedia →\nML, Sistem, Survei, IS/Aplikasi",
]
for i, r in enumerate(reasons):
    ty = 1.96 + i * 0.98
    add_rect(sl, 6.7, ty, 0.38, 0.38, fill=C_ORANGE)
    add_text(sl, str(i+1), 6.7, ty+0.01, 0.38, 0.36,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, r, 7.2, ty, 5.65, 0.88, size=13, color=C_DGRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 4 TEMPLATE EKSTRAKSI
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "03 · Empat Template Ekstraksi",
           "Pilih sesuai tipe paper — tersedia di file template-ekstraksi-paper.md")

templates = [
    ("A", "ML / Deep Learning",    C_NAVY,
     ["Arsitektur & klaim novelty", "Dataset, split, metrik", "Baseline & hasil (angka)", "Ablation study ada?"]),
    ("B", "Sistem / Implementasi", C_BLUE,
     ["Komponen & stack teknologi", "Metrik performa sistem", "User study ada?", "Open source / reproducible?"]),
    ("C", "Survei / Lit. Review",  RGBColor(0x6A,0x0D,0xAD),
     ["Jumlah & kriteria paper", "Taksonomi yang diusulkan", "Tabel perbandingan ada?", "Research gap teridentifikasi?"]),
    ("D", "IS / Aplikasi Bisnis",  RGBColor(0x2E,0x86,0x48),
     ["Domain & konteks organisasi", "Framework teoritis (TAM, dll.)", "User study & jumlah partisipan", "Implementasi nyata atau prototype?"]),
]
for i, (lbl, title, color, items) in enumerate(templates):
    col = i % 2; row = i // 2
    lx = 0.28 + col * 6.52
    ty = 1.32 + row * 2.92
    add_rect(sl, lx, ty, 6.22, 2.72, fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, lx, ty, 0.7, 2.72, fill=color)
    add_text(sl, lbl, lx+0.04, ty+1.0, 0.62, 0.72,
             size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Template "+lbl+" · "+title,
             lx+0.82, ty+0.12, 5.25, 0.5,
             size=15, bold=True, color=color)
    for j, item in enumerate(items):
        add_text(sl, "▸  "+item, lx+0.88, ty+0.65+j*0.5, 5.2, 0.44,
                 size=13, color=C_DGRAY)

callout(sl,
    "💡  Tambahkan poin kustomisasi sesuai keahlianmu "
    "(Computer Vision, NLP, Network Security, IoT/Embedded) di bagian Kustomisasi Template.",
    0.28, 7.06, 12.77, 0.22, size=12, bg=C_YELLOW, border=C_ORANGE)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — AI = TASK EXECUTOR
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "04 · AI adalah Pelaksana Tugas, Bukan Penjawab Pertanyaan",
           "Pergeseran mindset yang paling penting")

# Big contrast
add_rect(sl, 0.3, 1.32, 6.0, 5.72, fill=RGBColor(0xFF,0xEB,0xEB),
         line=C_RED, line_w=Pt(2))
add_text(sl, "❌  Pertanyaan (Lemah)", 0.5, 1.4, 5.6, 0.5,
         size=17, bold=True, color=C_RED)

weak = [
    '"Apa kelemahan paper ini?"',
    '"Bagaimana metodologinya?"',
    '"Apakah hasilnya valid?"',
    '"Review bagian referensi"',
]
for i, q in enumerate(weak):
    add_rect(sl, 0.45, 2.02+i*1.18, 5.7, 0.95, fill=C_WHITE)
    add_text(sl, q, 0.62, 2.1+i*1.18, 5.4, 0.78,
             size=14, italic=True, color=C_RED)

add_rect(sl, 6.62, 1.32, 6.42, 5.72, fill=RGBColor(0xE8,0xF5,0xE9),
         line=C_GREEN, line_w=Pt(2))
add_text(sl, "✅  Tugas (Kuat)", 6.8, 1.4, 6.1, 0.5,
         size=17, bold=True, color=C_GREEN)

strong = [
    "Cek apakah baseline (KNN, SVM, RF, LSTM)\nsudah mencakup metode 2022–2024.\nJika ada gap, sebutkan paper spesifik.",
    "Cek Section 3: ada penjelasan\nhyperparameter tuning?\nTulis komentar reviewer memintanya.",
    "NSL-KDD dirilis 2009. Buat komentar\nyang mengevaluasi relevansinya untuk\npaper IoT tahun 2025.",
    "Dari referensi: identifikasi paper >2020\nvs. <2018. Nilai komposisi ini wajar\nuntuk deep learning 2025.",
]
for i, q in enumerate(strong):
    add_rect(sl, 6.78, 2.02+i*1.18, 6.1, 0.95, fill=C_WHITE)
    add_text(sl, q, 6.92, 2.08+i*1.18, 5.85, 0.82,
             size=12.5, color=C_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — STRUKTUR TUGAS EFEKTIF
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "04 · Struktur Tugas yang Efektif",
           "Setiap tugas review = KONTEKS + INSTRUKSI SPESIFIK + FORMAT OUTPUT")

# Three pillars
pillars = [
    ("KONTEKS",              C_NAVY,
     "Berikan info dari hasil ekstraksi\nyang relevan untuk tugas ini",
     '"Baseline yang digunakan:\nKNN, SVM, RF, LSTM.\nPaper ditulis tahun 2025\nuntuk IoT IDS."'),
    ("INSTRUKSI\nSPESIFIK",  C_BLUE,
     "Satu tugas jelas, bukan\npertanyaan terbuka",
     '"Evaluasi apakah baseline\nup-to-date. Cari di referensi\napakah ada transformer/GNN\nuntuk IDS."'),
    ("FORMAT\nOUTPUT",       C_ORANGE,
     "Tentukan bentuk hasil yang\nkamu inginkan",
     '"Tulis:\n• Temuan tentang baseline\n• Paper yang terlewat\n• Draft komentar reviewer\n  (3-4 kalimat, Inggris)"'),
]
for i, (title, color, desc, ex) in enumerate(pillars):
    lx = 0.28 + i * 4.32
    add_rect(sl, lx, 1.32, 4.1, 1.08, fill=color)
    add_text(sl, title, lx+0.15, 1.38, 3.85, 0.96,
             size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, 2.4, 4.1, 1.38, fill=C_WHITE, line=color, line_w=Pt(1))
    add_text(sl, desc, lx+0.15, 2.5, 3.85, 1.18,
             size=14, color=C_DGRAY, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, 3.88, 4.1, 3.22, fill=RGBColor(0xF0,0xF6,0xFF),
             line=color, line_w=Pt(1))
    add_text(sl, "Contoh:", lx+0.15, 3.95, 3.85, 0.38,
             size=13, bold=True, color=color)
    add_text(sl, ex, lx+0.15, 4.3, 3.85, 2.55,
             size=12.5, italic=True, color=C_DGRAY)

# Arrow connectors (simple text arrows)
add_text(sl, "→", 4.38, 2.1, 0.45, 0.6,
         size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_text(sl, "→", 8.7,  2.1, 0.45, 0.6,
         size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — ALUR LENGKAP 5 LANGKAH
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "05 · Alur Lengkap: 5 Langkah",
           "Dari upload paper ke draft review siap submit")

flow = [
    ("1", "UPLOAD",    C_NAVY,   "Upload PDF ke\nProject ChatGPT Plus"),
    ("2", "EKSTRAKSI", C_BLUE,   "Kirim prompt dari\ntemplate-ekstraksi-paper.md"),
    ("3", "KOREKSI",   RGBColor(0x6A,0x0D,0xAD), "Baca & koreksi hasil\nextraksi — jangan skip!"),
    ("4", "TUGAS",     C_ORANGE, "Beri tugas spesifik\nper area review"),
    ("5", "EDIT",      C_GREEN,  "Edit output AI\n+ keahlian kamu → submit"),
]
box_w = 2.25
for i, (num, label, color, desc) in enumerate(flow):
    lx = 0.28 + i * 2.56
    # Connector arrow (except first)
    if i > 0:
        add_text(sl, "→", lx-0.42, 2.65, 0.38, 0.6,
                 size=22, bold=True, color=C_DGRAY, align=PP_ALIGN.CENTER)
    # Number circle (rect)
    add_rect(sl, lx+0.62, 1.32, 1.0, 0.9, fill=color)
    add_text(sl, num, lx+0.62, 1.34, 1.0, 0.86,
             size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Label
    add_rect(sl, lx, 2.22, box_w, 0.72, fill=color)
    add_text(sl, label, lx+0.05, 2.28, box_w-0.08, 0.6,
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Description
    add_rect(sl, lx, 2.94, box_w, 1.42, fill=C_WHITE, line=color, line_w=Pt(1.5))
    add_text(sl, desc, lx+0.1, 3.02, box_w-0.18, 1.26,
             size=13.5, color=C_DGRAY, align=PP_ALIGN.CENTER)

# Time estimate strip
add_rect(sl, 0.28, 4.55, 12.78, 2.55, fill=C_WHITE, line=C_LBLUE, line_w=Pt(1))
add_rect(sl, 0.28, 4.55, 12.78, 0.48, fill=C_NAVY)
add_text(sl, "Estimasi Waktu per Paper", 0.42, 4.58, 12.5, 0.42,
         size=15, bold=True, color=C_WHITE)

time_items = [
    ("Upload + Ekstraksi",    "10 menit"),
    ("Koreksi hasil",         "5 menit"),
    ("Tugas per area (4×)",   "15 menit"),
    ("Edit & finalisasi",     "10–15 menit"),
    ("TOTAL",                 "~40–45 menit"),
]
for i, (label, dur) in enumerate(time_items):
    lx = 0.55 + i * 2.52
    bg = C_NAVY if label == "TOTAL" else C_LGRAY
    fc = C_WHITE if label == "TOTAL" else C_DGRAY
    add_rect(sl, lx, 5.12, 2.28, 1.72, fill=bg)
    add_text(sl, label, lx+0.08, 5.2,  2.12, 0.72, size=12, color=fc, align=PP_ALIGN.CENTER)
    add_text(sl, dur,   lx+0.08, 5.9,  2.12, 0.72, size=16, bold=True, color=C_ORANGE if label=="TOTAL" else fc, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — BIAS, VARIASI, KONSISTENSI
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(sl, "06 · Jaga Bias, Variasi & Konsistensi",
           "Tiga hal yang menentukan kualitas review di seluruh paper")

sections = [
    ("⚠️  Bias", C_RED, [
        "ChatGPT cenderung mengonfirmasi framing kamu",
        "Pisahkan: minta kelemahan dulu, baru kelebihan",
        "Teks Inggris bagus → evaluasi lebih positif (salah!)",
        'Solusi: "Abaikan kualitas bahasa — fokus eksperimen"',
    ]),
    ("🔀  Variasi", C_BLUE, [
        "Jika ragu, minta 3 versi komentar:",
        '  • Versi Major Revision',
        '  • Versi Minor Revision',
        '  • Versi hanya klarifikasi',
        "Pilih yang sesuai penilaian expert kamu",
    ]),
    ("✅  Konsistensi", C_GREEN, [
        "Satu chat per paper — jangan ganti-ganti",
        "Gunakan Project agar style sama di semua paper",
        'Cross-check: "Apakah penilaian novelty konsisten\n  dengan metodologi?"',
        'Final check: "Adakah inkonsistensi di semua komentar?"',
    ]),
]
for i, (title, color, items) in enumerate(sections):
    lx = 0.28 + i * 4.32
    add_rect(sl, lx, 1.32, 4.1, 5.82, fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, lx, 1.32, 4.1, 0.6, fill=color)
    add_text(sl, title, lx+0.15, 1.37, 3.85, 0.5,
             size=18, bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        add_text(sl, item, lx+0.2, 2.1+j*0.82, 3.82, 0.74,
                 size=13.5, color=C_DGRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 3.35, 13.33, 4.15, fill=C_BLUE)
add_rect(sl, 0, 3.28, 13.33, 0.1, fill=C_ORANGE)

add_text(sl, "Kamu yang Ahli.", 0.6, 0.55, 12.1, 1.0,
         size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "AI yang Mengerjakan.", 0.6, 1.42, 12.1, 0.9,
         size=36, bold=False, color=C_LBLUE, align=PP_ALIGN.CENTER)
add_rect(sl, 3.0, 2.5, 7.33, 0.07, fill=C_ORANGE)
add_text(sl,
    "Encode keahlianmu sekali · Ekstrak dulu · Beri tugas yang jelas",
    0.6, 2.68, 12.1, 0.55,
    size=18, color=C_LBLUE, align=PP_ALIGN.CENTER, italic=True)

stats = [
    ("~40 menit", "per paper"),
    ("4 template", "ekstraksi siap pakai"),
    ("1× setup", "untuk semua paper ICICOS"),
]
for i, (big, small) in enumerate(stats):
    lx = 1.2 + i * 3.85
    add_rect(sl, lx, 3.65, 3.3, 1.82, fill=C_NAVY)
    add_text(sl, big,   lx+0.12, 3.72, 3.1, 0.88,
             size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, small, lx+0.12, 4.55, 3.1, 0.72,
             size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
    "Panduan lengkap: tutorials/panduan-review-paper-icicos-chatgpt.md\n"
    "Template ekstraksi: tutorials/template-ekstraksi-paper.md",
    0.6, 5.65, 12.1, 0.88,
    size=14, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(sl, "ICICOS 2026  ·  Juni 2026",
         0.6, 6.72, 12.1, 0.48,
         size=14, color=C_LBLUE, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────
OUT = "tutorials/ICICOS_Review_Tutorial.pptx"
prs.save(OUT)
print(f"✅  Saved: {OUT}  ({prs.slides.__len__()} slides)")

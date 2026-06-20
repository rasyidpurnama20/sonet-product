#!/usr/bin/env python3
"""
Generate ICICoS 2026 Reviewer Guide Presentation (7 slides).
Futuristic design with custom shapes, colors, and diagrams.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Color palette - futuristic dark theme
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)       # Deep navy
BG_MEDIUM = RGBColor(0x1B, 0x2B, 0x44)     # Medium navy
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)   # Bright cyan
ACCENT_PURPLE = RGBColor(0xA855, 0xF7, 0x00)[0:3] if False else RGBColor(0xA8, 0x55, 0xF7)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)  # Emerald
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C) # Warm orange
ACCENT_RED = RGBColor(0xF4, 0x3F, 0x5E)    # Coral red
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)    # Light gray
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)    # Muted gray
HIGHLIGHT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_text(slide, left, top, width, height, text, font_size=12,
                        font_color=TEXT_WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=None, align=PP_ALIGN.CENTER, line_color=None):
    """Add a shape with centered text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 font_color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=14,
                          default_color=TEXT_WHITE, line_spacing=1.2):
    """Add a text box with multiple formatted lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_info.get("text", "")
        p.text = text
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.alignment = line_info.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space_after", 4))
    return txBox


def create_slide_1_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    # Decorative accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    # Main title
    add_multiline_textbox(slide, Inches(1), Inches(1.8), Inches(8), Inches(2.5), [
        {"text": "PANDUAN REVIEWER ICICoS 2026", "size": 36, "bold": True,
         "color": TEXT_WHITE, "align": PP_ALIGN.CENTER, "space_after": 12},
        {"text": "Review Berbantuan AI", "size": 24, "color": ACCENT_CYAN,
         "align": PP_ALIGN.CENTER, "space_after": 20},
        {"text": "Prinsip, Nilai, dan Kontrol Penting", "size": 16,
         "color": TEXT_MUTED, "align": PP_ALIGN.CENTER, "space_after": 8},
    ])

    # Decorative shapes
    add_shape_with_text(slide, Inches(3.5), Inches(5.0), Inches(3), Inches(0.6),
                        "ICICoS 2026 Conference", font_size=11,
                        font_color=ACCENT_CYAN, line_color=ACCENT_CYAN)

    # Bottom accent
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), Inches(10), Inches(0.08))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_PURPLE
    bar2.line.fill.background()


def create_slide_2_comparison(prs):
    """Slide 2: Reviewer vs Reviewer+AI comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title
    add_multiline_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8), [
        {"text": "Reviewer Manual vs Reviewer + AI", "size": 24, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    # Left column - Traditional Reviewer
    add_shape_with_text(slide, Inches(0.3), Inches(1.2), Inches(4.4), Inches(0.55),
                        "REVIEWER MANUAL", font_size=13, bold=True,
                        font_color=BG_DARK, fill_color=ACCENT_ORANGE)

    manual_items = [
        "Membaca & menilai sendiri",
        "Format komentar bervariasi",
        "Risiko miss section tertentu",
        "Konsistensi tergantung mood/waktu",
        "Nada bisa subjektif",
        "Keputusan: 100% manusia",
    ]
    for i, item in enumerate(manual_items):
        add_shape_with_text(slide, Inches(0.4), Inches(1.9 + i * 0.75), Inches(4.2), Inches(0.6),
                            f"  {item}", font_size=11, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=ACCENT_ORANGE)

    # Right column - Reviewer + AI
    add_shape_with_text(slide, Inches(5.3), Inches(1.2), Inches(4.4), Inches(0.55),
                        "REVIEWER + AI", font_size=13, bold=True,
                        font_color=BG_DARK, fill_color=ACCENT_GREEN)

    ai_items = [
        "AI bantu checklist & struktur",
        "Format konsisten (14 template)",
        "Semua section tertinjau sistematis",
        "Konsistensi terjaga (AI checklist)",
        "Nada profesional (AI perapih)",
        "Keputusan: TETAP 100% manusia",
    ]
    for i, item in enumerate(ai_items):
        add_shape_with_text(slide, Inches(5.4), Inches(1.9 + i * 0.75), Inches(4.2), Inches(0.6),
                            f"  {item}", font_size=11, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=ACCENT_GREEN)

    # Bottom note
    add_multiline_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.5), [
        {"text": "AI = Asisten, BUKAN pengganti. Keputusan Accept/Reject tetap tanggung jawab reviewer.",
         "size": 10, "color": HIGHLIGHT_YELLOW, "align": PP_ALIGN.CENTER},
    ])


def create_slide_3_ethics(prs):
    """Slide 3: 7 Aturan Etika & Penggunaan AI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_multiline_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7), [
        {"text": "7 Aturan Etika Penggunaan AI dalam Review", "size": 22, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    rules = [
        ("DILARANG", "Upload manuskrip ke AI publik (rahasia!)", ACCENT_RED),
        ("DILARANG", "Delegasi keputusan akhir ke AI", ACCENT_RED),
        ("DILARANG", "Klaim/temuan fiktif dari AI (halusinasi)", ACCENT_RED),
        ("BOLEH", "AI sebagai asisten bahasa (tata bahasa, nada)", ACCENT_GREEN),
        ("BOLEH", "AI sebagai checklist/penstruktur review", ACCENT_GREEN),
        ("BOLEH", "AI untuk brainstorming pertanyaan (verifikasi manual!)", ACCENT_GREEN),
        ("WAJIB", "Akuntabilitas: komentar harus bisa ditelusuri ke paper", HIGHLIGHT_YELLOW),
    ]

    for i, (tag, desc, color) in enumerate(rules):
        y_pos = Inches(1.1 + i * 0.85)
        # Tag badge
        add_shape_with_text(slide, Inches(0.4), y_pos, Inches(1.5), Inches(0.55),
                            tag, font_size=10, bold=True,
                            font_color=BG_DARK, fill_color=color)
        # Description
        add_shape_with_text(slide, Inches(2.0), y_pos, Inches(7.5), Inches(0.55),
                            desc, font_size=12, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=color)


def create_slide_4_ai_weaknesses(prs):
    """Slide 4: Kontrol terhadap kelemahan AI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_multiline_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7), [
        {"text": "Kontrol Kelemahan AI dalam Peer Review", "size": 22, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    # Three columns for three weaknesses
    weaknesses = [
        {
            "title": "HALUSINASI",
            "icon_color": ACCENT_RED,
            "controls": [
                "Aturan Bukti-atau-Hapus",
                "Wajib tunjuk hal/baris",
                "Komentar tanpa rujukan",
                "= tidak sah",
            ]
        },
        {
            "title": "KERAHASIAAN",
            "icon_color": ACCENT_ORANGE,
            "controls": [
                "Jangan upload ke AI publik",
                "Matikan model training",
                "Gunakan Temporary Chat",
                "Kutip max 1-2 kalimat",
            ]
        },
        {
            "title": "DELEGASI\nKEPUTUSAN",
            "icon_color": ACCENT_PURPLE,
            "controls": [
                "Accept/Reject = manusia",
                "AI hanya brainstorm",
                "Reviewer verifikasi semua",
                "Akuntabilitas penuh",
            ]
        },
    ]

    col_width = Inches(3.0)
    for col, w in enumerate(weaknesses):
        x_start = Inches(0.4 + col * 3.2)
        # Header
        add_shape_with_text(slide, x_start, Inches(1.1), col_width, Inches(0.7),
                            w["title"], font_size=13, bold=True,
                            font_color=TEXT_WHITE, fill_color=w["icon_color"])
        # Controls
        for i, ctrl in enumerate(w["controls"]):
            add_shape_with_text(slide, x_start, Inches(2.0 + i * 0.7), col_width, Inches(0.55),
                                ctrl, font_size=11, font_color=TEXT_LIGHT,
                                align=PP_ALIGN.CENTER, line_color=w["icon_color"])

    # Bottom principle
    add_shape_with_text(slide, Inches(1.0), Inches(5.2), Inches(8.0), Inches(0.7),
                        "Prinsip Utama: AI memperkuat review, BUKAN menggantikan reviewer",
                        font_size=13, bold=True, font_color=HIGHLIGHT_YELLOW,
                        line_color=HIGHLIGHT_YELLOW)


def create_slide_5_format_comments(prs):
    """Slide 5: 14 Format Komentar (overview)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_multiline_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7), [
        {"text": "14 Format Komentar Terstruktur", "size": 22, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    formats_left = [
        "#1  Kutipan + Tanggapan + Saran",
        "#2  Klasifikasi (Major/Minor/Optional)",
        "#3  Pertanyaan (Clarification)",
        "#4  Strength - Weakness",
        "#5  Claim - Evidence - Gap",
        "#6  Actionable (ber-referensi baris)",
        "#7  Severity + Fix + Rationale",
    ]
    formats_right = [
        "#8  Reproducibility Checklist",
        "#9  Positioning Literatur",
        "#10 Probe Asumsi / Hipotesis",
        "#11 Rigor Metodologis / Statistik",
        "#12 Kejelasan & Presentasi",
        "#13 Consider (saran lunak)",
        "#14 Komentar Apresiatif",
    ]

    for i, fmt in enumerate(formats_left):
        color = ACCENT_GREEN if i < 4 else ACCENT_CYAN
        add_shape_with_text(slide, Inches(0.3), Inches(1.1 + i * 0.8), Inches(4.6), Inches(0.6),
                            fmt, font_size=11, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=color)

    for i, fmt in enumerate(formats_right):
        color = ACCENT_PURPLE if i < 4 else ACCENT_ORANGE
        add_shape_with_text(slide, Inches(5.1), Inches(1.1 + i * 0.8), Inches(4.6), Inches(0.6),
                            fmt, font_size=11, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=color)

    # Tag system note
    add_multiline_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.5), [
        {"text": "Tag: [ID-komentar] [Section] [Format#] [Severity] [Status: Open/Resolved]",
         "size": 10, "color": TEXT_MUTED, "align": PP_ALIGN.CENTER},
    ])


def create_slide_6_chatgpt_safety(prs):
    """Slide 6: Cara aman ChatGPT."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_multiline_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7), [
        {"text": "Pengaturan Keamanan ChatGPT untuk Review", "size": 22, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    # Steps box
    steps = [
        "1. Buka ChatGPT",
        "2. Klik profil/nama akun",
        "3. Masuk ke Settings",
        "4. Pilih Data Controls",
        "5. Matikan \"Improve the model for everyone\"",
    ]

    add_shape_with_text(slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(0.5),
                        "LANGKAH PENGATURAN DATA CONTROLS", font_size=12, bold=True,
                        font_color=BG_DARK, fill_color=ACCENT_GREEN)

    for i, step in enumerate(steps):
        add_shape_with_text(slide, Inches(0.5), Inches(1.6 + i * 0.6), Inches(5.5), Inches(0.5),
                            step, font_size=12, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=ACCENT_GREEN)

    # File upload info
    add_shape_with_text(slide, Inches(6.2), Inches(1.0), Inches(3.5), Inches(0.5),
                        "FILE UPLOAD", font_size=12, bold=True,
                        font_color=BG_DARK, fill_color=ACCENT_ORANGE)

    add_shape_with_text(slide, Inches(6.2), Inches(1.6), Inches(3.5), Inches(1.8),
                        "File & chat mengikuti\nsettings Data Controls.\n\nMatikan training\nSEBELUM upload\ndokumen penting.",
                        font_size=11, font_color=TEXT_LIGHT,
                        align=PP_ALIGN.CENTER, line_color=ACCENT_ORANGE)

    # Temporary Chat section
    add_shape_with_text(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.5),
                        "TEMPORARY CHAT (untuk percakapan sangat sensitif)",
                        font_size=12, bold=True, font_color=BG_DARK, fill_color=ACCENT_PURPLE)

    temp_features = [
        "Tidak muncul di history",
        "Tidak membuat memory",
        "Tidak digunakan untuk melatih model",
        "Catatan: salinan disimpan maks 30 hari (abuse monitoring)",
    ]
    for i, feat in enumerate(temp_features):
        icon = "+" if i < 3 else "!"
        color = ACCENT_GREEN if i < 3 else HIGHLIGHT_YELLOW
        add_shape_with_text(slide, Inches(0.5), Inches(5.2 + i * 0.55), Inches(9.0), Inches(0.45),
                            f"  [{icon}] {feat}", font_size=11, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=color)


def create_slide_7_summary(prs):
    """Slide 7: Summary / Key Takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    add_multiline_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7), [
        {"text": "Ringkasan: Prinsip & Nilai Utama", "size": 22, "bold": True,
         "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER},
    ])

    takeaways = [
        ("Kerahasiaan", "Jangan upload manuskrip ke AI publik, atur Data Controls",
         ACCENT_RED),
        ("Akuntabilitas", "Setiap komentar harus merujuk halaman/baris paper nyata",
         ACCENT_ORANGE),
        ("Keputusan Manusia", "Accept/Reject tidak boleh didelegasikan ke AI",
         ACCENT_PURPLE),
        ("Anti-Halusinasi", "Aturan Bukti-atau-Hapus: tanpa bukti = tidak sah",
         HIGHLIGHT_YELLOW),
        ("Nada Konstruktif", "Kritik ke naskah, bukan penulis. Spesifik & netral",
         ACCENT_GREEN),
        ("Format Terstruktur", "14 format komentar + tag tracking antar-ronde revisi",
         ACCENT_CYAN),
    ]

    for i, (title, desc, color) in enumerate(takeaways):
        y = Inches(1.1 + i * 0.95)
        # Title badge
        add_shape_with_text(slide, Inches(0.4), y, Inches(2.3), Inches(0.6),
                            title, font_size=11, bold=True,
                            font_color=BG_DARK, fill_color=color)
        # Description
        add_shape_with_text(slide, Inches(2.8), y, Inches(6.8), Inches(0.6),
                            desc, font_size=12, font_color=TEXT_LIGHT,
                            align=PP_ALIGN.LEFT, line_color=color)

    # Footer
    add_multiline_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.5), [
        {"text": "\"AI memperkuat kemampuan reviewer, tetapi integritas tetap di tangan manusia.\"",
         "size": 11, "color": TEXT_MUTED, "align": PP_ALIGN.CENTER},
    ])

    # Bottom bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), Inches(10), Inches(0.06))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_PURPLE
    bar2.line.fill.background()


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_slide_1_title(prs)
    create_slide_2_comparison(prs)
    create_slide_3_ethics(prs)
    create_slide_4_ai_weaknesses(prs)
    create_slide_5_format_comments(prs)
    create_slide_6_chatgpt_safety(prs)
    create_slide_7_summary(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "icicos-reviewer-guide.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

"""
Generate ICICoS Reviewer Guide PPT — max 7 slides, 16:9 standard size
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x1A, 0x2B, 0x4A)   # background utama
ACCENT_BLUE = RGBColor(0x1E, 0x6B, 0xBF)   # aksen/header
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xFF)   # background tabel header
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xFA)
MID_GRAY    = RGBColor(0x55, 0x65, 0x7A)
ACCENT_ORG  = RGBColor(0xF0, 0x8C, 0x00)   # highlight / badge

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helper functions ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(12), bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, font_size=Pt(11), bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False, space_before=Pt(0)):
    from pptx.oxml.ns import qn
    from lxml import etree
    para = tf.add_paragraph()
    para.alignment = align
    para.space_before = space_before
    run = para.add_run()
    run.text = text
    run.font.size   = font_size
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return para

def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color



# ── SLIDE 1 — Cover ──────────────────────────────────────────────────────────
def slide1_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, DARK_NAVY)

    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=ACCENT_BLUE)

    # decorative circle top-right
    c = slide.shapes.add_shape(9, Inches(10.8), Inches(-0.8), Inches(3.2), Inches(3.2))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT_BLUE
    c.line.fill.background()

    # badge
    badge = add_rect(slide, Inches(0.55), Inches(1.6), Inches(2.8), Inches(0.42), fill=ACCENT_ORG)
    add_textbox(slide, "ICICoS 2026", Inches(0.55), Inches(1.6), Inches(2.8), Inches(0.42),
                font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # main title
    add_textbox(slide, "Panduan Reviewer", Inches(0.55), Inches(2.2), Inches(9), Inches(0.85),
                font_size=Pt(42), bold=True, color=WHITE)
    add_textbox(slide, "Review Berbantuan AI",
                Inches(0.55), Inches(3.0), Inches(9), Inches(0.75),
                font_size=Pt(36), bold=True, color=RGBColor(0xA8, 0xC8, 0xFF))

    # subtitle line
    add_rect(slide, Inches(0.55), Inches(3.85), Inches(4.5), Inches(0.06), fill=ACCENT_ORG)

    # sub info
    add_textbox(slide, "Panduan Umum · Berlaku semua topik ICICoS",
                Inches(0.55), Inches(4.05), Inches(10), Inches(0.45),
                font_size=Pt(16), color=RGBColor(0xB0, 0xC8, 0xE8), italic=True)
    add_textbox(slide, "Versi 1.1  |  Bahasa Indonesia / English",
                Inches(0.55), Inches(4.55), Inches(10), Inches(0.4),
                font_size=Pt(13), color=MID_GRAY)

    # bottom tagline
    add_textbox(slide, "Bagian A  ·  B  ·  C  ·  D  ·  E",
                Inches(0.55), Inches(6.6), Inches(10), Inches(0.5),
                font_size=Pt(13), color=MID_GRAY)



# ── SLIDE 2 — Alur Kerja Review ───────────────────────────────────────────────
def slide2_alur(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT_GRAY)

    # header strip
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=DARK_NAVY)
    add_textbox(slide, "Alur Kerja Review", Inches(0.45), Inches(0.15),
                Inches(9), Inches(0.75), font_size=Pt(28), bold=True, color=WHITE)
    add_textbox(slide, "Ikuti urutan bagian secara berurutan sebelum mengisi rubrik penilaian",
                Inches(0.45), Inches(0.6), Inches(10), Inches(0.45),
                font_size=Pt(13), color=RGBColor(0xA8, 0xC8, 0xFF), italic=True)

    steps = [
        ("A", "Deklarasi Awal",      "Identitas paper,\nkeahlian & tingkat\nkeyakinan reviewer"),
        ("B", "Rubrik Penilaian",    "Overall Evaluation\n(7 level), Confidence\n(0–4), 5 Quality Criteria"),
        ("C", "Checklist Seksi",     "C.1 Abstract s.d.\nC.11 Tata Bahasa\n& Konsistensi Istilah"),
        ("D", "Komentar Terstruktur","5 Format komentar\n+ Definisi Severity\n(MAJOR/MINOR/SARAN)"),
        ("E", "Etika AI",            "4 aturan wajib\npenggunaan AI\ndalam proses review"),
    ]

    box_w = Inches(2.3)
    box_h = Inches(4.4)
    gap   = Inches(0.22)
    start_x = Inches(0.35)
    top_y = Inches(1.4)

    for i, (letter, title, desc) in enumerate(steps):
        x = start_x + i * (box_w + gap)

        # card bg
        add_rect(slide, x, top_y, box_w, box_h, fill=WHITE,
                 line_color=ACCENT_BLUE, line_width=Pt(1.2))

        # circle badge
        circ = slide.shapes.add_shape(9, x + Inches(0.75), top_y + Inches(0.18),
                                       Inches(0.8), Inches(0.8))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT_BLUE
        circ.line.fill.background()
        add_textbox(slide, letter, x + Inches(0.75), top_y + Inches(0.18),
                    Inches(0.8), Inches(0.8),
                    font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # title
        add_textbox(slide, title, x + Inches(0.1), top_y + Inches(1.1),
                    box_w - Inches(0.2), Inches(0.7),
                    font_size=Pt(13), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

        # separator
        add_rect(slide, x + Inches(0.3), top_y + Inches(1.85),
                 box_w - Inches(0.6), Inches(0.04), fill=ACCENT_BLUE)

        # desc
        add_textbox(slide, desc, x + Inches(0.1), top_y + Inches(2.0),
                    box_w - Inches(0.2), Inches(2.0),
                    font_size=Pt(11), color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

        # arrow between cards
        if i < len(steps) - 1:
            ax = x + box_w + gap * 0.15
            add_textbox(slide, "›", ax, top_y + Inches(1.9),
                        gap * 0.7, Inches(0.5),
                        font_size=Pt(22), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)



# ── SLIDE 3 — Bagian A & B ────────────────────────────────────────────────────
def slide3_ab(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT_GRAY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=DARK_NAVY)
    add_textbox(slide, "Bagian A — Deklarasi Awal  ·  Bagian B — Rubrik Penilaian",
                Inches(0.45), Inches(0.2), Inches(12), Inches(0.7),
                font_size=Pt(24), bold=True, color=WHITE)

    # ── LEFT: Bagian A ──
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(5.9), Inches(5.9), fill=WHITE,
             line_color=ACCENT_BLUE, line_width=Pt(1))
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(5.9), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(slide, "A  Deklarasi Awal", Inches(0.45), Inches(1.28),
                Inches(5.7), Inches(0.44), font_size=Pt(14), bold=True, color=WHITE)

    a_items = [
        ("paper_id", "Diisi sesuai sistem submission"),
        ("paper_title", "Sesuai judul naskah"),
        ("tanggal_review", "Format YYYY-MM-DD"),
        ("bidang_keahlian", "1–2 bidang utama reviewer"),
        ("justifikasi", "Mengapa reviewer kompeten?"),
        ("tingkat_keyakinan", "Expert(4) / High(3) / Medium(2) / Low(1)"),
    ]
    for idx, (k, v) in enumerate(a_items):
        y = Inches(1.9) + idx * Inches(0.72)
        add_rect(slide, Inches(0.45), y, Inches(1.9), Inches(0.6),
                 fill=LIGHT_BLUE, line_color=ACCENT_BLUE, line_width=Pt(0.5))
        add_textbox(slide, k, Inches(0.48), y + Inches(0.05), Inches(1.85), Inches(0.55),
                    font_size=Pt(9), bold=True, color=ACCENT_BLUE)
        add_textbox(slide, v, Inches(2.45), y + Inches(0.1), Inches(3.7), Inches(0.5),
                    font_size=Pt(10), color=MID_GRAY)

    # ── RIGHT: Bagian B ──
    add_rect(slide, Inches(6.55), Inches(1.25), Inches(6.45), Inches(5.9), fill=WHITE,
             line_color=ACCENT_BLUE, line_width=Pt(1))
    add_rect(slide, Inches(6.55), Inches(1.25), Inches(6.45), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(slide, "B  Rubrik Penilaian Resmi ICICoS", Inches(6.65), Inches(1.28),
                Inches(6.2), Inches(0.44), font_size=Pt(14), bold=True, color=WHITE)

    # Overall evaluation
    add_textbox(slide, "Overall Evaluation (pilih satu):", Inches(6.65), Inches(1.88),
                Inches(6.1), Inches(0.38), font_size=Pt(11), bold=True, color=DARK_NAVY)
    oe = ["Strong Accept", "Accept (revision req.)", "Weak Accept (revision req.)",
          "Borderline (revision req.)", "Weak Reject", "Reject", "Strong Reject"]
    for idx, label in enumerate(oe):
        y = Inches(2.28) + idx * Inches(0.35)
        add_rect(slide, Inches(6.65), y, Inches(0.28), Inches(0.26),
                 fill=None, line_color=ACCENT_BLUE, line_width=Pt(1))
        add_textbox(slide, label, Inches(7.0), y, Inches(5.8), Inches(0.32),
                    font_size=Pt(10), color=DARK_NAVY)

    # Confidence
    add_textbox(slide, "Reviewer's Confidence: Expert(4)  High(3)  Medium(2)  Low(1)  N/A(0)",
                Inches(6.65), Inches(4.82), Inches(6.1), Inches(0.45),
                font_size=Pt(10), bold=True, color=ACCENT_BLUE)

    # 5 quality criteria
    add_textbox(slide, "Quality of the Article (Excellent / Good / Adequate / Inadequate):",
                Inches(6.65), Inches(5.32), Inches(6.1), Inches(0.38),
                font_size=Pt(10), bold=True, color=DARK_NAVY)
    criteria = ["Novelty / Originality", "Significance of Topic",
                "Technical Quality", "Presentation", "Literature"]
    cols = ["", "", ""]
    row_y = Inches(5.72)
    for idx, cr in enumerate(criteria):
        col = idx % 3
        xi = Inches(6.65) + col * Inches(2.1)
        yi = row_y + (idx // 3) * Inches(0.38)
        add_textbox(slide, f"• {cr}", xi, yi, Inches(2.05), Inches(0.36),
                    font_size=Pt(10), color=MID_GRAY)



# ── SLIDE 4 — Bagian C Checklist ─────────────────────────────────────────────
def slide4_checklist(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT_GRAY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=DARK_NAVY)
    add_textbox(slide, "Bagian C — Checklist Review per Seksi",
                Inches(0.45), Inches(0.2), Inches(11), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_textbox(slide, "Item gagal → buat komentar di Bagian D",
                Inches(0.45), Inches(0.65), Inches(9), Inches(0.4),
                font_size=Pt(13), color=RGBColor(0xA8, 0xC8, 0xFF), italic=True)

    sections = [
        ("C.1",  "Abstract",                  "Klaim ↔ hasil, angka kunci, kontribusi jelas"),
        ("C.2",  "Introduction & Motivasi",   "Masalah jelas, kontribusi terukur, struktur dijelaskan"),
        ("C.3",  "Related Work",              "Gap eksplisit, sitasi terkini, posisi kontribusi jelas"),
        ("C.4",  "Data & Preprocessing",      "Sumber data, atribut, preprocessing, data leakage"),
        ("C.5",  "Metodologi & Model",        "Justifikasi alg., baseline, hyperparameter, validasi"),
        ("C.6",  "Eksperimen & Hasil",        "Metrik tepat, mean±std, uji statistik, tabel/gambar"),
        ("C.7",  "Diskusi",                   "Interpretasi, limitasi, ancaman validitas"),
        ("C.8",  "Kesimpulan & Future Work",  "Konsisten, tidak over-claim, future work realistis"),
        ("C.9",  "Referensi & Format",        "Sitasi lengkap, format ICICoS, batas halaman"),
        ("C.10", "Reproduksibilitas",         "Kode/data, seed, versi library, lingkungan eksperimen"),
        ("C.11", "Tata Bahasa & Konsistensi", "Gramatikal, istilah konsisten, notasi, tanda baca"),
    ]

    col_count = 2
    rows_per_col = 6
    card_w = Inches(6.3)
    card_h = Inches(0.78)
    gap_x  = Inches(0.35)
    gap_y  = Inches(0.1)
    start_x = Inches(0.35)
    start_y = Inches(1.25)

    for idx, (code, title, desc) in enumerate(sections):
        col = idx // rows_per_col
        row = idx % rows_per_col
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rect(slide, x, y, card_w, card_h, fill=WHITE,
                 line_color=ACCENT_BLUE, line_width=Pt(0.8))
        # code badge
        add_rect(slide, x, y, Inches(0.72), card_h, fill=ACCENT_BLUE)
        add_textbox(slide, code, x, y + Inches(0.22), Inches(0.72), Inches(0.36),
                    font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        add_textbox(slide, title, x + Inches(0.8), y + Inches(0.06),
                    card_w - Inches(0.9), Inches(0.34),
                    font_size=Pt(11), bold=True, color=DARK_NAVY)
        # desc
        add_textbox(slide, desc, x + Inches(0.8), y + Inches(0.38),
                    card_w - Inches(0.9), Inches(0.36),
                    font_size=Pt(9.5), color=MID_GRAY, italic=True)



# ── SLIDE 5 — Bagian D: Definisi Severity ────────────────────────────────────
def slide5_severity(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT_GRAY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=DARK_NAVY)
    add_textbox(slide, "Bagian D — Komentar Terstruktur: Definisi Severity",
                Inches(0.45), Inches(0.2), Inches(11), Inches(0.7),
                font_size=Pt(26), bold=True, color=WHITE)
    add_textbox(slide, "Setiap komentar MAJOR/MINOR wajib menyebut lokasi nyata (halaman/seksi/tabel/gambar)",
                Inches(0.45), Inches(0.68), Inches(12), Inches(0.38),
                font_size=Pt(12), color=RGBColor(0xA8, 0xC8, 0xFF), italic=True)

    severity = [
        (RGBColor(0xC0, 0x39, 0x2B), "[MAJOR]",
         "Klaim inti tidak dapat diverifikasi atau tidak valid",
         "Paper TIDAK DAPAT diterima dalam kondisi ini.",
         "Tidak ada baseline pembanding → klaim 'improving' belum terbukti"),
        (RGBColor(0xE6, 0x7E, 0x22), "[MINOR]",
         "Mengurangi kejelasan/kekuatan paper, tetapi tidak membatalkan klaim inti",
         "Paper masih dapat diterima jika poin ini diperbaiki.",
         "Seed & versi library tidak dilaporkan → replikasi terhambat"),
        (RGBColor(0x27, 0xAE, 0x60), "[SARAN]",
         "Perbaikan opsional; penulis bebas mengabaikan",
         "Tidak mempengaruhi keputusan Accept/Reject.",
         "Pertimbangkan SHAP untuk meningkatkan interpretabilitas"),
    ]

    card_w = Inches(12.5)
    card_h = Inches(1.68)
    start_x = Inches(0.4)
    start_y = Inches(1.3)
    gap = Inches(0.22)

    for idx, (color, tag, defn, impact, example) in enumerate(severity):
        y = start_y + idx * (card_h + gap)
        add_rect(slide, start_x, y, card_w, card_h, fill=WHITE,
                 line_color=color, line_width=Pt(1.5))
        # left color bar
        add_rect(slide, start_x, y, Inches(0.22), card_h, fill=color)
        # tag badge
        add_rect(slide, start_x + Inches(0.3), y + Inches(0.2),
                 Inches(1.3), Inches(0.44), fill=color)
        add_textbox(slide, tag,
                    start_x + Inches(0.3), y + Inches(0.2),
                    Inches(1.3), Inches(0.44),
                    font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # definition
        add_textbox(slide, defn,
                    start_x + Inches(1.75), y + Inches(0.12),
                    Inches(10.4), Inches(0.4),
                    font_size=Pt(12), bold=True, color=DARK_NAVY)
        # impact
        add_textbox(slide, f"↳ {impact}",
                    start_x + Inches(1.75), y + Inches(0.52),
                    Inches(10.4), Inches(0.38),
                    font_size=Pt(11), color=color, bold=True)
        # example
        add_textbox(slide, f"Contoh: {example}",
                    start_x + Inches(1.75), y + Inches(0.92),
                    Inches(10.4), Inches(0.42),
                    font_size=Pt(10), color=MID_GRAY, italic=True)



# ── SLIDE 6 — Bagian D: 5 Format Komentar ────────────────────────────────────
def slide6_formats(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, LIGHT_GRAY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=DARK_NAVY)
    add_textbox(slide, "Bagian D — 5 Format Komentar Terstruktur",
                Inches(0.45), Inches(0.2), Inches(11), Inches(0.7),
                font_size=Pt(27), bold=True, color=WHITE)

    formats = [
        ("#1  Analitis",
         "Masalah · Lokasi · Rekomendasi",
         "Masalah teknis/metodologis yang memerlukan tindakan konkret"),
        ("#2  Klarifikasi",
         "Pertanyaan · Konteks",
         "Bagian ambigu — butuh penjelasan sebelum menilai validitasnya"),
        ("#3  Posisi Literatur",
         "Klaim paper · Kesenjangan · Rekomendasi",
         "Klaim kebaruan tak didukung, sitasi hilang/usang"),
        ("#4  Reproduksibilitas",
         "Elemen hilang · Dampak · Rekomendasi",
         "Detail eksperimen tidak cukup untuk replikasi hasil"),
        ("#5  Apresiasi",
         "Aspek · Alasan",
         "Kekuatan genuine paper — minimal 1 per review"),
    ]

    card_w = Inches(2.35)
    card_h = Inches(5.3)
    gap    = Inches(0.21)
    start_x = Inches(0.38)
    start_y = Inches(1.28)

    colors = [ACCENT_BLUE, RGBColor(0x8E, 0x44, 0xAD),
              RGBColor(0x16, 0x7F, 0x39), RGBColor(0xC0, 0x39, 0x2B),
              RGBColor(0xE6, 0x7E, 0x22)]

    for idx, (title, fields, when) in enumerate(formats):
        x = start_x + idx * (card_w + gap)
        c = colors[idx]

        add_rect(slide, x, start_y, card_w, card_h, fill=WHITE,
                 line_color=c, line_width=Pt(1.5))
        # top color strip
        add_rect(slide, x, start_y, card_w, Inches(0.56), fill=c)
        # number + title
        add_textbox(slide, title, x + Inches(0.08), start_y + Inches(0.06),
                    card_w - Inches(0.16), Inches(0.46),
                    font_size=Pt(13), bold=True, color=WHITE)

        # "Kapan" label
        add_textbox(slide, "KAPAN DIGUNAKAN", x + Inches(0.1), start_y + Inches(0.72),
                    card_w - Inches(0.2), Inches(0.3),
                    font_size=Pt(8), bold=True, color=c)
        add_textbox(slide, when, x + Inches(0.1), start_y + Inches(1.0),
                    card_w - Inches(0.2), Inches(1.1),
                    font_size=Pt(9.5), color=DARK_NAVY)

        # separator
        add_rect(slide, x + Inches(0.1), start_y + Inches(2.15),
                 card_w - Inches(0.2), Inches(0.03), fill=c)

        # "Fields" label
        add_textbox(slide, "STRUKTUR FIELD", x + Inches(0.1), start_y + Inches(2.26),
                    card_w - Inches(0.2), Inches(0.3),
                    font_size=Pt(8), bold=True, color=c)
        for fi, field in enumerate(fields.split(" · ")):
            fy = start_y + Inches(2.6) + fi * Inches(0.55)
            add_rect(slide, x + Inches(0.1), fy, card_w - Inches(0.2), Inches(0.46),
                     fill=LIGHT_BLUE, line_color=c, line_width=Pt(0.5))
            add_textbox(slide, field, x + Inches(0.15), fy + Inches(0.06),
                        card_w - Inches(0.3), Inches(0.36),
                        font_size=Pt(10), bold=True, color=c)



# ── SLIDE 7 — Bagian E: Etika AI + Penutup ───────────────────────────────────
def slide7_etika(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, DARK_NAVY)

    # decorative circle
    c = slide.shapes.add_shape(9, Inches(10.5), Inches(4.5), Inches(4), Inches(4))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT_BLUE
    c.line.fill.background()

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=ACCENT_BLUE)
    add_textbox(slide, "Bagian E — Aturan Etika Penggunaan AI dalam Review",
                Inches(0.45), Inches(0.2), Inches(12), Inches(0.7),
                font_size=Pt(24), bold=True, color=WHITE)

    rules = [
        ("DILARANG", RGBColor(0xFF, 0x6B, 0x6B),
         "Mengunggah naskah ke AI publik (mis. ChatGPT.com)\n"
         "→ Naskah rahasia; pelanggaran terhadap etika COPE & IEEE"),
        ("DILARANG", RGBColor(0xFF, 0x6B, 0x6B),
         "Mendelegasikan keputusan akhir ke AI\n"
         "→ Rekomendasi Accept/Reject = tanggung jawab eksklusif reviewer"),
        ("DILARANG", RGBColor(0xFF, 0x6B, 0x6B),
         "Mengklaim kelemahan yang tidak ada dalam naskah (halusinasi AI)\n"
         "→ Setiap komentar harus dapat ditelusuri ke lokasi nyata"),
        ("DIIZINKAN", RGBColor(0x2E, 0xCC, 0x71),
         "AI sebagai asisten bahasa, penstruktur checklist, dan brainstorming\n"
         "→ Reviewer wajib memverifikasi setiap poin terhadap isi paper"),
    ]

    for idx, (tag, color, text) in enumerate(rules):
        y = Inches(1.28) + idx * Inches(1.38)
        add_rect(slide, Inches(0.4), y, Inches(12.4), Inches(1.22), fill=WHITE,
                 line_color=color, line_width=Pt(1.5))
        add_rect(slide, Inches(0.4), y, Inches(0.22), Inches(1.22), fill=color)
        add_rect(slide, Inches(0.7), y + Inches(0.16), Inches(1.55), Inches(0.4), fill=color)
        add_textbox(slide, tag, Inches(0.7), y + Inches(0.16), Inches(1.55), Inches(0.4),
                    font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, text, Inches(2.4), y + Inches(0.12), Inches(10.2), Inches(1.0),
                    font_size=Pt(11), color=DARK_NAVY)

    # footer
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), fill=ACCENT_BLUE)
    add_textbox(slide,
                "Selaras dengan: COPE Ethical Guidelines for Peer Reviewers (2017)  ·  "
                "IEEE PSPB Operations Manual  ·  ACM Policy on Authorship and Review",
                Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.5),
                font_size=Pt(10), color=WHITE, italic=True, align=PP_ALIGN.CENTER)


# ── MAIN ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()

    slide1_cover(prs)
    slide2_alur(prs)
    slide3_ab(prs)
    slide4_checklist(prs)
    slide5_severity(prs)
    slide6_formats(prs)
    slide7_etika(prs)

    out = "tutorials/Review_ICICoS_2026/ICICoS_Reviewer_Guide_2026.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")

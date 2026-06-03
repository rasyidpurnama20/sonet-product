"""
Generate a clear, decision-grade 5-slide product deck:
"Indonesia's Roadmap to Lead INTOSAI — Sustainable Financing Accountability (SGC + AI)"

Distilled from the 11-slide AI-professor framing deck into ONE coherent product
that guides Indonesia's leadership on the two INCOSAI 2028 themes:
  Theme 1 - The Sustainability Governance Continuum (SGC) ... WHAT to audit
  Theme 2 - AI-Enabled Innovation & Shared Knowledge ........ HOW audit becomes
            intelligent, predictive, and scalable.

Author framing: "AI professor" — rigorous, evidence-based, neat.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette
# ----------------------------------------------------------------------------
NAVY    = RGBColor(0x0B, 0x3D, 0x5C)   # primary
NAVY_DK = RGBColor(0x07, 0x2A, 0x40)   # deep
GREEN   = RGBColor(0x1F, 0x8A, 0x70)   # SGC / sustainability
BLUE    = RGBColor(0x2E, 0x6F, 0xB7)   # AI / technology
GOLD    = RGBColor(0xE0, 0xA9, 0x3B)   # accent
RED     = RGBColor(0xC8, 0x10, 0x2E)   # Indonesia accent (sparing)
LIGHT   = RGBColor(0xF4, 0xF7, 0xFA)   # page background
CARD    = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1B, 0x2A, 0x38)   # body text
MUTED   = RGBColor(0x5C, 0x6E, 0x7E)   # secondary text
LINE    = RGBColor(0xD7, 0xE0, 0xE8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_G  = RGBColor(0xE9, 0xF4, 0xF0)   # soft green tint
SOFT_B  = RGBColor(0xE8, 0xF0, 0xF9)   # soft blue tint

FONT = "Calibri"
FONT_H = "Calibri"

EMU_IN = 914400
PW, PH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(PW * EMU_IN))
prs.slide_height = Emu(int(PH * EMU_IN))
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
             rounded=False, radius=0.08):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    _no_shadow(s)
    if rounded:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    return s


def add_oval(slide, x, y, d, fill, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                               Inches(d), Inches(d))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    _no_shadow(s)
    return s


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=2,
             wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold,
    color, font, italic) tuples (font/italic optional)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for tup in para:
            text, size, bold, color = tup[0], tup[1], tup[2], tup[3]
            font = tup[4] if len(tup) > 4 else FONT
            italic = tup[5] if len(tup) > 5 else False
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
            r.font.color.rgb = color
    return tb


def page_bg(slide, color=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, kicker, title, accent=NAVY):
    # top accent strip
    add_rect(slide, 0, 0, PW, 0.16, fill=accent)
    add_rect(slide, 0, 0.16, PW, 0.03, fill=GOLD)
    add_text(slide, 0.6, 0.34, 11.5, 0.3,
             [[(kicker, 11.5, True, accent)]],
             space_after=0)
    add_text(slide, 0.6, 0.6, 12.1, 0.6,
             [[(title, 27, True, INK)]], space_after=0)
    # underline tick
    add_rect(slide, 0.62, 1.18, 0.85, 0.045, fill=GOLD)


def footer(slide, n):
    add_rect(slide, 0.6, 7.06, 12.13, 0.012, fill=LINE)
    add_text(slide, 0.6, 7.1, 9.0, 0.3,
             [[("Indonesia \u2022 INTOSAI Leadership Proposal", 8.5, False, MUTED)]],
             space_after=0)
    add_text(slide, 11.0, 7.1, 1.73, 0.3,
             [[("Menuju INCOSAI 2028  \u00b7  ", 8.5, False, MUTED),
               (f"{n} / 5", 8.5, True, NAVY)]],
             align=PP_ALIGN.RIGHT, space_after=0)


def chip(slide, x, y, w, h, label, fill, txt_color=WHITE, size=10.5, bold=True):
    add_rect(slide, x, y, w, h, fill=fill, rounded=True, radius=0.5)
    add_text(slide, x, y, w, h, [[(label, size, bold, txt_color)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


# ============================================================================
# SLIDE 1 — COVER / THESIS
# ============================================================================
s = prs.slides.add_slide(BLANK)
page_bg(s, NAVY)
# layered background bands
add_rect(s, 0, 0, PW, PH, fill=NAVY)
add_rect(s, 0, 0, 0.22, PH, fill=GOLD)
add_rect(s, 0, 5.95, PW, 0.04, fill=GREEN)
# subtle large faint number / motif block
add_rect(s, 8.7, 0, 4.63, PH, fill=NAVY_DK)

add_text(s, 0.85, 0.95, 9.0, 0.4,
         [[("PROPOSAL KEPEMIMPINAN INDONESIA DI INTOSAI", 13, True, GOLD)]])
add_text(s, 0.82, 1.45, 11.6, 1.7,
         [[("Sustainable Financing", 44, True, WHITE)],
          [("Accountability", 44, True, WHITE)]],
         line_spacing=1.0, space_after=0)
add_text(s, 0.85, 3.25, 7.6, 1.0,
         [[("Satu arsitektur akuntabilitas untuk dua tema INCOSAI 2028: ",
            15, False, RGBColor(0xD8, 0xE4, 0xEE)),
           ("SGC", 15, True, GOLD),
           (" sebagai kerangka tata kelola, ", 15, False, RGBColor(0xD8, 0xE4, 0xEE)),
           ("AI", 15, True, GOLD),
           (" sebagai infrastruktur kecerdasan.", 15, False, RGBColor(0xD8, 0xE4, 0xEE))]],
         line_spacing=1.15)

# three pillar chips
cy = 4.85
chip(s, 0.85, cy, 3.6, 0.95,
     "1  SGC", GREEN)
chip(s, 4.65, cy, 3.6, 0.95, "2  AI", BLUE)
chip(s, 8.45, cy, 4.05, 0.95, "3  OUTPUT", GOLD, txt_color=NAVY_DK)
# chip sublabels
add_text(s, 0.85, 4.55, 3.6, 0.3, [[("Apa yang diaudit", 9.5, True, RGBColor(0xBFD0, 0xBF00//256, 0xC8))]], align=PP_ALIGN.CENTER) if False else None
add_text(s, 0.85, 4.55, 3.6, 0.28, [[("APA yang diaudit", 9, True, RGBColor(0x9F, 0xC9, 0xBC))]], align=PP_ALIGN.CENTER, space_after=0)
add_text(s, 4.65, 4.55, 3.6, 0.28, [[("BAGAIMANA audit jadi cerdas", 9, True, RGBColor(0xA9, 0xC6, 0xE6))]], align=PP_ALIGN.CENTER, space_after=0)
add_text(s, 8.45, 4.55, 4.05, 0.28, [[("Toolkit \u2022 guidance \u2022 legacy", 9, True, GOLD)]], align=PP_ALIGN.CENTER, space_after=0)

add_text(s, 0.85, 6.55, 11.6, 0.5,
         [[("AI-professor framing deck", 10.5, True, GOLD),
           ("   \u2014  audit masa depan mengikuti uang, dampak, risiko, dan pengetahuan: etis, berbasis bukti, berbantuan AI.",
            10.5, False, RGBColor(0xC2, 0xD2, 0xDE))]])


# ============================================================================
# SLIDE 2 — DUA TEMA, SATU ARSITEKTUR
# ============================================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
header(s, "TESIS UTAMA", "Dua Tema, Satu Arsitektur Akuntabilitas")
add_text(s, 0.62, 1.32, 12.1, 0.4,
         [[("Dua tema dibaca sebagai ", 12.5, False, MUTED),
           ("satu sistem", 12.5, True, NAVY),
           (", bukan dua subtema terpisah \u2014 SGC menentukan ", 12.5, False, MUTED),
           ("apa", 12.5, True, GREEN),
           (", AI menentukan ", 12.5, False, MUTED),
           ("bagaimana", 12.5, True, BLUE),
           (".", 12.5, False, MUTED)]], space_after=0)

# Card A - SGC
ax, ay, aw, ah = 0.62, 1.95, 5.75, 3.55
add_rect(s, ax, ay, aw, ah, fill=CARD, line=LINE, line_w=1.0, rounded=True, radius=0.05)
add_rect(s, ax, ay, aw, 0.85, fill=GREEN, rounded=True, radius=0.05)
add_rect(s, ax, ay + 0.45, aw, 0.40, fill=GREEN)  # square off bottom of header
add_text(s, ax + 0.3, ay + 0.12, aw - 0.6, 0.7,
         [[("THEME 1", 10, True, RGBColor(0xCFE,  0x00//1 and 0xE8, 0xDD) if False else RGBColor(0xCF, 0xE8, 0xDD))],
          [("The Sustainability Governance Continuum (SGC)", 14.5, True, WHITE)]],
         space_after=0, line_spacing=1.0)
add_text(s, ax + 0.3, ay + 1.0, aw - 0.6, 0.5,
         [[("Menjawab ", 12, True, INK), ("APA yang harus diaudit.", 12, True, GREEN)]],
         space_after=0)
bullets_sgc = [
    "Membawa INTOSAI beyond UN 2030 Agenda \u2014 sustainability sebagai akuntabilitas berkelanjutan.",
    "Menghubungkan finance \u2192 governance \u2192 risk \u2192 impact dalam satu rangkaian.",
    "Fokus audit: climate finance, green budget, energy transition, public debt, PPP/blended finance.",
]
yb = ay + 1.55
for b in bullets_sgc:
    add_oval(s, ax + 0.32, yb + 0.06, 0.11, GREEN)
    add_text(s, ax + 0.55, yb - 0.04, aw - 0.85, 0.6,
             [[(b, 11, False, INK)]], line_spacing=1.04, space_after=0)
    yb += 0.62

# Card B - AI
bx = ax + aw + 0.34
add_rect(s, bx, ay, aw, ah, fill=CARD, line=LINE, line_w=1.0, rounded=True, radius=0.05)
add_rect(s, bx, ay, aw, 0.85, fill=BLUE, rounded=True, radius=0.05)
add_rect(s, bx, ay + 0.45, aw, 0.40, fill=BLUE)
add_text(s, bx + 0.3, ay + 0.12, aw - 0.6, 0.7,
         [[("THEME 2", 10, True, RGBColor(0xCD, 0xE0, 0xF4))],
          [("AI-Enabled Innovation & Shared Knowledge", 14.5, True, WHITE)]],
         space_after=0, line_spacing=1.0)
add_text(s, bx + 0.3, ay + 1.0, aw - 0.6, 0.5,
         [[("Menjawab ", 12, True, INK), ("BAGAIMANA audit jadi cerdas & skalabel.", 12, True, BLUE)]],
         space_after=0)
bullets_ai = [
    "AI sebagai intelligence infrastructure \u2014 bukan sekadar plugin audit.",
    "Knowledge hub lintas SAI: shared memory, multilingual search, automated learning.",
    "Predictive analytics, NLP, anomaly detection, dan real-time monitoring.",
]
yb = ay + 1.55
for b in bullets_ai:
    add_oval(s, bx + 0.32, yb + 0.06, 0.11, BLUE)
    add_text(s, bx + 0.55, yb - 0.04, aw - 0.85, 0.6,
             [[(b, 11, False, INK)]], line_spacing=1.04, space_after=0)
    yb += 0.62

# Convergence banner
add_rect(s, 0.62, 5.75, 12.11, 0.9, fill=NAVY, rounded=True, radius=0.12)
add_text(s, 0.62, 5.75, 12.11, 0.9,
         [[("INTEGRATED PARADIGM   ", 11, True, GOLD),
           ("AI-powered Sustainable Financing Accountability", 16, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
footer(s, 2)


# ============================================================================
# SLIDE 3 — THE PRODUCT: SGC-AI ACCOUNTABILITY CONTINUUM
# ============================================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
header(s, "MODEL INTI \u2014 PRODUK", "SGC\u2013AI Accountability Continuum")
add_text(s, 0.62, 1.32, 12.1, 0.4,
         [[("Satu alur audit dari komitmen pembiayaan hingga ketahanan pasca-2030, ", 12.5, False, MUTED),
           ("diperkuat lapisan AI.", 12.5, True, NAVY)]], space_after=0)

# "follow" tags
fx = 0.62
for lbl, col in [("Follow the money", GREEN), ("Follow the impact", GOLD), ("Follow the risk", RED)]:
    chip(s, fx, 1.78, 2.45, 0.42, lbl, col, size=10)
    fx += 2.62

stages = [
    ("1", "Commitment", "Komitmen pembiayaan", GREEN),
    ("2", "Allocation", "Alokasi & governance", GREEN),
    ("3", "Integrity", "Implementasi & integritas", GOLD),
    ("4", "Impact", "Outcome & dampak", GOLD),
    ("5", "Resilience", "Kesiapan pasca-2030", BLUE),
    ("6", "Learning", "Shared knowledge antar-SAI", BLUE),
]
n = len(stages)
left, right = 0.62, 12.73
gap = 0.30
bw = (right - left - gap * (n - 1)) / n
by, bh = 2.45, 2.05
for i, (num, title, desc, col) in enumerate(stages):
    x = left + i * (bw + gap)
    add_rect(s, x, by, bw, bh, fill=CARD, line=LINE, line_w=1.0, rounded=True, radius=0.10)
    add_rect(s, x, by, bw, 0.12, fill=col, rounded=True, radius=0.5)
    # number badge
    add_oval(s, x + bw / 2 - 0.27, by + 0.26, 0.54, col)
    add_text(s, x + bw / 2 - 0.27, by + 0.26, 0.54, 0.54,
             [[(num, 17, True, WHITE)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    add_text(s, x + 0.06, by + 0.92, bw - 0.12, 0.4,
             [[(title, 12.5, True, INK)]], align=PP_ALIGN.CENTER, space_after=0)
    add_text(s, x + 0.08, by + 1.32, bw - 0.16, 0.65,
             [[(desc, 9.5, False, MUTED)]], align=PP_ALIGN.CENTER,
             line_spacing=1.0, space_after=0)
    # connector arrow
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                Inches(x + bw + 0.02), Inches(by + bh / 2 - 0.13),
                                Inches(gap - 0.04), Inches(0.26))
        ar.fill.solid(); ar.fill.fore_color.rgb = LINE
        ar.line.fill.background(); _no_shadow(ar)

# AI layer bar
aly = by + bh + 0.35
add_rect(s, 0.62, aly, 12.11, 0.92, fill=NAVY, rounded=True, radius=0.10)
add_rect(s, 0.62, aly, 2.1, 0.92, fill=BLUE, rounded=True, radius=0.10)
add_rect(s, 1.6, aly, 1.12, 0.92, fill=BLUE)
add_text(s, 0.62, aly, 2.1, 0.92, [[("AI", 22, True, WHITE)], [("LAYER", 11, True, RGBColor(0xCD,0xE0,0xF4))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.95)
add_text(s, 2.95, aly, 9.6, 0.92,
         [[("Data integration  \u2022  NLP  \u2022  anomaly detection  \u2022  predictive modeling  \u2022  dashboards  \u2022  multilingual knowledge retrieval",
            12.5, True, WHITE)],
          [("AI bukan plugin audit, melainkan lapisan epistemik: mengubah data menjadi bukti, insight, dan pembelajaran kolektif.",
            10, False, RGBColor(0xC2, 0xD2, 0xDE), FONT, True)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)
footer(s, 3)


# ============================================================================
# SLIDE 4 — WARISAN INDONESIA (DELIVERABLES)
# ============================================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
header(s, "OUTPUT KEPEMIMPINAN", "Warisan Konkret Indonesia untuk INTOSAI")
add_text(s, 0.62, 1.32, 12.1, 0.4,
         [[("Dari gagasan menjadi standar, teknologi, dan pembelajaran \u2014 ", 12.5, False, MUTED),
           ("Indonesia sebagai arsitek, bukan sekadar tuan rumah.", 12.5, True, NAVY)]],
         space_after=0)

cols = [
    ("NORMATIVE LEGACY", GREEN, SOFT_G, [
        "Jakarta Declaration on Sustainable Financing Accountability",
        "INTOSAI Guidance on the SGC",
        "Post-2030 Sustainability Audit Principles",
    ]),
    ("TECHNOLOGICAL LEGACY", BLUE, SOFT_B, [
        "INTOSAI AI Audit Toolkit",
        "SGC Dashboard (real-time monitoring)",
        "AI-enabled Knowledge Hub",
        "Audit Data Commons",
    ]),
    ("CAPACITY LEGACY", GOLD, RGBColor(0xFB, 0xF1, 0xDD), [
        "AI Audit Certification",
        "Community of Practice lintas SAI",
        "Regional Centre of Excellence (Indonesia)",
        "Joint Audit Learning Network",
    ]),
]
cw, ch = 3.92, 3.55
cx0, cyy = 0.62, 1.95
gapc = 0.18
for j, (title, col, tint, items) in enumerate(cols):
    x = cx0 + j * (cw + gapc)
    add_rect(s, x, cyy, cw, ch, fill=tint, line=LINE, line_w=1.0, rounded=True, radius=0.06)
    add_rect(s, x, cyy, cw, 0.62, fill=col, rounded=True, radius=0.10)
    add_rect(s, x, cyy + 0.32, cw, 0.30, fill=col)
    add_text(s, x, cyy, cw, 0.62, [[(title, 12.5, True, WHITE if col != GOLD else NAVY_DK)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    yy = cyy + 0.85
    for it in items:
        add_rect(s, x + 0.28, yy + 0.07, 0.14, 0.14, fill=col)
        add_text(s, x + 0.52, yy - 0.02, cw - 0.78, 0.6,
                 [[(it, 11.5, False, INK)]], line_spacing=1.05, space_after=0)
        yy += 0.62

# guardrail banner
add_rect(s, 0.62, 5.72, 12.11, 0.95, fill=NAVY, rounded=True, radius=0.10)
add_rect(s, 0.62, 5.72, 0.16, 0.95, fill=GOLD)
add_text(s, 1.0, 5.72, 11.5, 0.95,
         [[("RESPONSIBLE AI  ", 11, True, GOLD),
           ("AI augments the professional judgment of auditors \u2014 it does not replace audit responsibility.",
            13.5, True, WHITE)],
          [("Explainability \u00b7 Fairness \u00b7 Privacy \u00b7 Human Oversight \u00b7 Quality Assurance \u00b7 Accountability",
            10.5, False, RGBColor(0xC2, 0xD2, 0xDE))]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.05)
footer(s, 4)


# ============================================================================
# SLIDE 5 — ROADMAP & POSITIONING
# ============================================================================
s = prs.slides.add_slide(BLANK)
page_bg(s)
header(s, "DARI PROPOSAL KE AGENDA", "Roadmap & Positioning Indonesia")

phases = [
    ("2026\u20132027", "Design & Coalition Building", GREEN,
     "Bentuk working group; finalisasi konsep SGC; pemetaan data; kurasi use case AI."),
    ("2028", "Adoption at INCOSAI", GOLD,
     "Jakarta Declaration; INTOSAI Guidance; toolkit beta; prototipe SGC Dashboard."),
    ("2029\u20132031", "Scaling & Institutionalization", BLUE,
     "Standar, training, peer learning, dan integrasi ke siklus audit SAI."),
]
pw_, ph_ = 3.92, 2.45
px0, py = 0.62, 1.95
gp = 0.18
for j, (yr, title, col, desc) in enumerate(phases):
    x = px0 + j * (pw_ + gp)
    add_rect(s, x, py, pw_, ph_, fill=CARD, line=LINE, line_w=1.0, rounded=True, radius=0.06)
    add_rect(s, x, py, pw_, 0.1, fill=col, rounded=True, radius=0.5)
    add_text(s, x + 0.28, py + 0.25, pw_ - 0.5, 0.5,
             [[(yr, 19, True, col)]], space_after=0)
    add_text(s, x + 0.28, py + 0.78, pw_ - 0.5, 0.5,
             [[(title, 13.5, True, INK)]], space_after=0, line_spacing=1.0)
    add_text(s, x + 0.28, py + 1.34, pw_ - 0.56, 1.0,
             [[(desc, 11, False, MUTED)]], line_spacing=1.1, space_after=0)
    # phase arrow
    if j < len(phases) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                Inches(x + pw_ + 0.0), Inches(py + ph_ / 2 - 0.16),
                                Inches(gp), Inches(0.32))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD
        ar.line.fill.background(); _no_shadow(ar)

# success metric strip
add_rect(s, 0.62, 4.62, 12.11, 0.62, fill=SOFT_G, line=GREEN, line_w=1.0, rounded=True, radius=0.18)
add_text(s, 0.62, 4.62, 12.11, 0.62,
         [[("SUCCESS METRIC   ", 10.5, True, GREEN),
           ("INTOSAI mengadopsi paradigma audit yang berulang untuk sustainable finance pasca-2030, didukung responsible AI.",
            12, True, INK)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

# positioning banner
add_rect(s, 0.62, 5.42, 12.11, 1.25, fill=NAVY, rounded=True, radius=0.08)
add_rect(s, 0.62, 5.42, 0.22, 1.25, fill=RED)
add_rect(s, 0.84, 5.42, 0.10, 1.25, fill=WHITE)
add_text(s, 1.2, 5.55, 11.3, 1.0,
         [[("Indonesia offers a post-2030 accountability architecture for INTOSAI.", 17, True, WHITE)],
          [("Menggabungkan Sustainable Financing Accountability, the Sustainability Governance Continuum, dan AI-enabled shared knowledge \u2014 ",
            11.5, False, RGBColor(0xC8, 0xD6, 0xE2)),
           ("kepemimpinan berbasis gagasan, teknologi, dan akuntabilitas global.",
            11.5, True, GOLD)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=4, line_spacing=1.08)
footer(s, 5)


out = "INTOSAI/Indonesia_INTOSAI_Leadership_Product_5Slides.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
